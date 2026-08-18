"""md2pptx のユニットテスト・結合テスト

モジュール構成に合わせて以下の単位でテストを分けている。
  - utils.py      : 描画・レイアウトのヘルパー関数
  - processors.py : Markdownの各タグに対応する処理
  - generator.py  : プレゼンテーション全体の組み立て
  - md2pptx.py    : CLI
  - gui_config.py : GUI（gui.py）の設定値の組み立て
"""

import base64
import io
import os
import sys
from dataclasses import replace
from datetime import datetime
import subprocess
from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock, patch

import markdown
import pytest
import yaml
from PIL import Image
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

import extractor
import font_metrics
import gui_deps
import gui_runner
import make_template
import md2pptx
import mermaid_renderer
import pptx2md
import processors
import utils
import text_metrics
from config_schema import validate_config
from extractor import extract, write_images
from generator import PPTXGenerator, TemplateError
from layout import SlideLayout
from text_metrics import (
    LINE_HEIGHT_RATIO,
    ParagraphMetrics,
    line_height_ratio,
    TableRowMetrics,
    estimate_row_heights_pt,
    estimate_table_height_pt,
    fit_table_scale,
    paginate_row_heights,
    char_width_ratio,
    estimate_height_pt,
    estimate_line_count,
    estimate_text_width_pt,
    fit_scale,
)
from mermaid_renderer import MermaidRenderError, mermaid_conf, render_mermaid
from md2pptx import apply_theme, load_config, main, parse_args, read_text_file
from gui_config import (
    FontSetting,
    GuiSettings,
    build_config,
    default_output_path,
    settings_from_config,
    validate_inputs,
)
from processors import (
    process_blockquote,
    process_code_or_mermaid,
    process_h3,
    process_heading,
    process_hr,
    process_image,
    process_text,
)
from table_render import cell_alignment, process_table, table_row_metrics
from utils import (
    DEFAULT_IMAGE_DPI,
    find_body_placeholder,
    inherited_bullet,
    apply_auto_numbering,
    disable_bullet,
    set_alt_text,
    add_runs_from_tag,
    append_text_block,
    create_code_textbox,
    apply_font_style,
    auto_shrink_text,
    downscale_image,
    hex_to_rgb,
    insert_image_fit,
    shrink_body_shape,
)

# 1x1ピクセルの最小PNG（画像挿入テスト用のダミーデータ）
TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAE"
    "hQGAhKmMIQAAAABJRU5ErkJggg=="
)


# --- フィクスチャ（テスト用の共通設定） ---


@pytest.fixture
def base_config():
    """テスト用の基本的なYAML設定のモック"""
    return {
        "slides": {"layout": "16:9", "show_slide_number": False},
        "fonts": {
            "title_h1": {"name": "Meiryo", "size_pt": 44, "bold": True},
            "body": {"name": "Meiryo", "size_pt": 20},
            "inline_code": {"name": "Consolas", "color_rgb": [220, 20, 60]},
        },
        "images": {"default_height_inches": 3.5},
    }


@pytest.fixture
def gen(base_config):
    """スライド未作成のジェネレーター"""
    return PPTXGenerator(base_config)


@pytest.fixture
def gen_with_slide(gen):
    """コンテンツスライドを1枚作成済みのジェネレーター"""
    process_heading(gen, parse_md("## 見出し").find("h2"))
    return gen


@pytest.fixture
def png_file(tmp_path):
    """ローカル画像ファイルのパス"""
    path = tmp_path / "tiny.png"
    path.write_bytes(TINY_PNG)
    return str(path)


@pytest.fixture
def mock_response():
    """requests.get の戻り値のモック"""
    response = MagicMock()
    response.content = TINY_PNG
    response.status_code = 200
    response.raise_for_status.return_value = None
    return response


def parse_md(md_text):
    """Markdown文字列をBeautifulSoupに変換するヘルパー"""
    html = markdown.markdown(md_text, extensions=["extra", "fenced_code", "sane_lists"])
    return BeautifulSoup(html, "html.parser")


def parse_html(html):
    """HTML文字列をBeautifulSoupに変換するヘルパー"""
    return BeautifulSoup(html, "html.parser")


def new_paragraph():
    """スタイル検証用の実物のparagraphオブジェクトを作る"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    return box.text_frame.paragraphs[0]


def textboxes_of(slide):
    """スライド上のテキストボックス（プレースホルダー以外）を返す"""
    return [s for s in slide.shapes if s.has_text_frame and not s.is_placeholder]


# =====================================================================
# utils.py
# =====================================================================


class TestHexToRgb:
    @pytest.mark.parametrize("value", ["#ff8000", "ff8000"])
    def test_parses_with_and_without_hash(self, value):
        """先頭の # の有無にかかわらず解釈できる"""
        assert hex_to_rgb(value) == RGBColor(255, 128, 0)

    @pytest.mark.parametrize("value", [None, ""])
    def test_empty_returns_none(self, value):
        """空の入力は None を返す"""
        assert hex_to_rgb(value) is None


class TestApplyFontStyle:
    def test_all_properties_applied(self):
        """name / size_pt / bold / color_rgb がすべて反映される"""
        run = new_paragraph().add_run()
        apply_font_style(
            run, {"name": "Meiryo", "size_pt": 24, "bold": True, "color_rgb": [1, 2, 3]}
        )
        assert run.font.name == "Meiryo"
        assert run.font.size == Pt(24)
        assert run.font.bold is True
        assert run.font.color.rgb == RGBColor(1, 2, 3)

    def test_partial_config_leaves_others_untouched(self):
        """指定されていないプロパティは変更されない"""
        run = new_paragraph().add_run()
        apply_font_style(run, {"name": "Meiryo"})
        assert run.font.name == "Meiryo"
        assert run.font.size is None
        assert run.font.bold is None

    @pytest.mark.parametrize("config", [None, {}])
    def test_empty_config_is_noop(self, config):
        """設定が空の場合は何も適用しない"""
        run = new_paragraph().add_run()
        apply_font_style(run, config)
        assert run.font.name is None


class TestDownscaleImage:
    """埋め込み画像の縮小（ファイルサイズ削減）"""

    def _png(self, tmp_path, width, height, name="big.png"):
        """指定サイズのグラデーション画像（圧縮しにくい＝縮小効果が出る）を作る"""
        img = Image.new("RGB", (width, height))
        img.putdata(
            [((x * 7) % 256, (y * 13) % 256, (x + y) % 256)
             for y in range(height) for x in range(width)]
        )
        path = tmp_path / name
        img.save(path)
        return str(path)

    @pytest.fixture
    def ignore_size_guard(self, mocker):
        """「再エンコードで大きくなるなら元画像を使う」判定を無効化する

        解像度の計算を検証したいテストでは、PNGの圧縮率（環境やPillowの版で
        変動する）に結果が左右されないよう、この判定を切り離す。
        判定そのものは test_original_is_kept_when_resize_grows で検証している。
        """
        mocker.patch("utils._source_size", return_value=10**9)

    def test_large_image_is_resampled(self, tmp_path, ignore_size_guard):
        """表示サイズに対して過大な画像は縮小される"""
        source = self._png(tmp_path, 1600, 800)
        result = downscale_image(source, Inches(4.0), Inches(3.0), dpi=100)

        with Image.open(result) as img:
            assert img.width == 400  # 4.0インチ × 100dpi
            assert img.height == 200  # 縦横比を維持

    def test_small_image_is_untouched(self, tmp_path):
        """必要解像度以下の画像は再エンコードせず、そのまま返す"""
        source = self._png(tmp_path, 100, 50)
        assert downscale_image(source, Inches(8.0), Inches(3.8), dpi=150) is source

    def test_never_upscales(self, tmp_path):
        """元より大きな表示枠でも拡大はしない"""
        source = self._png(tmp_path, 200, 100)
        result = downscale_image(source, Inches(10.0), Inches(10.0), dpi=300)
        assert result is source

    def test_file_size_is_reduced(self, tmp_path):
        """縮小によって実際にバイト数が減る"""
        source = self._png(tmp_path, 2000, 1000)
        result = downscale_image(source, Inches(4.0), Inches(3.0), dpi=100)

        assert isinstance(result, BytesIO)
        assert result.getbuffer().nbytes < os.path.getsize(source)

    def test_dpi_controls_resolution(self, tmp_path, ignore_size_guard):
        """dpiの指定が出力解像度に反映される"""
        source = self._png(tmp_path, 2000, 1000)

        with Image.open(downscale_image(source, Inches(4.0), Inches(3.0), dpi=50)) as low:
            with Image.open(downscale_image(source, Inches(4.0), Inches(3.0), dpi=200)) as high:
                assert low.width == 200   # 4.0インチ × 50dpi
                assert high.width == 800  # 4.0インチ × 200dpi

    def test_original_is_kept_when_resize_grows(self, tmp_path, mocker):
        """再エンコードで逆に大きくなる画像（図版など）は元データを使う"""
        source = self._png(tmp_path, 2000, 1000)
        mocker.patch("utils._source_size", return_value=1)

        assert downscale_image(source, Inches(4.0), Inches(3.0), dpi=100) is source

    def test_jpeg_format_is_preserved(self, tmp_path):
        """JPEGはJPEGのまま再エンコードされる"""
        img = Image.new("RGB", (1600, 800), (10, 120, 200))
        path = tmp_path / "big.jpg"
        img.save(path)

        with Image.open(downscale_image(str(path), Inches(4.0), Inches(3.0), dpi=100)) as out:
            assert out.format == "JPEG"

    def test_bytesio_source_is_rewound(self, tmp_path):
        """メモリ上の画像を渡した場合も、読み取り位置を先頭に戻して返す"""
        source = self._png(tmp_path, 100, 50)
        buffer = BytesIO(open(source, "rb").read())
        buffer.read()  # 読み取り位置を末尾へ

        result = downscale_image(buffer, Inches(8.0), Inches(3.8), dpi=150)
        assert result.tell() == 0

    def test_broken_image_falls_back(self, capsys):
        """画像として読めないデータは警告のみ出し、元データを返す"""
        broken = BytesIO(b"not an image")
        result = downscale_image(broken, Inches(4.0), Inches(3.0), dpi=100)

        assert result is broken
        assert "画像の縮小に失敗" in capsys.readouterr().out

    def test_bytesio_large_image_is_downscaled(self, tmp_path):
        """メモリ上の画像（ダウンロードした画像・Mermaid図）も縮小対象になる"""
        source = self._png(tmp_path, 2000, 1000)
        buffer = BytesIO(open(source, "rb").read())

        result = downscale_image(buffer, Inches(4.0), Inches(3.0), dpi=100)

        assert result is not buffer
        with Image.open(result) as img:
            assert img.width == 400

    def test_size_of_unmeasurable_source_is_none(self):
        """バイト数を取得できないデータでも例外にしない"""
        broken = MagicMock()
        broken.tell.side_effect = OSError("測定不能")
        assert utils._source_size(broken) is None


class TestImageDpiConfig:
    """images.downscale / images.dpi の設定"""

    def test_default_dpi(self, gen):
        """既定では縮小が有効"""
        assert processors.image_dpi(gen) == DEFAULT_IMAGE_DPI

    def test_dpi_from_config(self, base_config):
        """dpi を設定で変更できる"""
        base_config["images"]["dpi"] = 96
        assert processors.image_dpi(PPTXGenerator(base_config)) == 96

    def test_downscale_can_be_disabled(self, base_config):
        """downscale: false で縮小を無効化できる"""
        base_config["images"]["downscale"] = False
        assert processors.image_dpi(PPTXGenerator(base_config)) is None

    @patch("processors.insert_image_fit")
    @patch("processors.downscale_image")
    def test_place_image_downscales(self, mock_downscale, _mock_fit, gen_with_slide, png_file):
        """縮小が有効な場合は表示枠のサイズで縮小してから配置する"""
        processors.place_image(
            gen_with_slide, png_file, Inches(1.0), Inches(1.5), Inches(8.0), Inches(3.8)
        )
        mock_downscale.assert_called_once_with(
            png_file, Inches(8.0), Inches(3.8), DEFAULT_IMAGE_DPI
        )

    @patch("processors.insert_image_fit")
    @patch("processors.downscale_image")
    def test_place_image_skips_when_disabled(
        self, mock_downscale, mock_fit, base_config, png_file
    ):
        """縮小が無効な場合は元画像のまま配置する"""
        base_config["images"]["downscale"] = False
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        processors.place_image(
            gen, png_file, Inches(1.0), Inches(1.5), Inches(8.0), Inches(3.8)
        )
        mock_downscale.assert_not_called()
        assert mock_fit.call_args[0][1] is png_file


class TestInsertImageFit:
    def _mock_slide(self, width, height):
        mock_slide = MagicMock()
        mock_pic = MagicMock()
        mock_pic.width, mock_pic.height = width, height
        mock_slide.shapes.add_picture.return_value = mock_pic
        return mock_slide

    def test_small_image_is_scaled_up_to_cap(self):
        """小さい画像は最大1.5倍までしか拡大されない"""
        pic = insert_image_fit(self._mock_slide(100, 200), b"dummy", 0, 0, 500, 500)
        assert (pic.width, pic.height) == (150, 300)

    def test_large_image_is_shrunk_to_fit(self):
        """大きい画像はアスペクト比を保って枠内に縮小される"""
        pic = insert_image_fit(self._mock_slide(1000, 500), b"dummy", 0, 0, 500, 500)
        assert (pic.width, pic.height) == (500, 250)

    def test_image_is_centered_in_frame(self):
        """縮小後の画像は枠の中央に配置される"""
        pic = insert_image_fit(self._mock_slide(1000, 500), b"dummy", 100, 200, 500, 500)
        assert pic.left == 100
        assert pic.top == 200 + (500 - 250) / 2


class TestAddRunsFromTag:
    def _runs_for(self, gen, html):
        p = new_paragraph()
        add_runs_from_tag(
            parse_html(html).p, p, {"name": "Meiryo"}, gen.fonts_conf.get("inline_code")
        )
        return p.runs

    def test_bold_and_italic(self, gen):
        """strong/em がフォントスタイルに変換される"""
        runs = self._runs_for(gen, "<p><strong>太字</strong><em>斜体</em></p>")
        assert [r.text for r in runs] == ["太字", "斜体"]
        assert runs[0].font.bold is True
        assert runs[1].font.italic is True

    def test_inline_code_uses_code_font(self, gen):
        """codeタグにはinline_codeの設定が適用される"""
        (run,) = self._runs_for(gen, "<p><code>x = 1</code></p>")
        assert run.font.name == "Consolas"
        assert run.font.color.rgb == RGBColor(220, 20, 60)

    def test_inline_code_falls_back_to_defaults(self, base_config):
        """inline_code が未設定でも既定の等幅フォント・色になる"""
        base_config["fonts"].pop("inline_code")
        gen = PPTXGenerator(base_config)
        (run,) = self._runs_for(gen, "<p><code>x</code></p>")
        assert run.font.name == "Consolas"
        assert run.font.color.rgb == RGBColor(220, 20, 60)

    def test_plain_text_uses_default_font(self, gen):
        """装飾の無いテキストには既定フォントが適用される"""
        (run,) = self._runs_for(gen, "<p>ふつうの文字</p>")
        assert run.font.name == "Meiryo"

    def test_nested_container_is_flattened(self, gen):
        """spanなどのコンテナ要素は再帰的に展開される"""
        runs = self._runs_for(gen, "<p><span>外<strong>中</strong></span></p>")
        assert [r.text for r in runs] == ["外", "中"]
        assert runs[1].font.bold is True

    def test_block_tags_are_skipped(self, gen):
        """リストや表などのブロック要素はインライン処理から除外される"""
        runs = self._runs_for(gen, "<p>本文<ul><li>項目</li></ul><table></table></p>")
        assert [r.text for r in runs] == ["本文"]

    def test_newlines_are_replaced(self, gen):
        """改行は半角スペースに置換される"""
        runs = self._runs_for(gen, "<p>1行目\n2行目</p>")
        assert runs[0].text == "1行目 2行目"


class TestUtilsArePure:
    """utils はジェネレーターに依存しない（引数だけで完結する）"""

    def test_no_generator_import(self):
        """utils モジュールが generator を参照していない"""
        source = (Path(utils.__file__)).read_text(encoding="utf-8")
        assert "generator" not in source

    def test_add_runs_works_without_generator(self):
        """ジェネレーターを作らずにインライン装飾を描画できる"""
        p = new_paragraph()
        add_runs_from_tag(
            parse_html("<p>本文<code>x</code></p>").p,
            p,
            {"name": "Meiryo"},
            {"name": "Courier", "color_rgb": [1, 2, 3]},
        )

        assert [r.text for r in p.runs] == ["本文", "x"]
        assert p.runs[1].font.name == "Courier"
        assert p.runs[1].font.color.rgb == RGBColor(1, 2, 3)

    def test_inline_code_defaults_without_config(self):
        """インラインコードの設定を渡さなくても既定値が使われる"""
        p = new_paragraph()
        add_runs_from_tag(parse_html("<p><code>x</code></p>").p, p, None)

        assert p.runs[0].font.name == "Consolas"
        assert p.runs[0].font.color.rgb == RGBColor(220, 20, 60)

    def test_append_text_block_takes_a_text_frame(self):
        """本文枠を直接渡して段落を追加できる"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        text_frame = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(2)
        ).text_frame

        append_text_block(text_frame, parse_md("本文").find("p"), reuse_first_paragraph=True)
        assert text_frame.paragraphs[0].runs[0].text == "本文"

    def test_bullet_replaces_the_existing_mark(self):
        """行頭記号を切り替えると、以前の指定は残らない"""
        paragraph = new_paragraph()
        apply_auto_numbering(paragraph)
        disable_bullet(paragraph)

        xml = paragraph._element.xml
        assert "buNone" in xml
        assert "buAutoNum" not in xml

    def test_bullet_is_placed_before_its_successors(self):
        """行頭記号は defRPr など後続要素の前に挿入する（スキーマの順序）"""
        paragraph = new_paragraph()
        p_pr = paragraph._element.get_or_add_pPr()
        p_pr.append(OxmlElement("a:defRPr"))

        apply_auto_numbering(paragraph)

        tags = [child.tag.rsplit("}", 1)[-1] for child in p_pr]
        assert tags.index("buAutoNum") < tags.index("defRPr")

    def test_set_alt_text(self):
        """図形に代替テキストを設定できる"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        set_alt_text(box, "説明文")
        assert box._element._nvXxPr.cNvPr.get("descr") == "説明文"

    def test_set_alt_text_ignores_empty(self):
        """空文字を渡した場合は既存の値を書き換えない"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        set_alt_text(box, "元の説明")

        set_alt_text(box, "")
        assert box._element._nvXxPr.cNvPr.get("descr") == "元の説明"

    def test_create_code_textbox_takes_a_slide(self):
        """スライドと寸法を渡すだけでコード枠を作れる"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        create_code_textbox(
            slide, Inches(1), Inches(1), Inches(4), Inches(2),
            content="print(1)", language="python",
            font_conf={"name": "Consolas", "size_pt": 12}, background_rgb=[1, 2, 3],
        )

        box = textboxes_of(slide)[0]
        assert box.fill.fore_color.rgb == RGBColor(1, 2, 3)
        assert "print" in box.text_frame.text


class TestShrinkBodyShape:
    def test_width_is_applied(self, gen_with_slide):
        """指定した幅が本文枠に反映される"""
        shrink_body_shape(gen_with_slide.current_slide, Inches(4.8))
        assert gen_with_slide.current_slide.placeholders[1].width == Inches(4.8)

    def test_inherited_geometry_is_kept(self, gen_with_slide):
        """幅だけを変えてもレイアウトから継承した位置・高さが失われない

        プレースホルダーの寸法は継承値のため、一部だけ書き換えると残りが0になる。
        shrink_body_shape 内の left/top 自己代入が効いていることの回帰テスト。
        """
        shape = gen_with_slide.current_slide.placeholders[1]
        top, height = shape.top, shape.height

        shrink_body_shape(gen_with_slide.current_slide, Inches(4.8))

        assert shape.left > 0
        assert (shape.top, shape.height) == (top, height)

    def test_max_height_is_applied(self, gen_with_slide):
        """max_height を渡すと高さも縮む"""
        shrink_body_shape(gen_with_slide.current_slide, Inches(8.0), max_height=Inches(2.0))
        assert gen_with_slide.current_slide.placeholders[1].height == Inches(2.0)

    def test_without_slide_is_noop(self, gen):
        """スライド未作成でも例外にならない"""
        shrink_body_shape(gen.current_slide, Inches(4.8))


class TestAppendTextBlock:
    def test_first_paragraph_is_reused_then_appended(self, gen_with_slide):
        """最初の段落は空の既存段落を再利用し、以降は追加される"""
        append_text_block(
            gen_with_slide.current_body, parse_md("1つ目").find("p"),
            reuse_first_paragraph=True,
        )
        assert len(gen_with_slide.current_body.paragraphs) == 1

        gen_with_slide.slide_has_text = True
        append_text_block(
            gen_with_slide.current_body, parse_md("2つ目").find("p"),
            reuse_first_paragraph=False,
        )
        assert len(gen_with_slide.current_body.paragraphs) == 2

    def test_level_and_spacing(self, gen_with_slide):
        """レベルと行間・段落後余白が設定される"""
        append_text_block(
            gen_with_slide.current_body, parse_md("* 項目").find("li"),
            reuse_first_paragraph=True, level=2,
        )
        p = gen_with_slide.current_body.paragraphs[0]
        assert p.level == 2
        assert p.space_after == Pt(12)
        assert p.line_spacing == 1.2


class TestAppendCodeTextbox:
    def test_creates_textbox_with_background(self, gen_with_slide):
        """コードは背景色付きの独立したテキストボックスに描画される"""
        processors.append_code_textbox(gen_with_slide, "print(1)", language="python")

        boxes = textboxes_of(gen_with_slide.current_slide)
        assert len(boxes) == 1
        assert boxes[0].fill.fore_color.rgb == RGBColor(40, 44, 52)
        assert gen_with_slide.slide_has_text is True

    def test_background_color_from_theme(self, base_config):
        """theme.code_bg_color で背景色を変更できる"""
        base_config["theme"] = {"code_bg_color": [1, 2, 3]}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        processors.append_code_textbox(gen, "print(1)", language="python")
        assert textboxes_of(gen.current_slide)[0].fill.fore_color.rgb == RGBColor(1, 2, 3)

    def test_syntax_highlight_splits_runs(self, gen_with_slide):
        """シンタックスハイライトによりトークンごとにrunが分割される"""
        processors.append_code_textbox(gen_with_slide, "def f():\n    return 1\n", language="python")

        runs = textboxes_of(gen_with_slide.current_slide)[0].text_frame.paragraphs[0].runs
        assert len(runs) > 1
        assert "".join(r.text for r in runs) == "def f():\n    return 1\n"
        assert all(r.font.color.rgb is not None for r in runs)

    def test_unknown_language_falls_back_to_plain(self, gen_with_slide):
        """未知の言語指定でも例外にせずプレーンテキストとして描画する"""
        processors.append_code_textbox(gen_with_slide, "hello", language="no_such_language")
        runs = textboxes_of(gen_with_slide.current_slide)[0].text_frame.paragraphs[0].runs
        # pygmentsは末尾に改行を補うため、内容の一致のみを確認する
        assert "".join(r.text for r in runs).strip() == "hello"

    def test_language_is_guessed_when_omitted(self, gen_with_slide):
        """言語未指定でも推定してハイライトする"""
        processors.append_code_textbox(gen_with_slide, "def f():\n    return 1\n", language=None)
        runs = textboxes_of(gen_with_slide.current_slide)[0].text_frame.paragraphs[0].runs
        assert "".join(r.text for r in runs) == "def f():\n    return 1\n"

    def test_two_column_layout_shrinks_body(self, gen_with_slide):
        """テキストがある場合は本文枠を縮めて右側に配置する"""
        gen_with_slide.slide_has_text = True
        processors.append_code_textbox(gen_with_slide, "x", language="python")

        layout = gen_with_slide.layout
        assert gen_with_slide.current_slide.placeholders[1].width == layout.split_body_width
        assert textboxes_of(gen_with_slide.current_slide)[0].left == layout.split_right_left

    def test_center_layout(self, gen_with_slide):
        """forced_layout=center では中央寄せの枠になる"""
        gen_with_slide.forced_layout = "center"
        processors.append_code_textbox(gen_with_slide, "x", language="python")
        assert textboxes_of(gen_with_slide.current_slide)[0].left == gen_with_slide.layout.center_left

    def test_default_layout(self, gen_with_slide):
        """テキストが無い場合はスライド幅いっぱいに配置する"""
        processors.append_code_textbox(gen_with_slide, "x", language="python")
        box = textboxes_of(gen_with_slide.current_slide)[0]
        layout = gen_with_slide.layout
        assert (box.left, box.width) == (layout.content_left, layout.content_width)


class TestAutoShrinkText:
    def _fill_body(self, gen, line_count, text=None):
        for i in range(line_count):
            gen.slide_has_text = i > 0
            body = parse_md(text or f"行{i}").find("p")
            append_text_block(
                gen.current_body, body,
                reuse_first_paragraph=not gen.slide_has_text, font_conf={"size_pt": 20},
            )

    def _sizes(self, gen):
        return [
            run.font.size.pt for p in gen.current_body.paragraphs for run in p.runs
            if run.font.size
        ]

    def test_long_single_paragraph_is_shrunk(self, gen_with_slide):
        """折り返しの多い長い1段落も縮小される

        段落数だけを数えていた頃は「1段落」と判定され、
        実際には10行以上に折り返してはみ出していても縮小されなかった。
        """
        self._fill_body(gen_with_slide, 1, text="これは非常に長い一段落です。" * 30)
        auto_shrink_text(gen_with_slide.current_slide)

        assert all(size < 20 for size in self._sizes(gen_with_slide))

    def test_short_single_paragraph_is_untouched(self, gen_with_slide):
        """短い1段落は縮小しない"""
        self._fill_body(gen_with_slide, 1, text="短い本文です。")
        auto_shrink_text(gen_with_slide.current_slide)

        assert all(size == 20 for size in self._sizes(gen_with_slide))

    def test_fullwidth_text_shrinks_more_than_halfwidth(self, gen_with_slide, gen):
        """同じ文字数なら、全角のほうが幅を取るため強く縮小される"""
        self._fill_body(gen_with_slide, 1, text="あ" * 300)
        auto_shrink_text(gen_with_slide.current_slide)
        fullwidth = min(self._sizes(gen_with_slide))

        process_heading(gen, parse_md("## 見出し").find("h2"))
        self._fill_body(gen, 1, text="a" * 300)
        auto_shrink_text(gen.current_slide)
        halfwidth = min(self._sizes(gen))

        assert fullwidth < halfwidth

    def test_shrinks_when_too_many_lines(self, gen_with_slide):
        """行数が枠に収まらない場合はフォントサイズが縮小される"""
        self._fill_body(gen_with_slide, 12)
        auto_shrink_text(gen_with_slide.current_slide)

        sizes = [
            run.font.size for p in gen_with_slide.current_body.paragraphs for run in p.runs
        ]
        assert sizes and all(size < Pt(20) for size in sizes)

    def test_keeps_size_when_few_lines(self, gen_with_slide):
        """行数が少ない場合は縮小しない"""
        self._fill_body(gen_with_slide, 3)
        auto_shrink_text(gen_with_slide.current_slide)

        sizes = [
            run.font.size for p in gen_with_slide.current_body.paragraphs for run in p.runs
        ]
        assert sizes and all(size == Pt(20) for size in sizes)

    def test_heading_and_bullets_fit_in_the_frame(self, base_config, tmp_path):
        """h3見出しを含む本文がスライドからはみ出さない

        sample.md の「現状の分析と課題」で実際にはみ出していたケースの回帰テスト。
        h3 の space_before と、フォントの行高を考慮していなかったことが原因だった。
        """
        base_config["fonts"]["bullet_level_1"] = {"size_pt": 18}
        base_config["fonts"]["title_h3"] = {"size_pt": 20, "bold": True}
        gen = PPTXGenerator(base_config)

        md = (
            "## 現状の分析と課題\n\n現場では多くの課題が山積しています。\n\n"
            "### システムの老朽化\n* 既存システムの動作が遅い\n* メンテナンス担当者が不在\n\n"
            "### コミュニケーションの課題\n* 部署間での情報共有にタイムラグがある\n"
            "* 手作業によるデータ入力の負荷が高い\n"
        )
        gen.generate(md, str(tmp_path / "out.pptx"))

        body = gen.prs.slides[0].placeholders[1]
        tf = body.text_frame
        needed = estimate_height_pt(
            [utils._paragraph_metrics(p) for p in tf.paragraphs],
            Emu(int(body.width - tf.margin_left - tf.margin_right)).pt,
        )
        assert needed <= Emu(int(body.height - tf.margin_top - tf.margin_bottom)).pt

    def test_paragraph_level_font_size_is_also_shrunk(self, gen_with_slide):
        """段落レベルのフォントサイズ（h3が設定する基準値）も縮小対象にする"""
        self._fill_body(gen_with_slide, 1, text="あ" * 400)
        target = gen_with_slide.current_body.paragraphs[0]
        target.font.size = Pt(20)

        auto_shrink_text(gen_with_slide.current_slide)
        assert target.font.size.pt < 20

    def test_space_before_is_shrunk(self, gen_with_slide):
        """段落前の余白も縮小される"""
        self._fill_body(gen_with_slide, 1, text="あ" * 400)
        target = gen_with_slide.current_body.paragraphs[0]
        target.space_before = Pt(10)

        auto_shrink_text(gen_with_slide.current_slide)
        assert target.space_before.pt < 10

    def test_none_slide_is_noop(self):
        """スライドが None でも例外にならない"""
        auto_shrink_text(None)

    def test_unexpected_slide_structure_is_ignored(self):
        """想定外の構造のスライドを渡しても例外を投げない（縮小は best-effort）"""
        broken = MagicMock()
        broken.placeholders.__len__.side_effect = ValueError("broken")
        auto_shrink_text(broken)

    def test_slide_without_body_is_noop(self, gen):
        """本文プレースホルダーが無いレイアウトでも例外にならない"""
        blank = gen.prs.slides.add_slide(gen.prs.slide_layouts[6])
        auto_shrink_text(blank)


# =====================================================================
# processors.py
# =====================================================================


class TestProcessHeading:
    def test_h1_creates_title_slide(self, gen):
        """h1はタイトルスライド（レイアウト0）を作る"""
        process_heading(gen, parse_md("# タイトル").find("h1"))
        assert len(gen.prs.slides) == 1
        assert gen.current_slide.shapes.title.text == "タイトル"
        assert gen.current_slide.slide_layout == gen.prs.slide_layouts[0]

    def test_h2_creates_content_slide(self, gen):
        """h2はコンテンツスライド（レイアウト1）を作る"""
        process_heading(gen, parse_md("## 中身").find("h2"))
        assert gen.current_slide.slide_layout == gen.prs.slide_layouts[1]

    def test_title_font_is_applied(self, gen):
        """title_h1 の設定がタイトルに適用される"""
        process_heading(gen, parse_md("# タイトル").find("h1"))
        run = gen.current_slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(44)
        assert run.font.bold is True

    def test_generic_title_config_is_fallback(self, base_config):
        """title_h2 が無い場合は title の設定にフォールバックする"""
        base_config["fonts"]["title"] = {"size_pt": 18}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 中身").find("h2"))
        run = gen.current_slide.shapes.title.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(18)

    def test_state_is_reset_for_new_slide(self, gen):
        """スライド作成時に本文とテキスト有無フラグが初期化される"""
        gen.slide_has_text = True
        process_heading(gen, parse_md("## 中身").find("h2"))
        assert gen.slide_has_text is False
        assert gen.current_body.text == ""

    def test_body_height_is_clamped_to_slide(self, gen):
        """デフォルトテンプレートでは本文枠がスライド内に収まるよう補正される"""
        process_heading(gen, parse_md("## 中身").find("h2"))
        shape = gen.current_slide.placeholders[1]
        assert shape.top + shape.height == gen.prs.slide_height - Inches(0.5)
        assert shape.width > 0  # 継承値のリセット防止

    def test_previous_slide_is_shrunk(self, gen):
        """新しいスライドを作る前に直前のスライドの自動縮小が走る"""
        process_heading(gen, parse_md("## 1枚目").find("h2"))
        with patch("processors.auto_shrink_text") as mock_shrink:
            process_heading(gen, parse_md("## 2枚目").find("h2"))
        mock_shrink.assert_called_once()

    def test_body_correction_failure_is_ignored(self, gen, mocker):
        """本文枠の補正に失敗してもスライド生成は継続する"""
        mocker.patch.object(SlideLayout, "body_height_for", side_effect=ValueError("broken"))

        process_heading(gen, parse_md("## 中身").find("h2"))
        assert gen.current_slide.shapes.title.text == "中身"


class TestH3AsSlide:
    """slides.h3_as による h3 の扱いの切り替え"""

    MD = "## 現状の課題\n本文です。\n\n### システムの老朽化\n* 動作が遅い\n\n### 情報共有\n* タイムラグ\n"

    def _generate(self, base_config, tmp_path, h3_as=None):
        if h3_as is not None:
            base_config["slides"]["h3_as"] = h3_as
        gen = PPTXGenerator(base_config)
        gen.generate(self.MD, str(tmp_path / "out.pptx"))
        return gen

    @pytest.mark.parametrize("h3_as", [None, "subheading"])
    def test_default_is_subheading(self, base_config, tmp_path, h3_as):
        """既定ではスライド内の小見出しとして扱う（従来の挙動）"""
        gen = self._generate(base_config, tmp_path, h3_as)

        assert len(gen.prs.slides) == 1
        assert "システムの老朽化" in gen.current_body.text

    def test_slide_mode_creates_new_slides(self, base_config, tmp_path):
        """h3_as: slide では h3 ごとに新しいスライドを作る"""
        gen = self._generate(base_config, tmp_path, "slide")

        titles = [s.shapes.title.text for s in gen.prs.slides]
        assert titles == ["現状の課題", "システムの老朽化", "情報共有"]

    def test_slide_mode_uses_the_content_layout(self, base_config, tmp_path):
        """h3から作られるスライドは h2 と同じコンテンツ用レイアウトを使う"""
        gen = self._generate(base_config, tmp_path, "slide")

        assert gen.prs.slides[1].slide_layout == gen.prs.slide_layouts[1]

    def test_slide_mode_moves_content_to_the_new_slide(self, base_config, tmp_path):
        """h3の後の本文は、その h3 のスライドに入る"""
        gen = self._generate(base_config, tmp_path, "slide")

        bodies = [s.placeholders[1].text_frame.text for s in gen.prs.slides]
        assert bodies[0] == "本文です。"
        assert bodies[1] == "動作が遅い"

    def test_setting_is_case_insensitive(self, base_config, tmp_path):
        """大文字で書かれていても解釈する"""
        gen = self._generate(base_config, tmp_path, "SLIDE")
        assert len(gen.prs.slides) == 3

    def test_layout_comment_is_reset_on_new_slide(self, base_config, tmp_path):
        """h3でスライドが変わるとき、直前のレイアウト指定は引き継がない"""
        base_config["slides"]["h3_as"] = "slide"
        gen = PPTXGenerator(base_config)
        gen.generate("## 見出し\n\n<!-- layout: center -->\n\n### 次のスライド\n", str(tmp_path / "o.pptx"))

        assert gen.forced_layout is None

    def test_h3_before_any_heading(self, base_config, tmp_path):
        """先頭が h3 でも、slide 指定ならスライドが作られる"""
        base_config["slides"]["h3_as"] = "slide"
        gen = PPTXGenerator(base_config)
        gen.generate("### いきなり見出し\n本文\n", str(tmp_path / "o.pptx"))

        assert len(gen.prs.slides) == 1
        assert gen.prs.slides[0].shapes.title.text == "いきなり見出し"

    def test_h3_before_any_heading_is_ignored_as_subheading(self, base_config, tmp_path):
        """小見出しモードでは、配置先が無い h3 は無視する"""
        gen = PPTXGenerator(base_config)
        gen.generate("### いきなり見出し\n", str(tmp_path / "o.pptx"))

        assert len(gen.prs.slides) == 0


class TestProcessH3:
    BU_NONE = "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"

    def test_adds_subheading_paragraph(self, gen_with_slide):
        """h3は本文内の小見出し段落として追加される"""
        process_h3(gen_with_slide, parse_md("### 小見出し").find("h3"))

        assert gen_with_slide.current_body.paragraphs[-1].runs[0].text == "小見出し"
        assert gen_with_slide.slide_has_text is True

    def test_bullet_is_disabled(self, gen_with_slide):
        """箇条書き記号が無効化される（buNone要素が入る）"""
        process_h3(gen_with_slide, parse_md("### 小見出し").find("h3"))

        p_pr = gen_with_slide.current_body.paragraphs[-1]._element.get_or_add_pPr()
        assert p_pr.find(self.BU_NONE) is not None

    def test_font_config_is_applied(self, base_config):
        """title_h3 の設定が適用される"""
        base_config["fonts"]["title_h3"] = {"name": "Meiryo", "size_pt": 26, "bold": True}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_h3(gen, parse_md("### 小見出し").find("h3"))
        p = gen.current_body.paragraphs[-1]
        assert p.runs[0].font.size == Pt(26)
        assert p.font.size == Pt(26)

    def test_empty_heading_is_skipped(self, gen_with_slide):
        """空のh3は段落を追加しない"""
        before = len(gen_with_slide.current_body.paragraphs)
        process_h3(gen_with_slide, parse_html("<h3>  </h3>").h3)
        assert len(gen_with_slide.current_body.paragraphs) == before


class TestProcessHr:
    def test_creates_untitled_slide(self, gen_with_slide):
        """水平線はタイトルの無い新しいスライドを作る"""
        process_hr(gen_with_slide, parse_md("---").find("hr"))

        assert len(gen_with_slide.prs.slides) == 2
        assert gen_with_slide.current_slide.shapes.title is None
        assert gen_with_slide.slide_has_text is False

    def test_body_uses_full_height(self, gen_with_slide):
        """タイトルが無い分、本文枠を上部から広く使う"""
        process_hr(gen_with_slide, parse_md("---").find("hr"))

        shape = gen_with_slide.current_slide.placeholders[1]
        assert shape.top == Inches(0.5)
        assert shape.width > 0

    def test_body_correction_failure_is_ignored(self, gen_with_slide, mocker):
        """本文枠の補正に失敗してもスライド生成は継続する"""
        mocker.patch.object(SlideLayout, "body_height_for", side_effect=ValueError("broken"))

        process_hr(gen_with_slide, parse_md("---").find("hr"))
        assert len(gen_with_slide.prs.slides) == 2


class TestProcessBlockquote:
    def test_note_is_written(self, gen_with_slide):
        """引用ブロックはスピーカーノートに書き込まれる"""
        process_blockquote(gen_with_slide, parse_md("> メモ").find("blockquote"))
        assert gen_with_slide.current_slide.notes_slide.notes_text_frame.text == "メモ"

    def test_multiple_notes_are_appended(self, gen_with_slide):
        """複数の引用ブロックは空行を挟んで追記される"""
        process_blockquote(gen_with_slide, parse_md("> 1つ目").find("blockquote"))
        process_blockquote(gen_with_slide, parse_md("> 2つ目").find("blockquote"))

        notes = gen_with_slide.current_slide.notes_slide.notes_text_frame.text
        assert notes == "1つ目\n\n2つ目"

    def test_paragraphs_keep_their_separation(self, gen_with_slide):
        """引用内の段落が区切りなく連結されない

        get_text() をそのまま使っていた頃は
        「1段落目です。2段落目です。」のように繋がっていた。
        """
        md = "> 1段落目です。\n>\n> 2段落目です。\n"
        process_blockquote(gen_with_slide, parse_md(md).find("blockquote"))

        notes = gen_with_slide.current_slide.notes_slide.notes_text_frame.text
        assert notes == "1段落目です。\n\n2段落目です。"

    def test_line_breaks_within_a_paragraph_are_kept(self, gen_with_slide):
        """同じ段落内の改行は維持される（発表原稿の体裁を保つため）"""
        md = "> 1行目です。\n> 2行目です。\n"
        process_blockquote(gen_with_slide, parse_md(md).find("blockquote"))

        notes = gen_with_slide.current_slide.notes_slide.notes_text_frame.text
        assert notes == "1行目です。\n2行目です。"

    def test_list_items_become_separate_lines(self, gen_with_slide):
        """引用内の箇条書きは項目ごとに改行される"""
        md = "> 発表のポイント\n>\n> * 最初の3分で説明する\n> * 質問は別紙を参照\n"
        process_blockquote(gen_with_slide, parse_md(md).find("blockquote"))

        notes = gen_with_slide.current_slide.notes_slide.notes_text_frame.text
        assert notes == "発表のポイント\n最初の3分で説明する\n質問は別紙を参照"

    def test_empty_blockquote_is_skipped(self, gen_with_slide):
        """空の引用ではノートを作らない"""
        process_blockquote(gen_with_slide, parse_html("<blockquote><p>  </p></blockquote>").blockquote)
        assert gen_with_slide.current_slide.notes_slide.notes_text_frame.text == ""

    def test_blockquote_without_paragraph(self, gen_with_slide):
        """pタグを持たない引用でもテキストを取り出せる"""
        process_blockquote(gen_with_slide, parse_html("<blockquote>直書き</blockquote>").blockquote)
        assert gen_with_slide.current_slide.notes_slide.notes_text_frame.text == "直書き"


class TestMultipleImages:
    """同じスライドに複数の画像を置いたときの配置"""

    def _generate(self, gen, count, tmp_path, png_file, text=""):
        images = "\n\n".join(f"![図{n}]({png_file})" for n in range(count))
        gen.generate(f"## 図のスライド\n\n{text}\n\n{images}\n", str(tmp_path / "out.pptx"))
        return [
            s for s in gen.prs.slides[0].shapes
            if s.shape_type is not None and "PICTURE" in str(s.shape_type)
        ]

    @pytest.mark.parametrize("count", [2, 3, 4, 6])
    def test_images_do_not_overlap(self, gen, count, tmp_path, png_file):
        """複数の画像が重ならずに配置される

        従来は各画像が同じ枠に配置され、完全に重なって最後の1枚しか見えなかった。
        """
        pictures = self._generate(gen, count, tmp_path, png_file)

        assert len(pictures) == count
        boxes = [(p.left, p.top, p.width, p.height) for p in pictures]
        for i, a in enumerate(boxes):
            for b in boxes[i + 1:]:
                overlap_x = a[0] < b[0] + b[2] and b[0] < a[0] + a[2]
                overlap_y = a[1] < b[1] + b[3] and b[1] < a[1] + a[3]
                assert not (overlap_x and overlap_y)

    @pytest.mark.parametrize("count, expected", [(1, (1, 1)), (2, (2, 1)), (3, (3, 1)), (4, (2, 2)), (6, (3, 2)), (9, (3, 3))])
    def test_grid_shape(self, count, expected):
        """3枚までは横一列、4枚以上は正方形に近いグリッドにする"""
        assert processors.image_grid(count) == expected

    def test_no_images_is_noop(self, gen_with_slide):
        """画像が無いスライドで呼ばれても何もしない"""
        processors.arrange_images(gen_with_slide)
        assert gen_with_slide.current_images == []

    def test_images_stay_inside_the_content_area(self, gen, tmp_path, png_file):
        """並べた画像がコンテンツ領域からはみ出さない"""
        pictures = self._generate(gen, 4, tmp_path, png_file)
        layout = gen.layout

        for picture in pictures:
            assert picture.left >= layout.content_left
            assert picture.left + picture.width <= layout.content_left + layout.content_width
            assert picture.top + picture.height <= layout.content_top + layout.content_height

    def test_images_go_to_the_right_column_with_text(self, gen, tmp_path, png_file):
        """テキストがある場合、複数画像は右半分の中で並ぶ"""
        pictures = self._generate(gen, 2, tmp_path, png_file, text="説明文です。")

        for picture in pictures:
            assert picture.left >= gen.layout.split_right_left

    def test_single_image_is_centered_in_the_full_area(self, gen, tmp_path, png_file):
        """1枚だけの場合は従来どおりコンテンツ領域の中央に配置する"""
        (picture,) = self._generate(gen, 1, tmp_path, png_file)
        layout = gen.layout

        left_margin = picture.left - layout.content_left
        right_margin = (layout.content_left + layout.content_width) - (
            picture.left + picture.width
        )
        assert left_margin == pytest.approx(right_margin, abs=Inches(0.01))

    def test_images_are_reset_per_slide(self, gen, tmp_path, png_file):
        """スライドが変わると画像の並べ直しは新しいスライド内で完結する"""
        md = f"## 1枚目\n\n![a]({png_file})\n\n## 2枚目\n\n![b]({png_file})\n"
        gen.generate(md, str(tmp_path / "out.pptx"))

        assert len(gen.current_images) == 1

    @patch("processors.arrange_images")
    def test_fixed_position_is_not_rearranged(self, mock_arrange, base_config, tmp_path, png_file):
        """position_inches で明示配置した画像は自動配置の対象にしない"""
        base_config["images"]["position_inches"] = [1.0, 1.0]
        gen = PPTXGenerator(base_config)
        gen.generate(f"## 図\n\n![a]({png_file})\n", str(tmp_path / "out.pptx"))

        mock_arrange.assert_not_called()


class TestProcessImage:
    def test_local_path_is_inserted(self, gen_with_slide, png_file):
        """ローカル画像はダウンロードせずにそのまま挿入される"""
        process_image(gen_with_slide, parse_html(f'<img src="{png_file}">').img)
        assert len(gen_with_slide.current_slide.shapes) == 3  # title, body, picture

    @patch("requests.get")
    def test_remote_image_is_downloaded_with_timeout(
        self, mock_get, gen_with_slide, mock_response
    ):
        """URL画像はタイムアウト付きで取得される"""
        mock_get.return_value = mock_response
        process_image(gen_with_slide, parse_html('<img src="http://example.com/a.png">').img)

        mock_get.assert_called_once_with(
            "http://example.com/a.png", timeout=processors.HTTP_TIMEOUT_SEC
        )
        mock_response.raise_for_status.assert_called_once()

    @patch("processors.insert_image_fit")
    def test_fixed_position_from_config(self, mock_fit, base_config, png_file):
        """position_inches が指定されている場合はオートレイアウトを使わない"""
        base_config["images"]["position_inches"] = [5.2, 1.8]
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))
        gen.current_slide = MagicMock()

        process_image(gen, parse_html(f'<img src="{png_file}">').img)

        mock_fit.assert_not_called()
        gen.current_slide.shapes.add_picture.assert_called_once_with(
            png_file, Inches(5.2), Inches(1.8), height=Inches(3.5)
        )

    @patch("processors.insert_image_fit")
    def test_two_column_when_slide_has_text(self, mock_fit, gen_with_slide, png_file):
        """テキストがあるスライドでは本文枠を縮めて右側に配置する"""
        gen_with_slide.slide_has_text = True
        process_image(gen_with_slide, parse_html(f'<img src="{png_file}">').img)

        layout = gen_with_slide.layout
        args = mock_fit.call_args[0]
        assert (args[2], args[4]) == (layout.split_right_left, layout.split_right_width)
        assert gen_with_slide.current_slide.placeholders[1].width == layout.split_body_width

    @patch("processors.insert_image_fit")
    def test_forced_two_column_layout(self, mock_fit, gen_with_slide, png_file):
        """forced_layout=2-column ではテキストが無くても右側に配置する"""
        gen_with_slide.forced_layout = "2-column"
        process_image(gen_with_slide, parse_html(f'<img src="{png_file}">').img)
        assert mock_fit.call_args[0][2] == gen_with_slide.layout.split_right_left

    @patch("processors.insert_image_fit")
    def test_forced_center_layout(self, mock_fit, gen_with_slide, png_file):
        """forced_layout=center ではテキストの有無によらず中央に配置する"""
        gen_with_slide.slide_has_text = True
        gen_with_slide.forced_layout = "center"
        process_image(gen_with_slide, parse_html(f'<img src="{png_file}">').img)

        layout = gen_with_slide.layout
        assert (mock_fit.call_args[0][2], mock_fit.call_args[0][4]) == (
            layout.content_left,
            layout.content_width,
        )

    def _pictures_of(self, slide):
        return [
            s for s in slide.shapes
            if s.shape_type is not None and "PICTURE" in str(s.shape_type)
        ]

    def test_alt_text_is_carried_over(self, gen_with_slide, png_file):
        """Markdownの代替テキストが図形の説明として設定される"""
        tag = parse_html(f'<img src="{png_file}" alt="システム構成図">').img
        process_image(gen_with_slide, tag)

        picture = self._pictures_of(gen_with_slide.current_slide)[0]
        assert picture._element._nvXxPr.cNvPr.get("descr") == "システム構成図"

    def test_alt_text_for_fixed_position(self, base_config, png_file):
        """position_inches で配置した画像にも代替テキストが付く"""
        base_config["images"]["position_inches"] = [1.0, 1.0]
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_image(gen, parse_html(f'<img src="{png_file}" alt="固定配置の図">').img)

        picture = self._pictures_of(gen.current_slide)[0]
        assert picture._element._nvXxPr.cNvPr.get("descr") == "固定配置の図"

    def test_without_alt_the_default_is_kept(self, gen_with_slide, png_file):
        """alt が無い場合は python-pptx の既定値のままにする"""
        process_image(gen_with_slide, parse_html(f'<img src="{png_file}">').img)

        picture = self._pictures_of(gen_with_slide.current_slide)[0]
        assert picture._element._nvXxPr.cNvPr.get("descr") != ""

    def test_missing_src_is_skipped(self, gen_with_slide, capsys):
        """src属性が無い画像は警告を出してスキップする"""
        process_image(gen_with_slide, parse_html("<img>").img)

        assert "Warning" in capsys.readouterr().out
        assert len(gen_with_slide.current_slide.shapes) == 2

    @patch("requests.get", side_effect=OSError("network down"))
    def test_download_failure_is_reported(self, _mock_get, gen_with_slide, capsys):
        """取得に失敗しても処理は継続し、警告のみ表示する"""
        process_image(gen_with_slide, parse_html('<img src="https://example.com/a.png">').img)
        assert "画像の挿入に失敗しました" in capsys.readouterr().out


class TestProcessTable:
    MD_TABLE = "| 列A | 列B |\n|---|---|\n| 値1 | 値2 |"

    def _table_of(self, gen):
        return [s for s in gen.current_slide.shapes if s.has_table][0].table

    def test_rows_and_cells_are_converted(self, gen_with_slide):
        """行数・列数とセルのテキストが再現される"""
        process_table(gen_with_slide, parse_md(self.MD_TABLE).find("table"))
        table = self._table_of(gen_with_slide)

        assert (len(table.rows), len(table.columns)) == (2, 2)
        assert table.cell(0, 0).text == "列A"
        assert table.cell(1, 1).text == "値2"
        assert gen_with_slide.slide_has_text is True

    def test_header_style(self, gen_with_slide):
        """ヘッダー行は太字・専用サイズ・背景色が付く"""
        process_table(gen_with_slide, parse_md(self.MD_TABLE).find("table"))
        table = self._table_of(gen_with_slide)

        header = table.cell(0, 0)
        assert header.text_frame.paragraphs[0].runs[0].font.bold is True
        assert header.text_frame.paragraphs[0].runs[0].font.size == Pt(14)
        assert header.fill.fore_color.rgb == RGBColor(31, 73, 125)
        assert table.cell(1, 0).text_frame.paragraphs[0].runs[0].font.size == Pt(12)

    def test_inline_decoration_in_cell(self, gen_with_slide):
        """セル内の太字・インラインコードも反映される"""
        process_table(gen_with_slide, parse_md("| 見出し |\n|---|\n| **強調** |").find("table"))
        run = self._table_of(gen_with_slide).cell(1, 0).text_frame.paragraphs[0].runs[0]
        assert run.font.bold is True

    def test_column_count_uses_widest_row(self, gen_with_slide):
        """列数が揃っていない場合は最大列数に合わせる"""
        html = "<table><tr><td>a</td></tr><tr><td>b</td><td>c</td></tr></table>"
        process_table(gen_with_slide, parse_html(html).table)
        assert len(self._table_of(gen_with_slide).columns) == 2

    def test_empty_table_is_skipped(self, gen_with_slide):
        """行が無い表は何も生成しない"""
        process_table(gen_with_slide, parse_html("<table></table>").table)
        assert len(gen_with_slide.current_slide.shapes) == 2

    def test_layout_is_split_when_slide_has_text(self, gen_with_slide):
        """テキストがある場合は本文枠を縮めて表を下半分に置く"""
        gen_with_slide.slide_has_text = True
        process_table(gen_with_slide, parse_md(self.MD_TABLE).find("table"))

        shape = [s for s in gen_with_slide.current_slide.shapes if s.has_table][0]
        assert shape.top == gen_with_slide.layout.table_split_top
        assert (gen_with_slide.current_slide.placeholders[1].height
                == gen_with_slide.layout.table_split_body_height)

    def test_column_alignment_is_applied(self, gen_with_slide):
        """Markdownの列揃え指定（:---: など）がセルに反映される"""
        md = "| 左 | 中央 | 右 |\n| :--- | :---: | ---: |\n| a | b | c |"
        process_table(gen_with_slide, parse_md(md).find("table"))
        table = [s for s in gen_with_slide.current_slide.shapes if s.has_table][0].table

        expected = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT]
        for row_idx in range(2):  # 見出し行と本文行の両方
            actual = [
                table.cell(row_idx, col).text_frame.paragraphs[0].alignment for col in range(3)
            ]
            assert actual == expected

    def test_without_alignment_powerpoint_default_is_used(self, gen_with_slide):
        """列揃えの指定が無い場合はPowerPointの既定に任せる"""
        process_table(gen_with_slide, parse_md(self.MD_TABLE).find("table"))
        table = [s for s in gen_with_slide.current_slide.shapes if s.has_table][0].table
        assert table.cell(0, 0).text_frame.paragraphs[0].alignment is None

    @pytest.mark.parametrize(
        "html, expected",
        [
            ('<td style="text-align: center;">a</td>', PP_ALIGN.CENTER),
            ('<td style="text-align:right">a</td>', PP_ALIGN.RIGHT),
            ('<td align="left">a</td>', PP_ALIGN.LEFT),  # 古いmarkdownの出力形式
            ('<td>a</td>', None),
            ('<td style="color: red;">a</td>', None),
        ],
    )
    def test_cell_alignment_parsing(self, html, expected):
        """揃え指定の解釈（style属性・align属性の両方に対応する）"""
        assert cell_alignment(parse_html(html).td) is expected

    def test_layout_is_full_when_slide_is_empty(self, gen_with_slide):
        """テキストが無い場合は表を上部から配置する"""
        process_table(gen_with_slide, parse_md(self.MD_TABLE).find("table"))
        shape = [s for s in gen_with_slide.current_slide.shapes if s.has_table][0]
        assert shape.top == gen_with_slide.layout.content_top

    def _long_table(self, data_rows):
        body = "\n".join(f"| 項目{i} | 値{i} |" for i in range(data_rows))
        return f"| 見出しA | 見出しB |\n|---|---|\n{body}"

    def _tables_of(self, gen):
        return [
            (i, s) for i, slide in enumerate(gen.prs.slides)
            for s in slide.shapes if s.has_table
        ]

    def test_row_heights_are_explicit(self, gen_with_slide):
        """行の高さを明示する

        既定では総高さを行数で均等割りするため、行数が多いと1行あたりが
        極端に小さくなり、PowerPointが描画時に押し広げてはみ出していた。
        """
        process_table(gen_with_slide, parse_md(self._long_table(3)).find("table"))
        table = self._tables_of(gen_with_slide)[0][1].table

        # 12pt の本文なら 0.2インチ程度では収まらない
        assert all(row.height > Inches(0.25) for row in table.rows)

    def test_long_table_is_split_across_slides(self, gen_with_slide, capsys):
        """1枚に収まらない表は複数スライドに分割される"""
        process_table(gen_with_slide, parse_md(self._long_table(40)).find("table"))

        tables = self._tables_of(gen_with_slide)
        assert len(tables) > 1
        assert "分割しました" in capsys.readouterr().out

    def test_split_pages_stay_inside_the_slide(self, gen_with_slide):
        """分割後はどのページもスライドからはみ出さない"""
        process_table(gen_with_slide, parse_md(self._long_table(40)).find("table"))

        for _, shape in self._tables_of(gen_with_slide):
            bottom = shape.top + sum(row.height for row in shape.table.rows)
            assert bottom <= gen_with_slide.prs.slide_height

    def test_header_is_repeated_on_continuation_slides(self, gen_with_slide):
        """続きのスライドにも見出し行が繰り返される"""
        process_table(gen_with_slide, parse_md(self._long_table(40)).find("table"))

        for _, shape in self._tables_of(gen_with_slide):
            assert shape.table.cell(0, 0).text == "見出しA"

    def test_continuation_slide_title(self, gen_with_slide):
        """続きのスライドのタイトルには「（続き）」が付く"""
        process_table(gen_with_slide, parse_md(self._long_table(40)).find("table"))

        titles = [s.shapes.title.text for s in gen_with_slide.prs.slides]
        assert titles[0] == "見出し"
        assert titles[1] == "見出し（続き）"
        # 「（続き）（続き）」のように重ならない
        assert all(t.count("（続き）") <= 1 for t in titles)

    def test_short_table_is_not_split(self, gen_with_slide):
        """収まる表は分割しない"""
        process_table(gen_with_slide, parse_md(self._long_table(3)).find("table"))
        assert len(gen_with_slide.prs.slides) == 1

    def test_font_is_shrunk_before_splitting(self, base_config):
        """分割の前に、まず縮小で1枚に収めることを試みる"""
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))
        process_table(gen, parse_md(self._long_table(12)).find("table"))

        assert len(gen.prs.slides) == 1
        table = self._tables_of(gen)[0][1].table
        sizes = [
            r.font.size.pt for row in table.rows for c in row.cells
            for p in c.text_frame.paragraphs for r in p.runs if r.font.size
        ]
        assert min(sizes) < 12  # table_body の既定サイズより小さい


class TestProcessCodeOrMermaid:
    def test_code_block_without_language_does_not_crash(self, gen_with_slide):
        """言語指定の無いコードブロックでも落ちない（class属性がNoneになるケース）

        `code_tag.get('class')` は class 属性が無いと None を返すため、
        そのまま `in` 演算子に渡すと TypeError になる回帰テスト。
        """
        pre = parse_md("```\nplain text\n```").find("pre")
        assert pre.find("code").get("class") is None  # 前提条件の確認

        process_code_or_mermaid(gen_with_slide, pre)

        runs = textboxes_of(gen_with_slide.current_slide)[0].text_frame.paragraphs[0].runs
        assert "plain text" in "".join(r.text for r in runs)

    def test_language_is_extracted(self, gen_with_slide):
        """language-xxx クラスから言語名が取り出される"""
        with patch("processors.append_code_textbox") as mock_append:
            process_code_or_mermaid(gen_with_slide, parse_md("```python\nx=1\n```").find("pre"))
        assert mock_append.call_args[1]["language"] == "python"

    def test_code_block_is_rendered(self, gen_with_slide):
        """通常のコードブロックは背景色付きテキストボックスになる"""
        process_code_or_mermaid(gen_with_slide, parse_md("```python\nprint(1)\n```").find("pre"))

        assert len(textboxes_of(gen_with_slide.current_slide)) == 1
        assert gen_with_slide.slide_has_text is True

    @patch("processors.insert_image_fit")
    @patch("requests.get")
    def test_mermaid_is_rendered_as_image(
        self, mock_get, mock_fit, gen_with_slide, mock_response
    ):
        """Mermaid記法は画像として挿入される"""
        mock_get.return_value = mock_response
        md = "```mermaid\ngraph TD; A-->B;\n```"
        process_code_or_mermaid(gen_with_slide, parse_md(md).find("pre"))

        assert mock_get.call_args[0][0].startswith(mermaid_renderer.DEFAULT_KROKI_ENDPOINT)
        assert mock_get.call_args[1]["timeout"] == mermaid_renderer.HTTP_TIMEOUT_SEC
        mock_fit.assert_called_once()

    @patch("processors.insert_image_fit")
    @patch("requests.get")
    def test_mermaid_two_column_layout(self, mock_get, mock_fit, gen_with_slide, mock_response):
        """テキストがあるスライドではMermaid図を右側に配置する"""
        mock_get.return_value = mock_response
        gen_with_slide.slide_has_text = True
        process_code_or_mermaid(
            gen_with_slide, parse_md("```mermaid\ngraph TD; A-->B;\n```").find("pre")
        )
        assert mock_fit.call_args[0][2] == gen_with_slide.layout.split_right_left

    @patch("processors.insert_image_fit")
    @patch("requests.get")
    def test_falls_back_to_mermaid_ink(
        self, mock_get, mock_fit, gen_with_slide, mock_response, capsys
    ):
        """Krokiが失敗した場合はmermaid.inkにフォールバックする"""
        mock_get.side_effect = [OSError("kroki down"), mock_response]
        process_code_or_mermaid(
            gen_with_slide, parse_md("```mermaid\ngraph TD; A-->B;\n```").find("pre")
        )

        assert mock_get.call_count == 2
        assert mock_get.call_args[0][0].startswith(mermaid_renderer.MERMAID_INK_URL)
        assert "代替API" in capsys.readouterr().out
        mock_fit.assert_called_once()

    @patch("processors.place_image")
    def test_renderer_off_places_nothing(self, mock_place, base_config, capsys):
        """mermaid.renderer が off の場合、図を配置せずスライドはそのまま進む"""
        base_config["mermaid"] = {"renderer": "off"}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_code_or_mermaid(gen, parse_md("```mermaid\ngraph TD; A-->B;\n```").find("pre"))

        mock_place.assert_not_called()
        assert "スキップ" in capsys.readouterr().out

    @patch("requests.get", side_effect=OSError("timeout"))
    def test_both_apis_failing_is_reported(self, _mock_get, gen_with_slide, capsys):
        """両方のAPIが失敗しても警告のみ表示して処理を継続する"""
        process_code_or_mermaid(
            gen_with_slide, parse_md("```mermaid\ngraph TD; A-->B;\n```").find("pre")
        )
        assert "Mermaid図形の生成に失敗しました" in capsys.readouterr().out


class TestOrderedList:
    """番号付きリスト（ol）と箇条書き（ul）の区別"""

    AUTO_NUM = "buAutoNum"

    def _marks(self, gen):
        return [
            self.AUTO_NUM in p._element.xml for p in gen.current_body.paragraphs if p.text.strip()
        ]

    def test_ordered_list_gets_numbering(self, gen_with_slide):
        """番号付きリストにはPowerPointの自動採番を設定する

        従来は ul と ol で見た目が変わらず、手順書として読めなかった。
        """
        for item in parse_md("1. 手順1\n2. 手順2\n").find_all("li"):
            process_text(gen_with_slide, item)

        assert self._marks(gen_with_slide) == [True, True]

    def test_unordered_list_is_unchanged(self, gen_with_slide):
        """箇条書きは従来どおり（採番しない）"""
        for item in parse_md("* 項目A\n* 項目B\n").find_all("li"):
            process_text(gen_with_slide, item)

        assert self._marks(gen_with_slide) == [False, False]

    def test_mixed_lists(self, gen, tmp_path):
        """同じスライド内で箇条書きと番号付きリストが混在しても区別される"""
        gen.generate("## 混在\n\n* 箇条書き\n\n1. 手順\n", str(tmp_path / "out.pptx"))
        assert self._marks(gen) == [False, True]

    def test_nested_ordered_list(self, gen, tmp_path):
        """入れ子の番号付きリストにも採番される"""
        gen.generate("## 手順\n\n1. 手順1\n    1. 手順1-1\n", str(tmp_path / "out.pptx"))

        assert self._marks(gen) == [True, True]
        assert [p.level for p in gen.current_body.paragraphs] == [0, 1]

    def test_ordered_item_inside_unordered_list(self, gen, tmp_path):
        """直近の親がulなら、外側にolがあっても採番しない"""
        gen.generate("## 混在\n\n1. 手順\n    * 補足\n", str(tmp_path / "out.pptx"))
        assert self._marks(gen) == [True, False]

    @pytest.mark.parametrize(
        "html, expected",
        [
            ("<ol><li>a</li></ol>", True),
            ("<ul><li>a</li></ul>", False),
            ("<ol><li><ul><li>a</li></ul></li></ol>", False),  # 直近の親はul
        ],
    )
    def test_is_ordered_item(self, html, expected):
        item = parse_html(html).find_all("li")[-1]
        assert processors.is_ordered_item(item) is expected

    def test_ordered_font_can_be_configured(self, base_config):
        """ordered_level_N で番号付きリストだけ書式を変えられる"""
        base_config["fonts"]["bullet_level_1"] = {"size_pt": 18}
        base_config["fonts"]["ordered_level_1"] = {"size_pt": 14}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_text(gen, parse_md("* 箇条書き").find("li"))
        process_text(gen, parse_md("1. 手順").find("li"))

        sizes = [p.runs[0].font.size for p in gen.current_body.paragraphs]
        assert sizes == [Pt(18), Pt(14)]

    def test_ordered_falls_back_to_bullet_config(self, base_config):
        """ordered_level_N が無ければ bullet_level_N を使う（設定追加は任意）"""
        base_config["fonts"]["bullet_level_1"] = {"size_pt": 18}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_text(gen, parse_md("1. 手順").find("li"))
        assert gen.current_body.paragraphs[0].runs[0].font.size == Pt(18)


class TestProcessText:
    def test_paragraph_uses_body_font(self, gen_with_slide):
        """段落にはbodyのフォント設定が適用される"""
        process_text(gen_with_slide, parse_md("本文です").find("p"))

        run = gen_with_slide.current_body.paragraphs[0].runs[0]
        assert run.text == "本文です"
        assert run.font.size == Pt(20)
        assert gen_with_slide.slide_has_text is True

    def test_empty_text_is_skipped(self, gen_with_slide):
        """空の要素は段落を追加しない"""
        process_text(gen_with_slide, parse_html("<p>   </p>").p)

        assert gen_with_slide.slide_has_text is False
        assert len(gen_with_slide.current_body.paragraphs) == 1

    def test_bullet_level_follows_nesting(self, gen_with_slide):
        """ネストの深さがリストのレベルになる"""
        for item in parse_md("* 親\n    * 子\n        * 孫").find_all("li"):
            process_text(gen_with_slide, item)

        assert [p.level for p in gen_with_slide.current_body.paragraphs] == [0, 1, 2]

    def test_bullet_level_is_capped(self, gen_with_slide):
        """レベルはPowerPointの上限(8)でクランプされる"""
        html = "<ul><li>深い</li></ul>"
        for _ in range(12):
            html = f"<ul><li>{html}</li></ul>"

        process_text(gen_with_slide, parse_html(html).find_all("li")[-1])
        assert gen_with_slide.current_body.paragraphs[0].level == 8

    def test_bullet_level_font_config(self, base_config):
        """bullet_level_N の設定がレベルごとに適用される"""
        base_config["fonts"]["bullet_level_1"] = {"size_pt": 18}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_text(gen, parse_md("* 項目").find("li"))
        assert gen.current_body.paragraphs[0].runs[0].font.size == Pt(18)

    def test_nested_bullets_inherit_the_shallower_level(self, base_config):
        """未定義の階層は、より浅い階層の設定を引き継ぐ

        body へフォールバックしていた頃は、ネストした途端に
        書体が bullet_level_1 の指定から body の指定へ変わっていた。
        """
        base_config["fonts"]["bullet_level_1"] = {"name": "Yu Gothic", "size_pt": 18}
        base_config["fonts"]["body"] = {"name": "Meiryo", "size_pt": 24}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        for item in parse_md("* 親\n    * 子\n        * 孫").find_all("li"):
            process_text(gen, item)

        fonts = [(p.runs[0].font.name, p.runs[0].font.size) for p in gen.current_body.paragraphs]
        assert fonts == [("Yu Gothic", Pt(18))] * 3

    def test_explicit_deeper_level_wins(self, base_config):
        """深い階層に個別設定があればそちらが優先される"""
        base_config["fonts"]["bullet_level_1"] = {"size_pt": 18}
        base_config["fonts"]["bullet_level_2"] = {"size_pt": 14}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        for item in parse_md("* 親\n    * 子").find_all("li"):
            process_text(gen, item)

        sizes = [p.runs[0].font.size for p in gen.current_body.paragraphs]
        assert sizes == [Pt(18), Pt(14)]

    def test_bullet_falls_back_to_body_without_any_config(self, base_config):
        """箇条書きの設定が一切無い場合は body を使う"""
        base_config["fonts"].pop("bullet_level_1", None)
        base_config["fonts"]["body"] = {"size_pt": 22}
        gen = PPTXGenerator(base_config)
        process_heading(gen, parse_md("## 見出し").find("h2"))

        process_text(gen, parse_md("* 項目").find("li"))
        assert gen.current_body.paragraphs[0].runs[0].font.size == Pt(22)


# =====================================================================
# text_metrics.py
# =====================================================================


class TestCharWidth:
    @pytest.mark.parametrize("ch", ["あ", "漢", "全", "ー", "、"])
    def test_fullwidth_characters(self, ch):
        """日本語は全角幅として扱う"""
        assert char_width_ratio(ch) == 1.0

    @pytest.mark.parametrize("ch", ["a", "Z", "1", " ", "-"])
    def test_halfwidth_characters(self, ch):
        """英数字・記号は半角幅として扱う"""
        assert char_width_ratio(ch) == 0.5

    def test_width_of_mixed_text(self):
        """全角と半角が混在してもそれぞれの幅で合算される"""
        # 全角2文字(=2.0) + 半角4文字(=2.0) → 合計4.0文字ぶん
        assert estimate_text_width_pt("あいabcd", 10) == pytest.approx(40.0)


class TestEstimateLineCount:
    def test_text_wraps_by_available_width(self):
        """枠幅を超えるテキストは折り返して複数行になる"""
        # 全角20文字 × 10pt = 200pt を 100pt 幅に入れる → 2行
        assert estimate_line_count("あ" * 20, 10, 100) == 2

    def test_short_text_is_single_line(self):
        assert estimate_line_count("短い", 10, 500) == 1

    @pytest.mark.parametrize("text, width", [("", 100), ("あ", 0), ("あ", -10)])
    def test_degenerate_cases_are_one_line(self, text, width):
        """空文字や幅が取れない場合も1行として扱う（0除算を避ける）"""
        assert estimate_line_count(text, 10, width) == 1


class TestEstimateHeight:
    def _para(self, text, size=10, level=0, spacing=1.0, space_after=0.0, space_before=0.0):
        return ParagraphMetrics(text, size, level, spacing, space_after, space_before)

    def test_single_line_height(self):
        """1行の高さは フォントサイズ × フォントの行高 × 行送り"""
        # 10pt × 1.3(行高) × 1.2(行送り) = 15.6pt
        assert estimate_height_pt([self._para("あ", spacing=1.2)], 500) == pytest.approx(15.6)

    def test_line_height_exceeds_font_size(self):
        """行の高さはフォントサイズそのものではなく、フォントの行高を考慮する

        ここを 1.0 とみなしていた頃は必要な高さを約3割過小評価し、
        本文がスライドからはみ出していた。
        """
        height = estimate_height_pt([self._para("あ", size=10, spacing=1.0)], 500)
        assert height > 10
        assert height == pytest.approx(10 * text_metrics.LINE_HEIGHT_RATIO)

    def test_wrapped_lines_are_counted(self):
        """折り返した行数ぶんの高さになる"""
        # 全角30文字 × 10pt = 300pt を 100pt 幅 → 3行
        assert estimate_height_pt([self._para("あ" * 30)], 100) == pytest.approx(39.0)

    def test_space_after_is_included(self):
        assert estimate_height_pt([self._para("あ", space_after=5.0)], 500) == pytest.approx(18.0)

    def test_space_before_is_included(self):
        """段落前の余白（h3見出しなどが使う）も高さに含める"""
        assert estimate_height_pt([self._para("あ", space_before=10.0)], 500) == pytest.approx(23.0)

    def test_indent_reduces_available_width(self):
        """箇条書きのレベルが深いほど幅が狭まり、行数が増える"""
        flat = estimate_height_pt([self._para("あ" * 30, level=0)], 300)
        nested = estimate_height_pt([self._para("あ" * 30, level=3)], 300)
        assert nested > flat

    def test_scale_reduces_height(self):
        """縮小率を掛けると必要な高さが下がる"""
        paragraphs = [self._para("あ" * 30)]
        assert estimate_height_pt(paragraphs, 100, scale=0.5) < estimate_height_pt(
            paragraphs, 100
        )


class TestTableHeightEstimation:
    """表の高さの概算と分割"""

    def _rows(self, count, text="項目", size=12):
        return [TableRowMetrics([text, text], size) for _ in range(count)]

    def test_row_height_includes_margins(self):
        """行の高さはフォントの行高にセルの上下マージンを加えたもの"""
        (height,) = estimate_row_heights_pt(self._rows(1, size=10), 500, 7.2)
        assert height == pytest.approx(10 * text_metrics.LINE_HEIGHT_RATIO + 7.2)

    def test_wrapped_cell_makes_the_row_taller(self):
        """折り返すセルがある行はその分高くなる"""
        narrow = estimate_row_heights_pt([TableRowMetrics(["あ" * 30], 10)], 100, 0)
        wide = estimate_row_heights_pt([TableRowMetrics(["あ" * 30], 10)], 500, 0)
        assert narrow[0] > wide[0]

    def test_total_height_is_the_sum_of_rows(self):
        rows = self._rows(4)
        assert estimate_table_height_pt(rows, 500, 7.2) == pytest.approx(
            sum(estimate_row_heights_pt(rows, 500, 7.2))
        )

    def test_scale_shrinks_a_tall_table(self):
        """1枚に収まらない表は縮小率が1.0未満になる"""
        assert fit_table_scale(self._rows(30), 500, 7.2, 200) < 1.0

    def test_no_shrink_when_it_fits(self):
        assert fit_table_scale(self._rows(2), 500, 7.2, 1000) == 1.0

    def test_empty_table_needs_no_shrink(self):
        assert fit_table_scale([], 500, 7.2, 100) == 1.0


class TestPaginateRows:
    """表の行をページに分割する"""

    def test_single_page_when_it_fits(self):
        assert paginate_row_heights([10, 10, 10], 100) == [[0, 1, 2]]

    def test_splits_when_overflowing(self):
        assert paginate_row_heights([10, 10, 10, 10], 25) == [[0, 1], [2, 3]]

    def test_header_is_repeated_on_each_page(self):
        """見出し行は各ページの先頭に繰り返される"""
        pages = paginate_row_heights([10, 10, 10, 10, 10], 25, repeat_first_row=True)

        assert pages == [[0, 1], [0, 2], [0, 3], [0, 4]]
        assert all(page[0] == 0 for page in pages)

    def test_header_height_is_reserved_on_every_page(self):
        """繰り返す見出し行の高さも各ページの消費として数える

        枠20pt・行10pt の場合、見出しを繰り返すとデータは1ページ1行しか入らない。
        """
        assert paginate_row_heights([10, 10, 10, 10], 20) == [[0, 1], [2, 3]]
        assert paginate_row_heights([10, 10, 10, 10], 20, repeat_first_row=True) == [
            [0, 1], [0, 2], [0, 3]
        ]

    def test_oversized_row_still_gets_a_page(self):
        """1行だけで枠を超える場合でも、そのページに載せる（無限分割の防止）"""
        assert paginate_row_heights([100, 100], 10) == [[0], [1]]

    @pytest.mark.parametrize("heights, available", [([], 100), ([10], 0)])
    def test_degenerate_cases(self, heights, available):
        """行が無い、または高さが取れない場合も破綻しない"""
        pages = paginate_row_heights(heights, available)
        assert pages == ([] if not heights else [[0]])


class TestFitScale:
    def _paras(self, count, text="あ" * 20, size=20):
        return [ParagraphMetrics(text, size) for _ in range(count)]

    def test_no_shrink_when_it_fits(self):
        """収まっている場合は縮小しない"""
        assert fit_scale(self._paras(1), 500, 1000) == 1.0

    def test_shrinks_when_overflowing(self):
        """収まらない場合は1.0未満の縮小率を返す"""
        scale = fit_scale(self._paras(20), 500, 200)
        assert scale < 1.0

    def test_result_actually_fits(self):
        """返した縮小率で実際に枠へ収まる"""
        paragraphs = self._paras(10)
        scale = fit_scale(paragraphs, 500, 300)
        assert estimate_height_pt(paragraphs, 500, scale) <= 300

    def test_scale_has_a_floor(self):
        """極端に多い場合でも下限より小さくはしない（読めなくなるため）"""
        assert fit_scale(self._paras(500), 500, 50) == text_metrics.MIN_SHRINK_SCALE

    @pytest.mark.parametrize("paragraphs, height", [([], 100), ([ParagraphMetrics("あ", 10)], 0)])
    def test_degenerate_cases_do_not_shrink(self, paragraphs, height):
        """段落が無い・高さが取れない場合は縮小しない"""
        assert fit_scale(paragraphs, 500, height) == 1.0


# =====================================================================
# layout.py
# =====================================================================


class TestSlideLayout:
    """配置寸法の導出

    絶対値そのものより、要素どうしが揃っているかを検証する。
    寸法の微調整でテストが壊れず、崩れたときには確実に落ちるようにするため。
    """

    def _layout(self, width_in=10, height_in=5.625, **body):
        return SlideLayout(Inches(width_in), Inches(height_in), **body)

    def test_content_area_follows_the_body_placeholder(self):
        """コンテンツ領域は本文プレースホルダーに一致する

        画像・表だけが別の余白で配置されると左端が揃わないため、
        テンプレートが定めた本文の位置をそのまま基準にする。
        """
        layout = self._layout(
            body_left=Inches(0.5), body_top=Inches(1.75), body_width=Inches(9.0)
        )
        assert layout.content_left == Inches(0.5)
        assert layout.content_top == Inches(1.75)
        assert layout.content_width == Inches(9.0)

    def test_falls_back_without_a_placeholder(self):
        """プレースホルダーが無いテンプレートでは既定の余白を使う"""
        layout = self._layout()
        assert layout.content_left == Inches(0.5)
        assert layout.content_width == Inches(9.0)  # 10 - 0.5×2

    def test_built_from_presentation(self, base_config):
        """ジェネレーターは実際のテンプレートから寸法を取り込む"""
        gen = PPTXGenerator(base_config)
        body = gen.prs.slide_layouts[1].placeholders[1]

        assert gen.layout.content_left == body.left
        assert gen.layout.content_top == body.top
        assert gen.layout.content_width == body.width

    def test_template_without_body_placeholder(self):
        """本文プレースホルダーを持たないテンプレートでは既定値にフォールバックする"""
        prs = MagicMock()
        prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
        prs.slide_layouts.__getitem__.side_effect = KeyError("placeholder なし")

        layout = SlideLayout.from_presentation(prs)

        assert layout.content_left == Inches(0.5)
        assert layout.content_top == Inches(1.75)

    # --- 整合性 ---

    def test_full_width_and_split_share_the_right_edge(self):
        """全幅の要素と2カラム時の要素で右端が揃う

        従来は全幅が右余白1.0インチ、2カラム時の図が0.3インチと食い違っていた。
        """
        layout = self._layout(
            body_left=Inches(0.5), body_top=Inches(1.75), body_width=Inches(9.0)
        )
        assert layout.split_right_left + layout.split_right_width == (
            layout.content_left + layout.content_width
        )

    def test_split_columns_do_not_overlap(self):
        """2カラムの本文と図が重ならず、間隔が空く"""
        layout = self._layout(
            body_left=Inches(0.5), body_top=Inches(1.75), body_width=Inches(9.0)
        )
        body_right = layout.content_left + layout.split_body_width
        assert body_right < layout.split_right_left
        assert layout.split_right_left - body_right == Inches(0.3)

    def test_table_does_not_overlap_the_body(self):
        """表と本文枠が重ならない

        従来は本文枠の下端が表の上端より約0.95インチ下にあり、矩形が重なっていた。
        """
        layout = self._layout(
            body_left=Inches(0.5), body_top=Inches(1.75), body_width=Inches(9.0)
        )
        body_bottom = layout.content_top + layout.table_split_body_height
        assert body_bottom < layout.table_split_top

    def test_centered_box_is_symmetric(self):
        """中央寄せの枠は左右の余白が等しい"""
        layout = self._layout(
            body_left=Inches(0.5), body_top=Inches(1.75), body_width=Inches(9.0)
        )
        left_margin = layout.center_left - layout.content_left
        right_margin = (layout.content_left + layout.content_width) - (
            layout.center_left + layout.center_width
        )
        assert left_margin == pytest.approx(right_margin, abs=2)

    # --- スライドサイズへの追従 ---

    @pytest.mark.parametrize(
        "width_in, height_in",
        [(10, 5.625), (10, 7.5), (10, 6.25), (11.69, 8.27)],  # 16:9 / 4:3 / 16:10 / A4
    )
    def test_content_height_follows_slide_size(self, width_in, height_in):
        """コンテンツの高さはスライドサイズに追従し、下の余白は一定になる"""
        layout = self._layout(width_in, height_in, body_top=Inches(1.75))
        assert Emu(layout.content_height).inches == pytest.approx(
            height_in - 1.75 - 0.5, abs=0.001
        )

    @pytest.mark.parametrize("width_in", [10, 11.69, 13.333])
    def test_split_follows_slide_width(self, width_in):
        """画角が変わっても2カラムの右端はコンテンツ領域の右端に揃う"""
        layout = self._layout(width_in, body_width=Inches(width_in - 1.0))
        assert layout.split_right_left + layout.split_right_width == (
            layout.content_left + layout.content_width
        )

    def test_body_height_fits_in_slide(self):
        """本文枠の高さはスライド下端に余白を残す"""
        layout = self._layout(10, 7.5)
        assert Emu(layout.body_height_for(Inches(1.75))).inches == pytest.approx(5.25)



class TestLayoutAdaptsToSlideSize:
    """画角を変えたときに図・表が追従することの結合テスト"""

    MD = "## 図のスライド\n\n![img]({png})\n\n## 表のスライド\n\n| A | B |\n|---|---|\n| 1 | 2 |\n"

    def _generate(self, base_config, layout_name, tmp_path, png_file):
        base_config["slides"]["layout"] = layout_name
        gen = PPTXGenerator(base_config)
        gen.generate(self.MD.format(png=png_file), str(tmp_path / "out.pptx"))
        return gen

    @pytest.mark.parametrize("layout_name", ["16:9", "4:3", "A4"])
    def test_content_stays_inside_the_slide(self, base_config, layout_name, tmp_path, png_file):
        """どの画角でも図・表がスライドからはみ出さない"""
        gen = self._generate(base_config, layout_name, tmp_path, png_file)

        for slide in gen.prs.slides:
            for shape in slide.shapes:
                assert shape.left + shape.width <= gen.prs.slide_width
                assert shape.top + shape.height <= gen.prs.slide_height

    def test_wide_slide_uses_the_extra_width(self, base_config, tmp_path, png_file):
        """A4のような横長の画角では、表が広がった幅を使う"""
        narrow = self._generate(base_config, "16:9", tmp_path, png_file)
        wide = self._generate(base_config, "A4", tmp_path, png_file)

        def table_width(gen):
            for slide in gen.prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        return shape.width
            raise AssertionError("表が見つかりません")

        assert table_width(wide) > table_width(narrow)


# =====================================================================
# mermaid_renderer.py
# =====================================================================


class TestMermaidConfig:
    def test_conf_is_read_from_config(self, base_config):
        """config.yaml の mermaid セクションを読み取る"""
        base_config["mermaid"] = {"renderer": "off"}
        assert mermaid_conf(PPTXGenerator(base_config)) == {"renderer": "off"}

    def test_missing_section_is_empty(self, gen):
        """mermaid セクションが無い場合は空の設定になる"""
        assert mermaid_conf(gen) == {}

    @pytest.mark.parametrize(
        "endpoint, expected",
        [
            ("https://kroki.io", True),
            ("https://mermaid.ink/img/", True),
            ("http://kroki.internal.example.com", False),
            ("http://localhost:8000", False),
        ],
    )
    def test_public_endpoint_detection(self, endpoint, expected):
        """公開サービスかどうかを判定できる（警告とフォールバック可否の基準）"""
        assert mermaid_renderer.is_public_endpoint(endpoint) is expected


class TestRenderMermaid:
    MMD = "graph TD; A-->B;"

    def test_off_skips_generation(self, capsys):
        """renderer: off では図を生成せず None を返す"""
        assert render_mermaid({"renderer": "off"}, self.MMD) is None
        assert "スキップ" in capsys.readouterr().out

    def test_unknown_renderer_raises(self):
        """未知のレンダラー名はエラーにする（設定ミスを黙って無視しない）"""
        with pytest.raises(MermaidRenderError, match="未知の mermaid.renderer"):
            render_mermaid({"renderer": "magic"}, self.MMD)

    @patch("requests.get")
    def test_default_uses_public_kroki(self, mock_get, mock_response):
        """既定では公開Krokiを使う（従来どおりの動作）"""
        mock_get.return_value = mock_response
        assert render_mermaid({}, self.MMD) == TINY_PNG

        url = mock_get.call_args[0][0]
        assert url.startswith(f"{mermaid_renderer.DEFAULT_KROKI_ENDPOINT}/mermaid/png/")
        assert mock_get.call_args[1]["timeout"] == mermaid_renderer.HTTP_TIMEOUT_SEC

    @patch("requests.get")
    def test_warns_before_sending_to_public_service(self, mock_get, mock_response, capsys):
        """公開サービスへ送信する前に警告を表示する"""
        mock_get.return_value = mock_response
        render_mermaid({}, self.MMD)

        out = capsys.readouterr().out
        assert "外部サービス" in out
        assert "kroki.io" in out

    @patch("requests.get")
    def test_warning_can_be_disabled(self, mock_get, mock_response, capsys):
        """警告は warn_on_external: false で抑制できる"""
        mock_get.return_value = mock_response
        render_mermaid({"warn_on_external": False}, self.MMD)
        assert "外部サービス" not in capsys.readouterr().out

    @patch("requests.get")
    def test_self_hosted_endpoint_is_used(self, mock_get, mock_response, capsys):
        """自己ホストのKrokiを指定でき、その場合は警告を出さない"""
        mock_get.return_value = mock_response
        render_mermaid({"endpoint": "http://kroki.internal/"}, self.MMD)

        assert mock_get.call_args[0][0].startswith("http://kroki.internal/mermaid/png/")
        assert "外部サービス" not in capsys.readouterr().out

    @patch("requests.get")
    def test_public_endpoint_falls_back_to_mermaid_ink(self, mock_get, mock_response):
        """公開Kroki利用時は従来どおり mermaid.ink にフォールバックする"""
        mock_get.side_effect = [OSError("kroki down"), mock_response]

        assert render_mermaid({}, self.MMD) == TINY_PNG
        assert mock_get.call_args[0][0].startswith(mermaid_renderer.MERMAID_INK_URL)

    @patch("requests.get")
    def test_self_hosted_never_falls_back_to_public(self, mock_get):
        """自己ホスト指定時は公開APIへフォールバックしない（情報漏洩の防止）

        社内Krokiが落ちた際に、機密を含み得る図が公開サービスへ送られるのを防ぐ。
        """
        mock_get.side_effect = OSError("社内Krokiが停止")

        with pytest.raises(MermaidRenderError, match="情報漏洩"):
            render_mermaid({"endpoint": "http://kroki.internal"}, self.MMD)

        assert mock_get.call_count == 1  # mermaid.ink は呼ばれない

    @patch("requests.get")
    def test_fallback_can_be_forced_for_self_hosted(self, mock_get, mock_response):
        """明示的に許可した場合のみ、自己ホストでもフォールバックする"""
        mock_get.side_effect = [OSError("down"), mock_response]

        result = render_mermaid(
            {"endpoint": "http://kroki.internal", "fallback_to_public": True}, self.MMD
        )
        assert result == TINY_PNG
        assert mock_get.call_count == 2

    @patch("requests.get")
    def test_fallback_can_be_disabled_for_public(self, mock_get):
        """公開Kroki利用時でもフォールバックを禁止できる"""
        mock_get.side_effect = OSError("down")

        with pytest.raises(MermaidRenderError):
            render_mermaid({"fallback_to_public": False}, self.MMD)
        assert mock_get.call_count == 1


class TestLocalMermaidRendering:
    """mermaid-cli によるオフライン生成（外部送信なし）"""

    MMD = "graph TD; A-->B;"
    CONF = {"renderer": "local"}

    def _fake_run(self, png_bytes=TINY_PNG, returncode=0):
        """mmdc の代わりに出力ファイルを書くダミー"""
        def run(command, **kwargs):
            if returncode == 0:
                output_path = command[command.index("-o") + 1]
                with open(output_path, "wb") as f:
                    f.write(png_bytes)
            stderr = "パースに失敗しました".encode("utf-8")
            return subprocess.CompletedProcess(command, returncode, b"", stderr)
        return run

    def test_renders_without_network(self, mocker):
        """ローカル生成ではネットワークを一切使わない"""
        mocker.patch("subprocess.run", side_effect=self._fake_run())
        mock_get = mocker.patch("requests.get")

        assert render_mermaid(self.CONF, self.MMD) == TINY_PNG
        mock_get.assert_not_called()

    def test_diagram_source_is_passed_to_cli(self, mocker):
        """図のソースを一時ファイル経由でCLIへ渡す"""
        captured = {}

        def run(command, **kwargs):
            source_path = command[command.index("-i") + 1]
            with open(source_path, encoding="utf-8") as f:
                captured["source"] = f.read()
            with open(command[command.index("-o") + 1], "wb") as f:
                f.write(TINY_PNG)
            return subprocess.CompletedProcess(command, 0, b"", b"")

        mocker.patch("subprocess.run", side_effect=run)
        render_mermaid(self.CONF, self.MMD)
        assert captured["source"] == self.MMD

    def test_custom_cli_path(self, mocker):
        """cli_path で実行コマンドを差し替えられる"""
        mock_run = mocker.patch("subprocess.run", side_effect=self._fake_run())
        render_mermaid({"renderer": "local", "cli_path": "/opt/bin/mmdc"}, self.MMD)
        assert mock_run.call_args[0][0][0] == "/opt/bin/mmdc"

    def test_missing_cli_is_reported(self, mocker):
        """mermaid-cli が未インストールの場合は導入方法を案内する"""
        mocker.patch("subprocess.run", side_effect=FileNotFoundError)

        with pytest.raises(MermaidRenderError, match="mermaid-cli"):
            render_mermaid(self.CONF, self.MMD)

    def test_cli_failure_is_reported(self, mocker):
        """CLIが異常終了した場合は標準エラー出力を添えて報告する"""
        mocker.patch("subprocess.run", side_effect=self._fake_run(returncode=1))

        with pytest.raises(MermaidRenderError, match="失敗しました"):
            render_mermaid(self.CONF, self.MMD)

    def test_timeout_is_reported(self, mocker):
        """タイムアウトした場合も分かるメッセージにする"""
        mocker.patch(
            "subprocess.run", side_effect=subprocess.TimeoutExpired("mmdc", 60)
        )

        with pytest.raises(MermaidRenderError, match="タイムアウト"):
            render_mermaid(self.CONF, self.MMD)

    def test_missing_output_is_reported(self, mocker):
        """CLIが正常終了しても画像が無い場合はエラーにする"""
        mocker.patch(
            "subprocess.run",
            return_value=subprocess.CompletedProcess(["mmdc"], 0, b"", b""),
        )

        with pytest.raises(MermaidRenderError, match="画像を出力しませんでした"):
            render_mermaid(self.CONF, self.MMD)


# =====================================================================
# generator.py
# =====================================================================


class TestInitialization:
    @pytest.mark.parametrize(
        "layout, expected",
        [
            ("16:9", (Inches(10), Inches(5.625))),
            ("4:3", (Inches(10), Inches(7.5))),
            ("16:10", (Inches(10), Inches(6.25))),
            ("A4", (Inches(11.69), Inches(8.27))),
            ("unknown", (Inches(10), Inches(5.625))),
        ],
    )
    def test_get_slide_size(self, gen, layout, expected):
        """スライドサイズの計算ロジック（未知の値は16:9にフォールバック）"""
        assert gen._get_slide_size(layout) == expected

    def test_slide_size_applied_from_config(self, base_config):
        """config.yamlの画角がプレゼンテーションに反映される"""
        base_config["slides"]["layout"] = "4:3"
        gen = PPTXGenerator(base_config)
        assert (gen.prs.slide_width, gen.prs.slide_height) == (Inches(10), Inches(7.5))

    @pytest.mark.parametrize(
        "config", [None, {}, {"slides": None, "fonts": None, "images": None}]
    )
    def test_empty_config_is_accepted(self, config):
        """設定が空・Noneでも既定値で初期化できる"""
        assert PPTXGenerator(config).prs.slide_height == Inches(5.625)

    def test_template_is_loaded_when_exists(self, base_config, tmp_path):
        """テンプレートが存在する場合は読み込み、画角設定を上書きしない"""
        template = tmp_path / "template.pptx"
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        prs.save(str(template))

        base_config["slides"]["template_path"] = str(template)
        assert PPTXGenerator(base_config).prs.slide_width == Inches(13.333)

    def test_missing_template_falls_back_to_default(self, base_config):
        """テンプレートのパスが不正な場合はデフォルトのプレゼンテーションを使う"""
        base_config["slides"]["template_path"] = "no_such_template.pptx"
        assert PPTXGenerator(base_config).prs.slide_height == Inches(5.625)


class TestFrontMatter:
    FM = (
        '---\ntitle: "資料タイトル"\nsubtitle: "サブ"\n'
        'author: "著者"\ndate: "2026-05-10"\n---\n\n## 中身\n本文\n'
    )

    def test_title_slide_is_generated(self, gen, tmp_path):
        """フロントマターからタイトルスライドが自動生成される"""
        gen.generate(self.FM, str(tmp_path / "out.pptx"))

        title_slide = gen.prs.slides[0]
        assert title_slide.shapes.title.text == "資料タイトル"
        assert title_slide.placeholders[1].text == "サブ\n著者\n2026-05-10"
        assert len(gen.prs.slides) == 2

    def test_title_font_is_applied(self, gen, tmp_path):
        """タイトルには title_h1 の設定が適用される"""
        gen.generate(self.FM, str(tmp_path / "out.pptx"))
        run = gen.prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(44)

    def test_broken_front_matter_is_reported(self, gen, tmp_path, capsys):
        """壊れたフロントマターは警告を出して処理を続行する"""
        gen.generate("---\ntitle: [壊れた\n---\n\n## 中身\n", str(tmp_path / "out.pptx"))
        assert "フロントマターの解析に失敗しました" in capsys.readouterr().out

    def test_front_matter_without_title_is_ignored(self, gen, tmp_path):
        """titleが無いフロントマターではタイトルスライドを作らない"""
        gen.generate("---\nauthor: 著者\n---\n\n## 中身\n", str(tmp_path / "out.pptx"))
        assert len(gen.prs.slides) == 1


class TestLayoutComments:
    @patch("processors.insert_image_fit")
    def test_two_column_comment(self, mock_fit, gen, tmp_path, png_file):
        """<!-- layout: 2-column --> で画像が右側に寄る"""
        md = f"## 見出し\n\n<!-- layout: 2-column -->\n\n![img]({png_file})\n"
        gen.generate(md, str(tmp_path / "out.pptx"))

        assert gen.forced_layout == "2-column"
        assert mock_fit.call_args[0][2] == gen.layout.split_right_left

    @patch("processors.insert_image_fit")
    def test_center_comment(self, mock_fit, gen, tmp_path, png_file):
        """<!-- layout: center --> で画像が中央に配置される"""
        md = f"## 見出し\n\n本文\n\n<!-- layout: center -->\n\n![img]({png_file})\n"
        gen.generate(md, str(tmp_path / "out.pptx"))
        assert mock_fit.call_args[0][2] == gen.layout.content_left

    def test_layout_is_reset_on_new_slide(self, gen, tmp_path):
        """新しい見出しでレイアウト指定はリセットされる"""
        gen.generate("## 1枚目\n\n<!-- layout: center -->\n\n## 2枚目\n", str(tmp_path / "out.pptx"))
        assert gen.forced_layout is None


class TestFooter:
    """日付・文言・ページ番号のフッター"""

    MD = "# 表紙\n\n## 2枚目\n\n## 3枚目\n"

    def _generate(self, base_config, tmp_path, footer=None, show_number=True):
        base_config["slides"]["show_slide_number"] = show_number
        if footer is not None:
            base_config["slides"]["footer"] = footer
        gen = PPTXGenerator(base_config)
        gen.generate(self.MD, str(tmp_path / "out.pptx"))
        return gen

    def _texts(self, slide):
        return [s.text_frame.text for s in textboxes_of(slide) if s.text_frame.text]

    def test_footer_parts_are_placed(self, base_config, tmp_path):
        """左に日付、中央に文言、右にページ番号を置く（PowerPointの慣習）"""
        gen = self._generate(
            base_config, tmp_path, {"text": "社外秘", "date": "2026年5月10日"}
        )
        boxes = sorted(textboxes_of(gen.prs.slides[1]), key=lambda s: s.left)

        assert [b.text_frame.text for b in boxes] == ["2026年5月10日", "社外秘", "1"]

    def test_date_true_uses_the_conversion_date(self, base_config, tmp_path):
        """date: true は変換した日付を表示する"""
        gen = self._generate(base_config, tmp_path, {"date": True})

        expected = datetime.now().strftime("%Y-%m-%d")
        assert expected in self._texts(gen.prs.slides[1])

    def test_title_slide_is_excluded_by_default(self, base_config, tmp_path):
        """既定では表紙にフッターを出さない"""
        gen = self._generate(base_config, tmp_path, {"text": "社外秘"})
        assert self._texts(gen.prs.slides[0]) == []

    def test_show_on_title(self, base_config, tmp_path):
        """show_on_title で表紙にも文言・日付を出せる"""
        gen = self._generate(
            base_config, tmp_path, {"text": "社外秘", "show_on_title": True}
        )
        assert self._texts(gen.prs.slides[0]) == ["社外秘"]

    def test_title_slide_has_no_page_number(self, base_config, tmp_path):
        """表紙にはページ番号を振らない（本編を1ページ目として数える）"""
        gen = self._generate(
            base_config, tmp_path, {"text": "社外秘", "show_on_title": True}
        )

        assert "0" not in self._texts(gen.prs.slides[0])
        assert "1" in self._texts(gen.prs.slides[1])

    def test_no_footer_settings_adds_nothing_extra(self, base_config, tmp_path):
        """フッター設定が無ければページ番号だけになる（従来どおり）"""
        gen = self._generate(base_config, tmp_path)
        assert self._texts(gen.prs.slides[1]) == ["1"]

    def test_everything_disabled(self, base_config, tmp_path):
        """番号もフッターも無効なら何も追加しない"""
        gen = self._generate(base_config, tmp_path, {}, show_number=False)
        assert self._texts(gen.prs.slides[1]) == []

    def test_font_can_be_configured(self, base_config, tmp_path):
        """fonts.footer で書式を変えられる"""
        base_config["fonts"]["footer"] = {"size_pt": 9, "color_rgb": [1, 2, 3]}
        gen = self._generate(base_config, tmp_path, {"text": "社外秘"})

        run = textboxes_of(gen.prs.slides[1])[0].text_frame.paragraphs[0].runs[0]
        assert run.font.size == Pt(9)
        assert run.font.color.rgb == RGBColor(1, 2, 3)

    @pytest.mark.parametrize(
        "setting, expected", [(True, None), ("2026年5月10日", "2026年5月10日"), (False, ""), (None, "")]
    )
    def test_footer_date_text(self, setting, expected):
        """日付の指定の解釈"""
        actual = processors.footer_date_text(setting)
        assert actual == (datetime.now().strftime("%Y-%m-%d") if expected is None else expected)


class TestDarkTheme:
    """<!-- layout: dark-theme --> による配色の切り替え"""

    def _generate(self, gen, tmp_path, md):
        gen.generate(md, str(tmp_path / "out.pptx"))
        return gen.prs.slides[0]

    def test_background_and_text_are_inverted(self, gen, tmp_path):
        """背景が暗くなり、タイトルと本文の文字が明るくなる"""
        slide = self._generate(
            gen, tmp_path, "## 暗いスライド\n<!-- layout: dark-theme -->\n\n本文\n"
        )

        assert slide.background.fill.fore_color.rgb == RGBColor(30, 30, 30)
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        assert run.font.color.rgb == RGBColor(240, 240, 240)

    def test_colors_can_be_configured(self, base_config, tmp_path):
        """theme で配色を変えられる"""
        base_config["theme"] = {
            "dark_background_color": [1, 2, 3], "dark_text_color": [4, 5, 6],
        }
        gen = PPTXGenerator(base_config)
        slide = self._generate(gen, tmp_path, "## 暗い\n<!-- layout: dark-theme -->\n\n本文\n")

        assert slide.background.fill.fore_color.rgb == RGBColor(1, 2, 3)

    def test_can_be_combined_with_a_placement(self, gen, tmp_path, png_file):
        """配置の指定と併用できる（配色とは別軸のため）"""
        md = f"## 併用\n<!-- layout: 2-column, dark-theme -->\n\n本文\n\n![図]({png_file})\n"
        slide = self._generate(gen, tmp_path, md)

        assert gen.forced_layout == "2-column"
        assert slide.background.fill.fore_color.rgb == RGBColor(30, 30, 30)

    def test_applies_only_to_the_marked_slide(self, gen, tmp_path):
        """指定したスライドだけが対象になる"""
        gen.generate(
            "## 暗い\n<!-- layout: dark-theme -->\n\n本文\n\n## 通常\n本文\n",
            str(tmp_path / "out.pptx"),
        )

        assert gen.prs.slides[0].background.fill.type is not None
        assert gen.dark_slide is False  # 次のスライドでリセットされる

    def test_unknown_directive_warns(self, gen, tmp_path, capsys):
        """未知の指定は警告して無視する"""
        gen.generate("## 見出し\n<!-- layout: fancy-mode -->\n\n本文\n", str(tmp_path / "o.pptx"))
        assert "不明なレイアウト指定" in capsys.readouterr().out

    def test_empty_directive_is_ignored(self, gen, tmp_path, capsys):
        """区切りだけの指定（余分なカンマ）は黙って読み飛ばす"""
        gen.generate("## 見出し\n<!-- layout: center, -->\n\n本文\n", str(tmp_path / "o.pptx"))

        assert gen.forced_layout == "center"
        assert "不明なレイアウト指定" not in capsys.readouterr().out

    def test_without_a_slide_is_noop(self, gen):
        """スライドが無い状態でダークテーマを適用しても落ちない"""
        processors.apply_dark_theme(gen)


class TestFullImageLayout:
    """<!-- layout: full-image --> による全面配置"""

    def test_image_covers_the_slide(self, gen, tmp_path, png_file):
        """余白を取らずスライドいっぱいに広げる"""
        gen.generate(
            f"## 全面\n<!-- layout: full-image -->\n\n![図]({png_file})\n",
            str(tmp_path / "out.pptx"),
        )
        picture = [
            s for s in gen.prs.slides[0].shapes
            if s.shape_type is not None and "PICTURE" in str(s.shape_type)
        ][0]

        # 縦横比を保つため、片方の軸はスライド全体に一致する
        fills_width = picture.width == gen.prs.slide_width
        fills_height = picture.height == gen.prs.slide_height
        assert fills_width or fills_height
        assert picture.left >= 0 and picture.top >= 0

    def test_multiple_images_share_the_whole_slide(self, gen, tmp_path, png_file):
        """複数枚でもスライド全体を使って並べる"""
        md = f"## 全面\n<!-- layout: full-image -->\n\n![a]({png_file})\n\n![b]({png_file})\n"
        gen.generate(md, str(tmp_path / "out.pptx"))

        pictures = [
            s for s in gen.prs.slides[0].shapes
            if s.shape_type is not None and "PICTURE" in str(s.shape_type)
        ]
        assert len(pictures) == 2
        assert min(p.left for p in pictures) < gen.layout.content_left


class TestSlideNumbers:
    def test_numbers_are_added_except_first_slide(self, base_config, tmp_path):
        """先頭スライドを除いてページ番号が挿入される"""
        base_config["slides"]["show_slide_number"] = True
        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n\n## 2枚目\n\n## 3枚目\n", str(tmp_path / "out.pptx"))

        assert textboxes_of(gen.prs.slides[0]) == []
        assert textboxes_of(gen.prs.slides[1])[0].text_frame.text == "1"
        assert textboxes_of(gen.prs.slides[2])[0].text_frame.text == "2"

    def test_numbers_can_be_disabled(self, base_config, tmp_path):
        """show_slide_number: false で無効化できる"""
        base_config["slides"]["show_slide_number"] = False
        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n\n## 2枚目\n", str(tmp_path / "out.pptx"))
        assert textboxes_of(gen.prs.slides[1]) == []

    def test_enabled_by_default(self, tmp_path):
        """設定が無い場合はページ番号を表示する"""
        gen = PPTXGenerator({"slides": {"layout": "16:9"}})
        gen.generate("# 表紙\n\n## 2枚目\n", str(tmp_path / "out.pptx"))
        assert textboxes_of(gen.prs.slides[1])[0].text_frame.text == "1"


class TestGenerate:
    @patch("processors.insert_image_fit")
    @patch("requests.get")
    def test_markdown_integration(
        self, mock_get, mock_insert_image, base_config, mock_response, tmp_path
    ):
        """Markdownのパースからスライド生成までの一連の結合テスト"""
        mock_get.return_value = mock_response

        md_content = """
# タイトルスライド
ここはタイトルです。

## テキストと装飾のテスト
### 小見出し
* 箇条書き1
* **太字** と `インラインコード`

> これはスピーカーノートです。

## 画像のテスト
![テスト画像](http://example.com/test.png)

## 表のテスト
| 列A | 列B |
|---|---|
| 値1 | 値2 |

## Mermaidのテスト
```mermaid
graph TD; A-->B;
```

## コードのテスト
```python
print(1)
```
"""
        gen = PPTXGenerator(base_config)
        output_path = tmp_path / "test_output.pptx"
        gen.generate(md_content, str(output_path))

        # ファイルが生成され、h1/h2の数だけスライドが作られる
        assert os.path.exists(output_path)
        assert len(gen.prs.slides) == 6

        # ノートが正しいスライドに追加される
        notes = gen.prs.slides[1].notes_slide.notes_text_frame.text
        assert "これはスピーカーノートです。" in notes

        # 外部通信と画像挿入は 画像1回 + Mermaid1回
        assert mock_get.call_count == 2
        assert mock_insert_image.call_count == 2

    def test_content_before_first_heading_is_ignored(self, gen, tmp_path):
        """見出しより前の本文は配置先が無いため無視される"""
        gen.generate("本文だけ\n\n# タイトル\n中身", str(tmp_path / "out.pptx"))

        assert len(gen.prs.slides) == 1
        assert gen.current_body.text == "中身"

    def test_list_items_are_converted(self, gen, tmp_path):
        """箇条書きは1項目につき1段落になる"""
        gen.generate("# タイトル\n\n* 項目1\n* 項目2\n", str(tmp_path / "out.pptx"))
        assert [p.text for p in gen.current_body.paragraphs] == ["項目1", "項目2"]

    def test_list_item_paragraph_is_not_duplicated(self, gen, tmp_path):
        """段落を含むリスト（loose list）でもpタグが二重に出力されない"""
        gen.generate("# タイトル\n\n* 項目1\n\n* 項目2\n", str(tmp_path / "out.pptx"))
        assert [p.text.strip() for p in gen.current_body.paragraphs] == ["項目1", "項目2"]

    def test_blockquote_paragraph_is_not_duplicated(self, gen, tmp_path):
        """引用内のpタグは本文に出力されない"""
        gen.generate("# タイトル\n\n> ノート\n", str(tmp_path / "out.pptx"))

        assert gen.current_body.text == ""
        assert gen.current_slide.notes_slide.notes_text_frame.text == "ノート"

    def test_blockquote_list_is_not_duplicated(self, gen, tmp_path):
        """引用内の箇条書きが本文に漏れ出さない

        除外条件が p だけだった頃は、ノートに書いた箇条書きが
        スライド本文にもそのまま出力されていた。
        """
        md = "# タイトル\n\n> 発表のポイント\n>\n> * 最初に結論を述べる\n> * 質問は最後に受ける\n"
        gen.generate(md, str(tmp_path / "out.pptx"))

        assert gen.current_body.text == ""
        notes = gen.current_slide.notes_slide.notes_text_frame.text
        assert notes == "発表のポイント\n最初に結論を述べる\n質問は最後に受ける"

    @pytest.mark.parametrize(
        "content", ["| A | B |\n> |---|---|\n> | 1 | 2 |", "```\ncode\n```"]
    )
    def test_other_blocks_in_blockquote_stay_in_notes(self, gen, tmp_path, content):
        """引用内の表やコードブロックも本文には出力されない"""
        quoted = "\n".join(f"> {line}" for line in content.split("\n"))
        gen.generate(f"# タイトル\n\n> 前置き\n>\n{quoted}\n", str(tmp_path / "out.pptx"))

        assert gen.current_body.text == ""
        assert len(gen.current_slide.shapes) == 2  # タイトルと本文のみ

    def test_hr_creates_additional_slide(self, gen, tmp_path):
        """水平線でスライドが追加される"""
        gen.generate("## 1枚目\n\n本文\n\n---\n\n続き\n", str(tmp_path / "out.pptx"))
        assert len(gen.prs.slides) == 2


# =====================================================================
# md2pptx.py（CLI）
# =====================================================================


class TestConfigValidation:
    """config.yaml の検証"""

    def _errors(self, config):
        return validate_config(config).errors

    def _warnings(self, config):
        return validate_config(config).warnings

    def test_valid_config_passes(self):
        """同梱の config.yaml は警告もエラーも出ない"""
        config = yaml.safe_load(Path("config.yaml").read_text(encoding="utf-8"))
        result = validate_config(config)

        assert result.errors == []
        assert result.warnings == []

    @pytest.mark.parametrize("config", [None, {}])
    def test_empty_config_is_valid(self, config):
        """設定が空でも既定値で動くため、エラーにはしない"""
        assert validate_config(config).is_valid

    def test_non_mapping_is_rejected(self):
        """最上位がマッピングでない場合はエラーにする"""
        assert "最上位" in self._errors(["slides"])[0]

    # --- 型の誤り ---

    def test_size_must_be_a_number(self):
        """文字列のサイズ指定を分かりやすく指摘する

        従来は python-pptx まで届き
        「Exceeds the limit (4300 digits) for integer string conversion」
        という、設定ミスとは読み取れない例外になっていた。
        """
        (error,) = self._errors({"fonts": {"body": {"size_pt": "20"}}})
        assert "fonts.body.size_pt" in error
        assert "数値で指定" in error

    @pytest.mark.parametrize("size", [0, -5])
    def test_size_must_be_positive(self, size):
        (error,) = self._errors({"fonts": {"body": {"size_pt": size}}})
        assert "0より大きい" in error

    def test_bold_must_be_boolean(self):
        (error,) = self._errors({"fonts": {"body": {"bold": "yes"}}})
        assert "true か false" in error

    @pytest.mark.parametrize(
        "color, expected",
        [
            ([10, 20], "3つの値が必要"),
            ("#0070C0", "[赤, 緑, 青]"),
            ([300, 0, 0], "0〜255"),
            ([0, "x", 0], "整数で指定"),
        ],
    )
    def test_color_validation(self, color, expected):
        """色指定の誤りを、どの成分が悪いかまで含めて指摘する"""
        (error,) = self._errors({"fonts": {"body": {"color_rgb": color}}})
        assert expected in error

    def test_theme_color_validation(self):
        """テーマ色も同じ基準で検証する"""
        (error,) = self._errors({"theme": {"accent_color": "#0070C0"}})
        assert "theme.accent_color" in error

    def test_image_settings(self):
        errors = self._errors({"images": {"dpi": "high", "downscale": "yes"}})
        assert len(errors) == 2

    @pytest.mark.parametrize("position", [[1.0], "5.2, 1.8", [1.0, "x"]])
    def test_position_inches_validation(self, position):
        (error,) = self._errors({"images": {"position_inches": position}})
        assert "images.position_inches" in error

    def test_mermaid_renderer_choices(self):
        """選択肢の誤りは、正しい候補を添えて指摘する"""
        (error,) = self._errors({"mermaid": {"renderer": "kroky"}})
        assert "kroki / local / off" in error
        assert "'kroki' の誤りではありませんか？" in error

    def test_nested_value_must_be_a_mapping(self):
        (error,) = self._errors({"slides": "16:9"})
        assert "入れ子" in error

    # --- 綴り違いの警告（従来は黙って無視されていた） ---

    def test_unknown_top_level_key_suggests_correction(self):
        (warning,) = self._warnings({"font": {}})
        assert "'fonts' の誤りではありませんか？" in warning

    def test_unknown_font_key_suggests_correction(self):
        (warning,) = self._warnings({"fonts": {"bodyy": {"size_pt": 18}}})
        assert "'body' の誤りではありませんか？" in warning

    def test_unknown_font_field_suggests_correction(self):
        (warning,) = self._warnings({"fonts": {"body": {"sise_pt": 18}}})
        assert "'size_pt' の誤りではありませんか？" in warning

    def test_bullet_levels_are_accepted(self):
        """bullet_level_N は任意の階層を指定できる"""
        config = {"fonts": {f"bullet_level_{n}": {"size_pt": 18} for n in range(1, 6)}}
        assert validate_config(config).warnings == []

    @pytest.mark.parametrize(
        "footer, expected",
        [
            ({"text": 123}, "文字列で指定"),
            ({"show_on_title": "yes"}, "true か false"),
            ({"date": 123}, "true（変換日を表示）"),
        ],
    )
    def test_footer_validation(self, footer, expected):
        """フッター設定の型の誤りを指摘する"""
        (error,) = self._errors({"slides": {"footer": footer}})
        assert expected in error

    def test_footer_must_be_a_mapping(self):
        (error,) = self._errors({"slides": {"footer": "社外秘"}})
        assert "入れ子" in error

    def test_footer_unknown_key(self):
        (warning,) = self._warnings({"slides": {"footer": {"txt": "a"}}})
        assert "'text' の誤りではありませんか？" in warning

    def test_dark_theme_colors_are_validated(self):
        """ダークテーマの配色も色として検証される"""
        (error,) = self._errors({"theme": {"dark_background_color": "#000000"}})
        assert "theme.dark_background_color" in error

    def test_h3_as_choices(self):
        """h3_as の選択肢の誤りを候補付きで指摘する"""
        (error,) = self._errors({"slides": {"h3_as": "slides"}})
        assert "subheading / slide" in error
        assert "'slide' の誤りではありませんか？" in error

    def test_unknown_layout_warns_about_fallback(self):
        """未対応の画角は、代わりに使われる値まで示す"""
        (warning,) = self._warnings({"slides": {"layout": "16:9ワイド"}})
        assert "16:9 が使われます" in warning

    def test_unrelated_key_warns_without_suggestion(self):
        """似た候補が無い場合は候補を出さない"""
        (warning,) = self._warnings({"まったく無関係": 1})
        assert "誤りではありませんか" not in warning

    @pytest.mark.parametrize(
        "config, path",
        [
            ({"slides": {"template_path": 123}}, "slides.template_path"),
            ({"mermaid": {"endpoint": 123}}, "mermaid.endpoint"),
            ({"mermaid": {"cli_path": []}}, "mermaid.cli_path"),
            ({"fonts": {"body": {"name": 123}}}, "fonts.body.name"),
        ],
    )
    def test_string_fields(self, config, path):
        """文字列を期待する項目に他の型を書いた場合を指摘する"""
        (error,) = self._errors(config)
        assert path in error
        assert "文字列で指定" in error

    @pytest.mark.parametrize(
        "section", ["slides", "fonts", "images", "theme", "mermaid"]
    )
    def test_empty_section_is_accepted(self, section):
        """セクションを書いて中身が空（None）でも既定値で動くため許容する"""
        assert validate_config({section: None}).is_valid

    @pytest.mark.parametrize(
        "config", [
            {"fonts": "Meiryo"},
            {"fonts": {"body": "Meiryo"}},
            {"images": "3.5"},
            {"theme": "dark"},
            {"mermaid": "kroki"},
        ],
    )
    def test_sections_must_be_mappings(self, config):
        """セクションの値がマッピングでない場合を指摘する"""
        (error,) = self._errors(config)
        assert "入れ子" in error

    def test_boolean_mermaid_flags(self):
        errors = self._errors(
            {"mermaid": {"warn_on_external": "yes", "fallback_to_public": 1}}
        )
        assert len(errors) == 2

    def test_all_problems_are_reported_at_once(self):
        """1件目で止めず、まとめて指摘する"""
        config = {
            "fonts": {"body": {"size_pt": "20"}, "title_h1": {"color_rgb": [300, 0, 0]}},
            "images": {"dpi": "high"},
        }
        assert len(self._errors(config)) == 3


class TestApplyTheme:
    def test_accent_and_text_colors(self):
        """テーマ色が見出し系・本文系のフォント設定に展開される"""
        config = apply_theme(
            {"theme": {"accent_color": [1, 2, 3], "text_color": [4, 5, 6]}, "fonts": {}}
        )

        assert config["fonts"]["title_h1"]["color_rgb"] == [1, 2, 3]
        assert config["fonts"]["table_header"]["color_rgb"] == [1, 2, 3]
        assert config["fonts"]["body"]["color_rgb"] == [4, 5, 6]
        assert config["fonts"]["table_body"]["color_rgb"] == [4, 5, 6]

    def test_existing_font_settings_are_kept(self):
        """既存のフォント設定（サイズ等）は保持したまま色だけ上書きする"""
        config = apply_theme(
            {
                "theme": {"accent_color": [1, 2, 3]},
                "fonts": {"title_h1": {"name": "Meiryo", "size_pt": 44}},
            }
        )

        assert config["fonts"]["title_h1"]["name"] == "Meiryo"
        assert config["fonts"]["title_h1"]["color_rgb"] == [1, 2, 3]

    def test_without_fonts_section(self):
        """fontsセクションが無くても展開できる"""
        config = apply_theme({"theme": {"accent_color": [1, 2, 3]}})
        assert config["fonts"]["title_h1"]["color_rgb"] == [1, 2, 3]

    @pytest.mark.parametrize("config", [{}, {"theme": None}, {"theme": {}}])
    def test_without_theme_is_noop(self, config):
        """テーマが無い場合は設定を変更しない"""
        assert apply_theme(config) == config


class TestCli:
    def test_parse_args_defaults(self):
        """出力先と設定ファイルの既定値"""
        args = parse_args(["input.md"])
        assert (args.input, args.output, args.config) == (
            "input.md",
            "output.pptx",
            "config.yaml",
        )

    def test_parse_args_options(self):
        """オプション指定が反映される"""
        args = parse_args(["in.md", "-o", "out.pptx", "-c", "my.yaml"])
        assert (args.output, args.config) == ("out.pptx", "my.yaml")

    def test_load_config(self, tmp_path):
        """YAML設定ファイルを辞書として読み込む"""
        path = tmp_path / "c.yaml"
        path.write_text("slides:\n  layout: '4:3'\n", encoding="utf-8")
        assert load_config(str(path)) == {"slides": {"layout": "4:3"}}

    def test_load_config_empty_file(self, tmp_path):
        """空のYAMLは空の辞書として扱う（Noneを返さない）"""
        path = tmp_path / "empty.yaml"
        path.write_text("", encoding="utf-8")
        assert load_config(str(path)) == {}

    def test_read_text_file(self, tmp_path):
        """UTF-8のテキストを読み込む"""
        path = tmp_path / "a.md"
        path.write_text("# 日本語", encoding="utf-8")
        assert read_text_file(str(path)) == "# 日本語"

    def _write_project(self, tmp_path, config_text="slides:\n  layout: '16:9'\n"):
        md = tmp_path / "in.md"
        md.write_text("# タイトル\n本文\n", encoding="utf-8")
        conf = tmp_path / "c.yaml"
        conf.write_text(config_text, encoding="utf-8")
        return str(md), str(conf), str(tmp_path / "out.pptx")

    def test_main_success(self, tmp_path, capsys):
        """正常系では0を返しファイルを生成する"""
        md, conf, out = self._write_project(tmp_path)

        assert main([md, "-o", out, "-c", conf]) == 0
        assert os.path.exists(out)
        assert "Success" in capsys.readouterr().out

    def test_main_with_empty_config(self, tmp_path):
        """空の設定ファイルでも変換できる"""
        md, conf, out = self._write_project(tmp_path, config_text="")
        assert main([md, "-o", out, "-c", conf]) == 0

    def test_main_applies_theme(self, tmp_path, mocker):
        """テーマ設定がジェネレーターに渡る設定へ反映される"""
        md, conf, out = self._write_project(
            tmp_path, config_text="theme:\n  accent_color: [1, 2, 3]\n"
        )
        spy = mocker.spy(md2pptx, "PPTXGenerator")

        assert main([md, "-o", out, "-c", conf]) == 0
        assert spy.call_args[0][0]["fonts"]["title_h1"]["color_rgb"] == [1, 2, 3]

    def test_main_rejects_invalid_config(self, tmp_path, capsys):
        """設定に誤りがあれば、変換を始める前に中止する"""
        md, _, out = self._write_project(tmp_path)
        conf = tmp_path / "bad.yaml"
        conf.write_text('fonts:\n  body:\n    size_pt: "20"\n', encoding="utf-8")

        assert main([md, "-o", out, "-c", str(conf)]) == 1

        captured = capsys.readouterr().out
        assert "fonts.body.size_pt" in captured
        assert "数値で指定" in captured
        assert not os.path.exists(out)  # 途中まで書き出さない

    def test_main_warns_but_continues(self, tmp_path, capsys):
        """警告のみの場合は指摘したうえで変換を続ける"""
        md, _, out = self._write_project(tmp_path)
        conf = tmp_path / "typo.yaml"
        conf.write_text("font:\n  body:\n    size_pt: 20\n", encoding="utf-8")

        assert main([md, "-o", out, "-c", str(conf)]) == 0

        assert "'fonts' の誤りではありませんか？" in capsys.readouterr().out
        assert os.path.exists(out)

    def test_main_uses_sys_argv_by_default(self, tmp_path, mocker):
        """引数を渡さない場合はsys.argvを解析する"""
        md, conf, out = self._write_project(tmp_path)
        mocker.patch.object(md2pptx.sys, "argv", ["md2pptx.py", md, "-o", out, "-c", conf])
        assert main() == 0

    def test_main_missing_input(self, tmp_path, capsys):
        """入力ファイルが無い場合は1を返す"""
        _, conf, out = self._write_project(tmp_path)

        assert main(["no_such.md", "-o", out, "-c", conf]) == 1
        assert "入力ファイル" in capsys.readouterr().out

    def test_main_missing_config(self, tmp_path, capsys):
        """設定ファイルが無い場合は1を返す"""
        md, _, out = self._write_project(tmp_path)

        assert main([md, "-o", out, "-c", "no_such.yaml"]) == 1
        assert "設定ファイル" in capsys.readouterr().out

    def test_main_permission_error(self, tmp_path, mocker, capsys):
        """出力先に書き込めない場合は専用のメッセージを表示する"""
        md, conf, out = self._write_project(tmp_path)
        mocker.patch.object(PPTXGenerator, "generate", side_effect=PermissionError)

        assert main([md, "-o", out, "-c", conf]) == 1
        assert "書き込めません" in capsys.readouterr().out

    def test_main_unexpected_error(self, tmp_path, mocker, capsys):
        """想定外の例外はメッセージとトレースバックを表示して1を返す"""
        md, conf, out = self._write_project(tmp_path)
        mocker.patch.object(PPTXGenerator, "generate", side_effect=ValueError("boom"))

        assert main([md, "-o", out, "-c", conf]) == 1
        captured = capsys.readouterr()
        assert "予期せぬエラー" in captured.out
        assert "ValueError" in captured.err


# =====================================================================
# extractor.py / pptx2md.py（PPTX → Markdown の逆変換）
# =====================================================================


@pytest.fixture
def converted(base_config, tmp_path, png_file):
    """Markdownから変換したPPTXを、逆変換して読み戻す"""
    def convert(md, **config_overrides):
        base_config["slides"].update(config_overrides)
        gen = PPTXGenerator(base_config)
        path = tmp_path / "src.pptx"
        gen.generate(md.replace("{png}", png_file), str(path))
        return extract(Presentation(str(path)))
    return convert


class TestExtractText:
    """見出し・本文・リストの復元"""

    def test_headings(self, converted):
        result = converted("# タイトル\n\n## 中身\n本文\n")

        assert "# タイトル" in result.markdown
        assert "## 中身" in result.markdown

    def test_subheading(self, converted):
        """h3（行頭記号を消した段落）は ### に戻る"""
        result = converted("## 見出し\n\n### 小見出し\n本文\n")
        assert "### 小見出し" in result.markdown

    def test_nested_list_indentation(self, converted):
        """入れ子の階層はインデントで表現される"""
        result = converted("## 見出し\n\n* 親\n    * 子\n        * 孫\n")

        assert "* 親" in result.markdown
        assert "    * 子" in result.markdown
        assert "        * 孫" in result.markdown

    def test_ordered_list(self, converted):
        """自動採番の段落は番号付きリストに戻る"""
        result = converted("## 手順\n\n1. 最初\n2. 次\n")
        assert result.markdown.count("1. ") == 2  # Markdownでは連番でなくてよい

    def test_plain_paragraph_stays_plain(self, converted):
        """平文の段落は箇条書きにならない

        以前は本文枠のすべての段落が箇条書きとして書き出されていた。
        """
        result = converted("## 見出し\n\n平文の段落です。\n\n* 箇条書き\n")

        assert "\n平文の段落です。\n" in result.markdown
        assert "* 平文の段落です。" not in result.markdown

    def test_bullet_stays_a_bullet(self, converted):
        """同じ本文の中で、箇条書きは箇条書きのまま戻る"""
        result = converted("## 見出し\n\n平文の段落です。\n\n* 箇条書き\n")

        assert "* 箇条書き" in result.markdown

    def test_inline_decorations(self, converted):
        """太字・インラインコードが復元される"""
        result = converted("## 見出し\n\n**太字** と `コード` です\n")

        assert "**太字**" in result.markdown
        assert "`コード`" in result.markdown

    def test_slide_without_title_becomes_a_separator(self, converted):
        """タイトルの無いスライド（--- 由来）は水平線に戻る"""
        result = converted("## 見出し\n\n本文\n\n---\n\n続き\n")
        assert "\n---\n" in result.markdown


class TestExtractBlocks:
    """表・コード・画像・ノートの復元"""

    def test_table_with_alignment(self, converted):
        """表は列揃えの指定ごと復元される"""
        result = converted("## 表\n\n| 左 | 中 | 右 |\n| :--- | :---: | ---: |\n| a | b | c |\n")

        assert "| 左 | 中 | 右 |" in result.markdown
        assert "| :--- | :---: | ---: |" in result.markdown
        assert "| a | b | c |" in result.markdown

    def test_table_header_is_not_double_emphasized(self, converted):
        """見出し行はMarkdown側で強調されるため ** を付けない"""
        result = converted("## 表\n\n| 見出し |\n|---|\n| 値 |\n")
        assert "| **見出し** |" not in result.markdown

    def test_code_block(self, converted):
        """コード枠はフェンス付きコードブロックに戻る"""
        result = converted("## コード\n\n```python\nprint(1)\n```\n")

        assert "```" in result.markdown
        assert "print(1)" in result.markdown

    def test_image_and_alt_text(self, converted):
        """画像は書き出され、代替テキストが記法に戻る"""
        result = converted("## 図\n\n![システム構成]({png})\n")

        assert "![システム構成](images/image1.png)" in result.markdown
        assert len(result.images) == 1
        assert result.images[0].filename == "image1.png"

    def test_same_image_shares_one_file(self, converted):
        """同じ画像を複数回使っている場合は1ファイルにまとめ、同じ名前を参照する

        連番を振ってしまうと、まとめた結果として存在しないファイルを
        参照することになる。
        """
        result = converted("## 図\n\n![a]({png})\n\n## 図2\n\n![b]({png})\n")

        assert len(result.images) == 1
        assert result.markdown.count("images/image1.png") == 2

    def test_speaker_notes(self, converted):
        """スピーカーノートは引用記法に戻る"""
        result = converted("## 見出し\n\n> メモです\n")
        assert "> メモです" in result.markdown

    def test_footer_is_not_extracted(self, converted):
        """フッターやページ番号は内容ではないので取り出さない"""
        result = converted(
            "## 見出し\n\n本文\n", footer={"text": "社外秘", "date": "2026-05-10"}
        )

        assert "社外秘" not in result.markdown
        assert "2026-05-10" not in result.markdown


class TestRoundTrip:
    """Markdown → PPTX → Markdown の往復"""

    MD = (
        "# タイトル\n\n## 中身\n\n* 項目1\n    * 入れ子\n\n1. 手順\n\n"
        "### 小見出し\n\n**太字** と `コード`\n\n"
        "| A | B |\n| :--- | ---: |\n| 1 | 2 |\n\n"
        "```\ncode\n```\n\n![図]({png})\n\n> ノート\n"
    )

    def test_round_trip_is_stable(self, base_config, tmp_path, png_file):
        """2回往復させても内容が変わらない

        1回目の復元結果を再び変換して復元し、同じになることを確かめる。
        """
        md = self.MD.replace("{png}", png_file)

        first = self._cycle(base_config, tmp_path, md, "1")
        # 1回目の復元結果は images/ を参照するため、画像を書き出してから再変換する
        write_images(first.images, str(tmp_path / "images"))
        second = self._cycle(base_config, tmp_path, first.markdown, "2", cwd=tmp_path)

        assert first.markdown == second.markdown

    def _cycle(self, base_config, tmp_path, md, suffix, cwd=None):
        import os
        original = os.getcwd()
        if cwd:
            os.chdir(cwd)
        try:
            gen = PPTXGenerator(base_config)
            path = tmp_path / f"cycle{suffix}.pptx"
            gen.generate(md, str(path))
            return extract(Presentation(str(path)))
        finally:
            os.chdir(original)

    def test_content_survives_the_trip(self, base_config, tmp_path, png_file):
        """主要な要素が往復後も残っている"""
        gen = PPTXGenerator(base_config)
        path = tmp_path / "src.pptx"
        gen.generate(self.MD.replace("{png}", png_file), str(path))

        markdown = extract(Presentation(str(path))).markdown
        for expected in ["# タイトル", "## 中身", "### 小見出し", "**太字**",
                         "`コード`", "| A | B |", "```", "> ノート"]:
            assert expected in markdown


class TestPptx2mdCli:
    """逆変換のCLI"""

    def _pptx(self, base_config, tmp_path, md="## 見出し\n\n本文\n"):
        gen = PPTXGenerator(base_config)
        path = tmp_path / "src.pptx"
        gen.generate(md, str(path))
        return str(path)

    def test_parse_args_defaults(self):
        args = pptx2md.parse_args(["deck.pptx"])
        assert (args.input, args.output, args.image_dir) == ("deck.pptx", "output.md", None)

    def test_main_success(self, base_config, tmp_path, capsys):
        source = self._pptx(base_config, tmp_path)
        output = tmp_path / "out.md"

        assert pptx2md.main([source, "-o", str(output)]) == 0
        assert "## 見出し" in output.read_text(encoding="utf-8")
        assert "Success" in capsys.readouterr().out

    def test_images_are_written(self, base_config, tmp_path, png_file, capsys):
        source = self._pptx(base_config, tmp_path, f"## 図\n\n![a]({png_file})\n")
        output = tmp_path / "out.md"

        pptx2md.main([source, "-o", str(output)])

        assert (tmp_path / "images" / "image1.png").exists()
        assert "画像 1枚" in capsys.readouterr().out

    def test_custom_image_dir(self, base_config, tmp_path, png_file):
        source = self._pptx(base_config, tmp_path, f"## 図\n\n![a]({png_file})\n")
        image_dir = tmp_path / "assets"

        pptx2md.main([source, "-o", str(tmp_path / "o.md"), "--image-dir", str(image_dir)])

        assert (image_dir / "image1.png").exists()

    def test_missing_input(self, tmp_path, capsys):
        assert pptx2md.main(["no_such.pptx", "-o", str(tmp_path / "o.md")]) == 1
        assert "入力ファイル" in capsys.readouterr().out

    def test_broken_input(self, tmp_path, capsys):
        broken = tmp_path / "broken.pptx"
        broken.write_text("これはPPTXではありません", encoding="utf-8")

        assert pptx2md.main([str(broken), "-o", str(tmp_path / "o.md")]) == 1
        assert "予期せぬエラー" in capsys.readouterr().out

    def test_permission_error(self, base_config, tmp_path, mocker, capsys):
        source = self._pptx(base_config, tmp_path)
        # 読み込みには影響しないよう、書き出し側の open だけを差し替える
        mocker.patch("pptx2md.open", side_effect=PermissionError, create=True)

        assert pptx2md.main([source, "-o", str(tmp_path / "o.md")]) == 1
        assert "書き込めません" in capsys.readouterr().out


class TestExtractEdgeCases:
    """逆変換の細かい挙動"""

    def test_italic(self, converted):
        """斜体も復元される"""
        result = converted("## 見出し\n\n*斜体* です\n")
        assert "*斜体*" in result.markdown

    def test_title_slide_subtitle_is_not_a_list(self, converted):
        """タイトルスライドのサブタイトルは箇条書きにしない"""
        result = converted("# タイトル\nサブタイトル\n")

        assert "サブタイトル" in result.markdown
        assert "* サブタイトル" not in result.markdown

    def test_image_without_alt_text(self, converted):
        """代替テキストが無い画像は空の alt になる"""
        result = converted("## 図\n\n![]({png})\n")
        assert "![](images/image1.png)" in result.markdown

    def test_empty_table_is_skipped(self):
        """行の無い表は何も出力しない"""
        table = MagicMock()
        table.rows = []
        assert extractor.table_to_markdown(table) == []

    def test_slide_without_notes(self):
        """ノートが無いスライドは引用を出さない"""
        slide = MagicMock()
        slide.has_notes_slide = False
        assert extractor.notes_to_markdown(slide) == []

    def test_empty_notes_are_skipped(self):
        """ノート枠はあるが空の場合も引用を出さない"""
        slide = MagicMock()
        slide.has_notes_slide = True
        slide.notes_slide.notes_text_frame.text = "   "

        assert extractor.notes_to_markdown(slide) == []

    def test_consecutive_blank_lines_are_collapsed(self):
        """空行が3つ以上続く箇所は2つにまとめる"""
        prs = MagicMock()
        prs.slide_layouts = []
        prs.slides = []
        # 行の生成を差し替えて、空行が多い状態を作る
        import extractor as ex
        original = ex.slide_to_markdown
        try:
            prs.slides = [MagicMock()]
            ex.slide_to_markdown = lambda *a, **k: ["A", "", "", "", "B"]
            assert ex.extract(prs).markdown == "A\n\nB\n"
        finally:
            ex.slide_to_markdown = original

    def test_shape_without_fill_is_not_a_code_box(self):
        """塗りつぶしを調べられない図形はコード枠とみなさない"""
        shape = MagicMock()
        shape.is_placeholder = False
        shape.has_text_frame = True
        type(shape).fill = property(lambda self: (_ for _ in ()).throw(ValueError()))

        assert extractor.is_code_box(shape) is False

    def test_blank_lines_are_collapsed(self, converted):
        """空行が続く箇所はまとめられる"""
        result = converted("## 見出し\n\n本文1\n\n本文2\n")
        assert "\n\n\n" not in result.markdown


# =====================================================================
# テンプレート（テーマ）の利用
# =====================================================================


@pytest.fixture
def template_factory(tmp_path):
    """検証用のテンプレートPPTXを作る"""
    def make(name="template.pptx", reorder=None, width_in=13.333, height_in=7.5,
             body_idx=None):
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(width_in), Inches(height_in)
        if body_idx is not None:
            # 本文枠の idx を振り直す（Googleスライドから書き出したテンプレートを想定）
            for layout in prs.slide_layouts:
                for ph in layout.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph._element.nvSpPr.nvPr.ph.set('idx', str(body_idx))
        if reorder:
            # レイアウトの並び順を入れ替える（社内テンプレートを想定）
            id_list = prs.slide_master.slide_layouts._sldLayoutIdLst
            ids = list(id_list)
            id_list.remove(ids[reorder[0]])
            id_list.insert(reorder[0], ids[reorder[1]])
        path = tmp_path / name
        prs.save(str(path))
        return str(path)
    return make


class TestTemplate:
    """テンプレートから引き継がれるもの"""

    def _generate(self, base_config, tmp_path, template, **slides):
        base_config["slides"].update({"template_path": template, **slides})
        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n\n## 中身\n本文\n", str(tmp_path / "out.pptx"))
        return gen

    def test_slide_size_comes_from_the_template(self, base_config, tmp_path, template_factory):
        """スライドサイズはテンプレートに従う（slides.layout より優先）"""
        base_config["slides"]["layout"] = "16:9"
        gen = self._generate(base_config, tmp_path, template_factory())

        assert gen.prs.slide_width == Inches(13.333)

    def test_content_area_follows_the_template(self, base_config, tmp_path, template_factory):
        """画像・表の配置がテンプレートの本文枠に追従する"""
        gen = self._generate(base_config, tmp_path, template_factory())
        body = gen.prs.slide_layouts[1].placeholders[1]

        assert gen.layout.content_left == body.left

    def test_layouts_can_be_selected_by_name(self, base_config, tmp_path, template_factory):
        """レイアウトを名前で指定できる

        レイアウトの並び順はテンプレートによって異なるため、位置で決め打つと
        意図しないレイアウトが使われる。
        """
        gen = self._generate(
            base_config, tmp_path, template_factory(),
            layouts={"title": "Title Slide", "content": "Two Content"},
        )

        assert {s.slide_layout.name for s in gen.prs.slides} == {"Title Slide", "Two Content"}

    def test_unknown_layout_name_is_reported(self, base_config, tmp_path, template_factory):
        """存在しないレイアウト名は、使用できる一覧を添えて知らせる"""
        base_config["slides"].update({
            "template_path": template_factory(), "layouts": {"content": "無いレイアウト"},
        })

        with pytest.raises(TemplateError, match="無いレイアウト"):
            PPTXGenerator(base_config)

    def test_missing_layout_index_is_reported(self, base_config, tmp_path, mocker):
        """レイアウトが足りないテンプレートも分かるように知らせる"""
        gen_cls = PPTXGenerator
        mocker.patch.object(
            gen_cls, "_resolve_layouts",
            side_effect=lambda self=None: None,
        )
        # 直接確認する（レイアウトが1つも無いテンプレートは作れないため）
        gen = PPTXGenerator(base_config)
        mocker.stopall()
        gen.prs.slide_layouts._sldLayoutIdLst.clear()

        with pytest.raises(TemplateError, match="レイアウトが足りません"):
            gen._layout_by_index(1)


class TestFindBodyPlaceholder:
    """本文・副題を書き込む枠の選び方"""

    def test_idx1_is_preferred(self):
        layout = Presentation().slide_layouts[0]      # idx=0/1/10/11/12

        assert find_body_placeholder(layout).placeholder_format.idx == 1

    def test_falls_back_to_another_text_placeholder(self):
        """Googleスライド由来のテンプレートのように idx=1 が無い場合"""
        layout = Presentation().slide_layouts[3]      # idx=0/1/2/10/11/12
        body = layout.placeholders[1]
        body._element.getparent().remove(body._element)

        assert find_body_placeholder(layout).placeholder_format.idx == 2

    def test_title_and_footers_are_not_used(self):
        layout = Presentation().slide_layouts[5]      # タイトルのみ（idx=0/10/11/12）

        assert find_body_placeholder(layout) is None


class TestLayoutsAcrossMasters:
    """スライドマスターが複数あるテンプレート

    python-pptx の prs.slide_layouts は1つ目のマスターしか返さないため、
    本文用のレイアウトが2つ目以降にあると見つけられなかった。
    """

    def _layout(self, name, idxs=()):
        layout = MagicMock()
        layout.name = name
        layout.placeholders = [
            MagicMock(**{"placeholder_format.idx": idx}) for idx in idxs
        ]
        return layout

    def _generator(self, base_config, masters):
        gen = PPTXGenerator(base_config)
        gen.prs = MagicMock()
        gen.prs.slide_masters = [MagicMock(slide_layouts=layouts) for layouts in masters]
        gen.slide_layouts = {}
        return gen

    def test_layout_name_is_found_in_the_second_master(self, base_config):
        empty = self._layout("DEFAULT")
        body = self._layout("TITLE_AND_BODY", idxs=(0, 1))
        gen = self._generator(base_config, [[empty], [body]])
        gen.slides_conf = {"layouts": {"title": "DEFAULT", "content": "TITLE_AND_BODY"}}

        gen._resolve_layouts()

        assert gen.slide_layouts["content"].name == "TITLE_AND_BODY"

    def test_index_selection_continues_into_the_next_master(self, base_config):
        first = self._layout("DEFAULT", idxs=(0,))
        second = self._layout("TITLE_AND_BODY", idxs=(0, 1))
        gen = self._generator(base_config, [[first], [second]])
        gen.slides_conf = {}

        gen._resolve_layouts()

        assert gen.slide_layouts["title"].name == "DEFAULT"
        assert gen.slide_layouts["content"].name == "TITLE_AND_BODY"

    def test_body_fallback_looks_at_every_master(self, base_config, capsys):
        """本文枠が2つ目のマスターにしか無い場合でも見つける"""
        empty = self._layout("DEFAULT")
        another = self._layout("1_DEFAULT")
        body = self._layout("TITLE_AND_BODY", idxs=(0, 1))
        gen = self._generator(base_config, [[empty, another], [body]])
        gen.slides_conf = {}

        gen._resolve_layouts()

        assert gen.slide_layouts["content"].name == "TITLE_AND_BODY"
        assert "本文枠が無いため" in capsys.readouterr().out

    def test_error_lists_layouts_from_every_master(self, base_config):
        gen = self._generator(base_config, [[self._layout("DEFAULT")], [self._layout("SECTION")]])
        gen.slides_conf = {"layouts": {"content": "無い名前"}}

        with pytest.raises(TemplateError, match="DEFAULT / SECTION"):
            gen._resolve_layouts()


class TestTitleSlidePlaceholders:
    """表紙の副題（subtitle / author / date）の書き込み先"""

    FRONT_MATTER = (
        "---\n"
        'title: "表紙"\n'
        'subtitle: "副題"\n'
        "---\n\n"
        "## 中身\n本文\n"
    )

    def _title_slide(self, base_config, tmp_path, template, **slides):
        base_config["slides"].update({"template_path": template, **slides})
        gen = PPTXGenerator(base_config)
        gen.generate(self.FRONT_MATTER, str(tmp_path / "out.pptx"))
        return gen.prs.slides[0]

    def test_subtitle_uses_another_placeholder_when_idx1_is_missing(
        self, base_config, tmp_path
    ):
        """副題がidx=1でないテンプレート（Googleスライド由来など）でも書き込む"""
        prs = Presentation()
        layout = prs.slide_layouts[3]           # idx=0/1/2 を持つレイアウト
        body = layout.placeholders[1]
        body._element.getparent().remove(body._element)   # idx=1 だけ取り除く
        path = tmp_path / "no_idx1.pptx"
        prs.save(str(path))

        slide = self._title_slide(
            base_config, tmp_path, str(path),
            layouts={"title": layout.name, "content": "Title and Content"},
        )

        subtitle = slide.placeholders[2]
        assert subtitle.text_frame.text == "副題"

    def test_missing_subtitle_placeholder_warns_instead_of_failing(
        self, base_config, tmp_path, template_factory, capsys
    ):
        """副題を書ける枠が無くても変換は止まらない"""
        slide = self._title_slide(
            base_config, tmp_path, template_factory(),
            layouts={"title": "Title Only", "content": "Title and Content"},
        )

        assert slide.shapes.title.text_frame.text == "表紙"
        assert "副題" in capsys.readouterr().out


class TestBodyWithoutPlaceholder:
    """本文枠が無いレイアウトが使われた場合でも本文を捨てない"""

    def test_h1_slide_gets_a_textbox(self, base_config, tmp_path, template_factory):
        base_config["slides"].update({
            "template_path": template_factory(),
            "layouts": {"title": "Title Only", "content": "Title and Content"},
        })

        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n本文のテキスト\n", str(tmp_path / "out.pptx"))

        slide = gen.prs.slides[0]
        assert slide.shapes.title.text_frame.text == "表紙"
        textboxes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        assert any("本文のテキスト" in box.text_frame.text for box in textboxes)


class TestTemplateWithSlides:
    """中身が入ったファイルをテンプレートに指定した場合"""

    def test_existing_slides_are_reported(self, base_config, tmp_path, capsys):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / "deck.pptx"
        prs.save(str(path))
        base_config["slides"]["template_path"] = str(path)

        gen = PPTXGenerator(base_config)
        gen.generate("## 見出し\n本文\n", str(tmp_path / "out.pptx"))

        captured = capsys.readouterr().out
        assert "既存のスライドが1枚あります" in captured
        # 既存のスライドは消さず、生成分をその後ろに足す
        assert len(gen.prs.slides) == 2


class TestBodyPlaceholderFallback:
    """本文枠を持たないレイアウトが選ばれた場合"""

    def test_falls_back_with_a_warning(self, base_config, tmp_path, template_factory, capsys):
        """本文枠が無ければ、持っているレイアウトに切り替えて知らせる

        以前は本文を書き込む段階で
        「no placeholder on this slide with idx == 1」という
        原因の分かりにくいエラーになっていた。
        """
        # レイアウト[1]を「Title Only」（本文枠なし）に差し替える
        template = template_factory(reorder=(1, 5))
        base_config["slides"]["template_path"] = template

        gen = PPTXGenerator(base_config)
        gen.generate("## 中身\n本文\n", str(tmp_path / "out.pptx"))

        captured = capsys.readouterr().out
        assert "本文枠が無いため" in captured
        assert gen.slide_layouts["content"].name != "Title Only"

    def test_error_when_no_layout_has_a_body(self, base_config, mocker):
        """本文枠を持つレイアウトが1つも無ければエラーにする"""
        mocker.patch("generator._has_body_placeholder", return_value=False)

        with pytest.raises(TemplateError, match="本文を書き込めるレイアウトがありません"):
            PPTXGenerator(base_config)


class TestNonStandardBodyIndex:
    """本文枠の idx が 1 でないテンプレート（Googleスライドからの書き出しなど）

    以前は本文枠の探し方が箇所ごとに違い、idx=1 を決め打ちしている場所が
    残っていた。そのため PPTXGenerator の生成時に弾かれたり、水平線の
    スライドで KeyError になったりしていた。
    """

    def _config(self, base_config, template_factory):
        base_config["slides"]["template_path"] = template_factory(body_idx=12)
        return base_config

    def test_template_is_accepted(self, base_config, template_factory):
        """idx=1 が無くてもテンプレートとして受け付ける"""
        gen = PPTXGenerator(self._config(base_config, template_factory))

        assert find_body_placeholder(gen.slide_layouts["content"]) is not None

    def test_body_text_is_written(self, base_config, tmp_path, template_factory):
        """本文が idx=1 以外のプレースホルダーに書き込まれる"""
        gen = PPTXGenerator(self._config(base_config, template_factory))
        gen.generate("## 中身\n本文\n", str(tmp_path / "out.pptx"))

        body = find_body_placeholder(gen.prs.slides[0])
        assert body.placeholder_format.idx == 12
        assert "本文" in body.text_frame.text

    def test_horizontal_rule_slide_is_created(self, base_config, tmp_path, template_factory):
        """水平線のスライドでも落ちない（以前は KeyError になっていた）"""
        gen = PPTXGenerator(self._config(base_config, template_factory))
        gen.generate("## 中身\n本文\n\n---\n\n続き\n", str(tmp_path / "out.pptx"))

        bodies = [find_body_placeholder(s).text_frame.text for s in gen.prs.slides]
        assert "続き" in bodies[1]

    def test_content_area_follows_the_template(self, base_config, template_factory):
        """画像・表の配置も idx=1 以外の本文枠に追従する"""
        gen = PPTXGenerator(self._config(base_config, template_factory))
        body = find_body_placeholder(gen.slide_layouts["content"])

        assert gen.layout.content_left == body.left


class TestBodyPlaceholderMissingAtRuntime:
    """本文枠が見つからないままスライドを組み立てる場合

    本文はテキストボックスに書き出されるため、プレースホルダーを
    前提にした寸法補正は何もせずに諦める（例外にしない）。
    """

    def test_start_slide_skips_the_size_fix(self, gen, mocker):
        """見出しスライドの枠補正を飛ばす"""
        mocker.patch("processors.find_body_placeholder", return_value=None)

        gen.generate("## 中身\n本文\n", "/dev/null")

        assert gen.prs.slides[0].shapes  # 例外にならずスライドは作られる

    def test_hr_slide_skips_the_size_fix(self, gen, mocker):
        """水平線スライドの枠補正を飛ばす"""
        mocker.patch("processors.find_body_placeholder", return_value=None)

        gen.generate("## 中身\n本文\n\n---\n\n続き\n", "/dev/null")

        assert len(gen.prs.slides) == 2

    def test_shrink_body_shape_does_nothing(self, gen_with_slide, mocker):
        """本文枠の縮小も何もしない"""
        mocker.patch("utils.find_body_placeholder", return_value=None)

        shrink_body_shape(gen_with_slide.current_slide, Inches(4.0))  # 例外にならない


class TestUseTemplateFonts:
    """テンプレート（テーマ）のフォントに任せる"""

    def _title_run(self, gen):
        return gen.prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]

    def test_config_fonts_are_applied_by_default(self, base_config, tmp_path):
        """既定では config.yaml のフォントが当たる"""
        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n", str(tmp_path / "out.pptx"))

        assert self._title_run(gen).font.name == "Meiryo"

    def test_template_fonts_are_kept(self, base_config, tmp_path):
        """use_template_fonts: true では書式を指定せず、テンプレートに任せる"""
        base_config["slides"]["use_template_fonts"] = True
        gen = PPTXGenerator(base_config)
        gen.generate("# 表紙\n", str(tmp_path / "out.pptx"))

        run = self._title_run(gen)
        assert run.font.name is None
        assert run.font.size is None

    def test_code_fonts_are_preserved(self, base_config, tmp_path):
        """等幅フォントの指定は可読性のために残す"""
        base_config["slides"]["use_template_fonts"] = True
        gen = PPTXGenerator(base_config)

        assert "inline_code" in gen.fonts_conf
        assert "body" not in gen.fonts_conf


class TestTemplateConfigValidation:
    """テンプレート関連の設定の検証"""

    def _result(self, config):
        return validate_config(config)

    def test_layouts_must_be_strings(self):
        (error,) = self._result({"slides": {"layouts": {"title": 1}}}).errors
        assert "slides.layouts.title" in error

    def test_layouts_must_be_a_mapping(self):
        (error,) = self._result({"slides": {"layouts": "Title Slide"}}).errors
        assert "入れ子" in error

    def test_unknown_layout_kind_warns(self):
        (warning,) = self._result({"slides": {"layouts": {"titel": "a"}}}).warnings
        assert "'title' の誤りではありませんか？" in warning

    def test_use_template_fonts_must_be_boolean(self):
        (error,) = self._result({"slides": {"use_template_fonts": "yes"}}).errors
        assert "true か false" in error

    def test_valid_settings_pass(self):
        result = self._result({
            "slides": {
                "layouts": {"title": "Title Slide", "content": "Title and Content"},
                "use_template_fonts": True,
            }
        })
        assert result.errors == [] and result.warnings == []


class TestTemplateErrorInCli:
    """テンプレートの問題をCLIで分かりやすく伝える"""

    def test_main_reports_template_error(self, tmp_path, capsys, template_factory):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")
        conf = tmp_path / "c.yaml"
        conf.write_text(
            "slides:\n"
            f"  template_path: '{template_factory()}'\n"
            "  layouts:\n"
            "    content: '無いレイアウト'\n",
            encoding="utf-8",
        )

        assert main([str(md), "-o", str(tmp_path / "o.pptx"), "-c", str(conf)]) == 1

        captured = capsys.readouterr().out
        assert "無いレイアウト" in captured
        assert "使用できるレイアウト" in captured

    def test_error_points_at_the_template_not_the_config(self, tmp_path, capsys, template_factory):
        """設定ファイルではなくテンプレートの場所を示す（GUIでは設定が一時ファイルのため）"""
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")
        template = template_factory()
        conf = tmp_path / "tmp1234.yaml"
        conf.write_text(
            f"slides:\n  template_path: '{template}'\n  layouts:\n    content: '無いレイアウト'\n",
            encoding="utf-8",
        )

        assert main([str(md), "-o", str(tmp_path / "o.pptx"), "-c", str(conf)]) == 1

        captured = capsys.readouterr().out
        assert str(template) in captured
        assert "tmp1234.yaml" not in captured


# =============================================================================
# GUI (gui_config.py / gui_deps.py / gui_runner.py)
#
# gui.py 本体は tkinter に依存し、環境によっては import すらできないため、
# 画面に依存しない処理をこれらのモジュールに分けてテストする。
# =============================================================================


class TestGuiBuildConfig:
    """画面の入力内容から config.yaml 相当の設定を組み立てる"""

    def test_default_settings_produce_valid_config(self):
        config = build_config(GuiSettings())

        result = validate_config(config)
        assert result.errors == [] and result.warnings == []

    def test_template_settings_produce_valid_config(self, template_factory):
        settings = GuiSettings(
            use_template=True,
            template_path=str(template_factory()),
            layout_title="Title Slide",
            layout_content="Title and Content",
            footer_text="社外秘",
            footer_date=True,
        )

        config = build_config(settings)

        assert validate_config(config).errors == []
        assert config["slides"]["template_path"] == settings.template_path
        assert config["slides"]["layouts"] == {
            "title": "Title Slide",
            "content": "Title and Content",
        }
        # テンプレートの画角を尊重するため、画角は書き出さない
        assert "layout" not in config["slides"]

    def test_aspect_is_written_without_template(self):
        config = build_config(GuiSettings(aspect="4:3"))

        assert config["slides"]["layout"] == "4:3"
        assert "template_path" not in config["slides"]

    def test_template_path_ignored_when_checkbox_off(self):
        config = build_config(GuiSettings(use_template=False, template_path="company.pptx"))

        assert "template_path" not in config["slides"]

    def test_layouts_omitted_when_names_are_blank(self):
        config = build_config(GuiSettings(use_template=True, template_path="t.pptx"))

        assert "layouts" not in config["slides"]

    def test_fonts_omitted_when_using_template_fonts(self):
        settings = GuiSettings(
            use_template=True, template_path="t.pptx", use_template_fonts=True
        )

        assert "fonts" not in build_config(settings)

    def test_blank_font_fields_are_omitted(self):
        settings = GuiSettings()
        settings.fonts["body"] = FontSetting(name="  ", size_pt=0)

        fonts = build_config(settings)["fonts"]

        assert "body" not in fonts

    def test_inline_code_follows_code_block_font(self):
        fonts = build_config(GuiSettings())["fonts"]

        assert fonts["inline_code"] == {"name": fonts["code_block"]["name"]}

    def test_footer_omitted_when_empty(self):
        assert "footer" not in build_config(GuiSettings())["slides"]

    def test_footer_written_when_specified(self):
        settings = GuiSettings(footer_text=" 社外秘 ", footer_date=True, footer_on_title=True)

        footer = build_config(settings)["slides"]["footer"]

        assert footer == {"text": "社外秘", "date": True, "show_on_title": True}

    def test_endpoint_omitted_unless_kroki(self):
        config = build_config(GuiSettings(mermaid_renderer="off"))

        assert config["mermaid"] == {"renderer": "off"}

    def test_theme_colors_are_written_as_lists(self):
        config = build_config(GuiSettings(accent_color=(1, 2, 3)))

        assert config["theme"]["accent_color"] == [1, 2, 3]


class TestGuiSettingsFromConfig:
    """config.yaml を画面の入力内容に読み込む"""

    def test_round_trip_keeps_values(self):
        original = GuiSettings(
            use_template=True,
            template_path="company.pptx",
            h3_as="slide",
            show_slide_number=False,
            use_template_fonts=False,
            layout_title="表紙",
            layout_content="本文",
            footer_text="社外秘",
            footer_date=True,
            footer_on_title=True,
            accent_color=(10, 20, 30),
            text_color=(40, 50, 60),
            code_bg_color=(70, 80, 90),
            mermaid_renderer="local",
            image_height_inches=4.0,
            image_downscale=False,
            image_dpi=96,
        )

        restored = settings_from_config(build_config(original))

        for field_name in (
            "use_template", "template_path", "h3_as", "show_slide_number",
            "layout_title", "layout_content", "footer_text", "footer_date",
            "footer_on_title", "accent_color", "text_color", "code_bg_color",
            "mermaid_renderer", "image_height_inches", "image_downscale", "image_dpi",
        ):
            assert getattr(restored, field_name) == getattr(original, field_name), field_name

    def test_reads_actual_config_yaml(self):
        with open("config.yaml", encoding="utf-8") as f:
            settings = settings_from_config(yaml.safe_load(f))

        assert settings.aspect == "16:9"
        assert settings.fonts["title_h1"].name == "Meiryo"
        assert settings.fonts["title_h1"].size_pt == 44
        assert settings.accent_color == (0, 112, 192)

    def test_invalid_values_fall_back_to_defaults(self):
        config = {
            "slides": {"layout": "A5", "h3_as": 3, "show_slide_number": "yes"},
            "theme": {"accent_color": [0, 300, 0]},
            "images": {"dpi": "高い"},
            "mermaid": {"renderer": "unknown"},
            "fonts": {"body": {"size_pt": "18"}},
        }

        settings = settings_from_config(config)
        defaults = GuiSettings()

        assert settings.aspect == defaults.aspect
        assert settings.h3_as == defaults.h3_as
        assert settings.show_slide_number == defaults.show_slide_number
        assert settings.accent_color == defaults.accent_color
        assert settings.image_dpi == defaults.image_dpi
        assert settings.mermaid_renderer == defaults.mermaid_renderer
        assert settings.fonts["body"].size_pt == defaults.fonts["body"].size_pt

    def test_empty_config_keeps_defaults(self):
        assert settings_from_config(None) == GuiSettings()

    def test_footer_date_as_text_turns_the_checkbox_on(self):
        settings = settings_from_config({"slides": {"footer": {"date": "2026-08-18"}}})

        assert settings.footer_date is True


class TestGuiInputValidation:
    """変換を始める前の入力チェック"""

    def test_default_output_path_follows_input(self):
        assert default_output_path("docs/資料.md") == "docs/資料.pptx"
        assert default_output_path("/tmp/no.extension/report") == "/tmp/no.extension/report.pptx"
        assert default_output_path("") == "output.pptx"

    def test_valid_settings_have_no_errors(self, tmp_path):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")

        settings = GuiSettings(input_path=str(md), output_path=str(tmp_path / "out.pptx"))

        assert validate_inputs(settings) == []

    def test_missing_input_is_reported(self):
        errors = validate_inputs(GuiSettings(output_path="out.pptx"))

        assert any("Markdownファイル" in message for message in errors)

    def test_missing_file_is_reported(self, tmp_path):
        settings = GuiSettings(
            input_path=str(tmp_path / "無い.md"), output_path="out.pptx"
        )

        assert any("見つかりません" in message for message in validate_inputs(settings))

    def test_missing_template_is_reported(self, tmp_path):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")
        settings = GuiSettings(
            input_path=str(md), output_path="out.pptx", use_template=True
        )

        assert any("テンプレート" in message for message in validate_inputs(settings))

    def test_missing_template_file_is_reported(self, tmp_path):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")
        settings = GuiSettings(
            input_path=str(md),
            output_path="out.pptx",
            use_template=True,
            template_path=str(tmp_path / "無い.pptx"),
        )

        assert any("見つかりません" in message for message in validate_inputs(settings))

    def test_output_extension_is_checked(self, tmp_path):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")

        settings = GuiSettings(input_path=str(md), output_path="out.ppt")

        assert any(".pptx" in message for message in validate_inputs(settings))

    def test_blank_output_is_reported(self, tmp_path):
        md = tmp_path / "in.md"
        md.write_text("## 見出し\n", encoding="utf-8")

        settings = GuiSettings(input_path=str(md), output_path="  ")

        assert any("出力ファイル名" in message for message in validate_inputs(settings))


# =============================================================================
# テンプレート作成 (make_template.py)
# =============================================================================


class TestMakeTemplate:
    """完成した資料から書式だけのテンプレートを作る"""

    def _deck(self, tmp_path, slides=3, name="deck.pptx"):
        prs = Presentation()
        for _ in range(slides):
            prs.slides.add_slide(prs.slide_layouts[1])
        path = tmp_path / name
        prs.save(str(path))
        return path

    def test_slides_are_removed(self, tmp_path):
        deck = self._deck(tmp_path)
        out = tmp_path / "template.pptx"

        removed = make_template.create_template(str(deck), str(out))

        assert removed == 3
        assert len(Presentation(str(out)).slides) == 0

    def test_layouts_and_slide_size_are_kept(self, tmp_path):
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        prs.slides.add_slide(prs.slide_layouts[0])
        deck = tmp_path / "deck.pptx"
        prs.save(str(deck))
        out = tmp_path / "template.pptx"

        make_template.create_template(str(deck), str(out))

        template = Presentation(str(out))
        assert len(template.slide_layouts) == len(prs.slide_layouts)
        assert template.slide_width == Inches(13.333)

    def test_generated_template_is_usable(self, base_config, tmp_path):
        """作ったテンプレートで変換すると、生成した分だけのスライドになる"""
        deck = self._deck(tmp_path, slides=5)
        template = tmp_path / "template.pptx"
        make_template.create_template(str(deck), str(template))

        base_config["slides"]["template_path"] = str(template)
        gen = PPTXGenerator(base_config)
        gen.generate("## 見出し\n本文\n", str(tmp_path / "out.pptx"))

        assert len(gen.prs.slides) == 1

    def test_overwriting_the_source_is_refused(self, tmp_path):
        deck = self._deck(tmp_path)

        with pytest.raises(ValueError, match="上書き"):
            make_template.create_template(str(deck), str(deck))

        assert len(Presentation(str(deck)).slides) == 3

    def test_default_output_path(self):
        assert make_template.default_output_path("資料.pptx") == "資料-template.pptx"
        assert make_template.default_output_path("deck") == "deck-template.pptx"

    def test_cli_reports_success(self, tmp_path, capsys):
        deck = self._deck(tmp_path, slides=2)

        assert make_template.main([str(deck)]) == 0

        assert "スライド2枚" in capsys.readouterr().out
        assert (tmp_path / "deck-template.pptx").exists()

    def test_cli_accepts_output_option(self, tmp_path):
        deck = self._deck(tmp_path)
        out = tmp_path / "別名.pptx"

        assert make_template.main([str(deck), "-o", str(out)]) == 0
        assert out.exists()

    def test_cli_reports_missing_input(self, tmp_path, capsys):
        assert make_template.main([str(tmp_path / "無い.pptx")]) == 1

        assert "見つかりません" in capsys.readouterr().out

    def test_cli_refuses_to_overwrite_the_source(self, tmp_path, capsys):
        deck = self._deck(tmp_path)

        assert make_template.main([str(deck), "-o", str(deck)]) == 1

        assert "上書き" in capsys.readouterr().out

    def test_cli_reports_broken_file(self, tmp_path, capsys):
        broken = tmp_path / "broken.pptx"
        broken.write_text("これはPPTXではありません", encoding="utf-8")

        assert make_template.main([str(broken)]) == 1

        assert "テンプレートを作成できませんでした" in capsys.readouterr().out

    def test_cli_reports_locked_output(self, tmp_path, capsys, monkeypatch):
        deck = self._deck(tmp_path)
        monkeypatch.setattr(
            make_template, "create_template",
            MagicMock(side_effect=PermissionError()),
        )

        assert make_template.main([str(deck)]) == 1

        assert "書き込めません" in capsys.readouterr().out


class TestGuiDependencyCheck:
    """起動前の依存ライブラリ確認（gui_deps.py）"""

    def test_nothing_is_missing_in_this_environment(self):
        """テストが動く環境では、必要なライブラリは揃っている"""
        assert gui_deps.missing_dependencies() == []

    def test_reports_pip_names(self, mocker):
        """import名ではなく、pipで指定する名前を返す"""
        mocker.patch("gui_deps.importlib.util.find_spec", return_value=None)

        assert "python-pptx" in gui_deps.missing_dependencies()
        assert "pptx" not in gui_deps.missing_dependencies()

    def test_message_shows_how_to_install(self):
        """案内文に、実行中のPythonでの導入手順が入る"""
        message = gui_deps.dependency_message(["PyYAML"], "/app", "/usr/bin/python3")

        assert "PyYAML" in message
        assert "/usr/bin/python3 -m pip install -r /app/requirements.txt" in message
        assert "-m venv .venv" in message

    def test_exit_if_missing_passes_when_complete(self):
        """揃っていれば何もしない"""
        gui_deps.exit_if_missing("/app")

    def test_exit_if_missing_stops_with_guidance(self, mocker):
        """足りなければ案内を出して終了する"""
        mocker.patch("gui_deps.missing_dependencies", return_value=["Pillow"])
        stream = io.StringIO()

        with pytest.raises(SystemExit) as excinfo:
            gui_deps.exit_if_missing("/app", "/usr/bin/python3", stream)

        assert excinfo.value.code == 1
        assert "Pillow" in stream.getvalue()


class TestGuiColorHelpers:
    """画面の色見本と設定値の変換（gui_runner.py）"""

    def test_rgb_to_hex(self):
        assert gui_runner.rgb_to_hex((0, 112, 192)) == "#0070c0"

    def test_hex_to_rgb(self):
        assert gui_runner.hex_to_rgb("#0070c0") == (0, 112, 192)

    def test_round_trip(self):
        """カラーピッカーとの往復で値が変わらない"""
        assert gui_runner.hex_to_rgb(gui_runner.rgb_to_hex((40, 44, 52))) == (40, 44, 52)


class TestGuiOpenInFileManager:
    """出力先をOSのファイラーで開く（gui_runner.py）"""

    def _command(self, mocker, platform, os_name):
        mocker.patch("gui_runner.sys.platform", platform)
        mocker.patch("gui_runner.os.name", os_name)
        run = mocker.patch("gui_runner.subprocess.run")

        gui_runner.open_in_file_manager(os.path.join("out", "deck.pptx"))

        return run.call_args[0][0]

    def test_macos_selects_the_file(self, mocker):
        assert self._command(mocker, "darwin", "posix")[:2] == ["open", "-R"]

    def test_windows_selects_the_file(self, mocker):
        assert self._command(mocker, "win32", "nt")[0] == "explorer"

    def test_linux_opens_the_folder(self, mocker):
        """xdg-open にファイルを渡すと関連付けアプリが開いてしまうため、フォルダを渡す"""
        command = self._command(mocker, "linux", "posix")

        assert command == ["xdg-open", "out"]

    def test_linux_falls_back_to_current_folder(self, mocker):
        mocker.patch("gui_runner.sys.platform", "linux")
        mocker.patch("gui_runner.os.name", "posix")
        run = mocker.patch("gui_runner.subprocess.run")

        gui_runner.open_in_file_manager("deck.pptx")

        assert run.call_args[0][0] == ["xdg-open", "."]


class TestGuiConversionSetup:
    """CLIを呼び出す準備（gui_runner.py）"""

    def test_template_path_becomes_absolute(self, tmp_path, monkeypatch):
        """相対パスのままだと、作業フォルダを移す変換でテンプレートを見失う"""
        monkeypatch.chdir(tmp_path)
        settings = GuiSettings(use_template=True, template_path="theme.pptx")

        config = gui_runner.settings_to_config(settings)

        assert config["slides"]["template_path"] == str(tmp_path / "theme.pptx")

    def test_template_is_omitted_when_unused(self):
        settings = GuiSettings(use_template=False, template_path="theme.pptx")

        assert "template_path" not in gui_runner.settings_to_config(settings)["slides"]

    def test_temp_config_is_readable_yaml(self):
        """日本語を含む設定が、そのまま読み戻せる形で書かれる"""
        path = gui_runner.write_temp_config({"slides": {"footer": {"text": "社外秘"}}})
        try:
            with open(path, encoding="utf-8") as f:
                assert yaml.safe_load(f)["slides"]["footer"]["text"] == "社外秘"
        finally:
            os.unlink(path)

    def test_command_uses_the_running_python(self):
        """仮想環境から起動された場合に、その環境のライブラリを使う"""
        command = gui_runner.build_command("in.md", "out.pptx", "conf.yaml")

        assert command[0] == sys.executable
        assert command[1].endswith("md2pptx.py")
        assert command[2:] == ["in.md", "-o", "out.pptx", "-c", "conf.yaml"]


class TestGuiRunConversion:
    """GUIからCLIを実際に呼び出す（gui_runner.py）"""

    def _settings(self, tmp_path, text="## 中身\n本文\n"):
        source = tmp_path / "deck.md"
        source.write_text(text, encoding="utf-8")
        return GuiSettings(
            input_path=str(source),
            output_path=str(tmp_path / "out.pptx"),
            mermaid_renderer="off",
        )

    def test_converts_and_reports_progress(self, tmp_path):
        """変換が成功し、CLIの出力が1行ずつ渡る"""
        lines = []

        returncode, output_path = gui_runner.run_conversion(
            self._settings(tmp_path), lines.append
        )

        assert returncode == 0
        assert os.path.isfile(output_path)
        assert any("Success" in line for line in lines)
        assert not any(line.endswith("\n") for line in lines)

    def test_reports_failure_from_the_cli(self, tmp_path):
        """CLIが失敗した場合、終了コードとメッセージが伝わる"""
        settings = self._settings(tmp_path)
        settings = replace(settings, input_path=str(tmp_path / "missing.md"))
        lines = []

        returncode, _ = gui_runner.run_conversion(settings, lines.append)

        assert returncode == 1
        assert any("見つかりません" in line for line in lines)

    def test_temp_config_is_removed(self, tmp_path, mocker):
        """一時設定ファイルは変換後に消える"""
        created = []
        original = gui_runner.write_temp_config
        mocker.patch(
            "gui_runner.write_temp_config",
            side_effect=lambda config: created.append(original(config)) or created[-1],
        )

        gui_runner.run_conversion(self._settings(tmp_path), lambda _line: None)

        assert created and not os.path.exists(created[0])

    def test_unexpected_failure_does_not_raise(self, tmp_path, mocker):
        """予期せぬ失敗でも例外にせず、画面を操作不能にしない"""
        mocker.patch("gui_runner.subprocess.Popen", side_effect=OSError("起動できません"))
        lines = []

        returncode, _ = gui_runner.run_conversion(self._settings(tmp_path), lines.append)

        assert returncode == 1
        assert "変換を実行できませんでした" in lines[0]

    def test_temp_config_removal_failure_is_ignored(self, tmp_path, mocker):
        """一時ファイルを消せなくても、変換結果は返す"""
        mocker.patch("gui_runner.os.unlink", side_effect=OSError())

        returncode, _ = gui_runner.run_conversion(self._settings(tmp_path), lambda _l: None)

        assert returncode == 0


class TestBulletMarking:
    """箇条書きの行頭記号を段落へ明示的に書き込む

    本文プレースホルダーでは行頭記号がレイアウトから継承され、平文の段落にも
    同じ記号が付く。XML上は区別できないため、逆変換で平文と箇条書きが
    見分けられなかった。継承値と同じ記号を書き戻すことで、見た目を変えずに
    「箇条書きである」と残す。
    """

    def _bullets(self, gen):
        """本文の各段落の (階層, テキスト, 行頭記号の文字) を並べる"""
        result = []
        for p in gen.current_body.paragraphs:
            if not p.text.strip():
                continue
            p_pr = p._element.find(qn("a:pPr"))
            char = None
            if p_pr is not None:
                element = p_pr.find(qn("a:buChar"))
                char = element.get("char") if element is not None else None
            result.append((p.level, p.text, char))
        return result

    def _master_char(self, gen, level):
        """マスターが定める、その階層の行頭記号"""
        body_style = gen.prs.slide_master._element.find(qn("p:txStyles")).find(
            qn("p:bodyStyle")
        )
        return body_style.find(qn(f"a:lvl{level + 1}pPr")).find(qn("a:buChar")).get("char")

    def test_list_items_are_marked(self, gen_with_slide):
        """箇条書きにだけ行頭記号が書き込まれる"""
        gen_with_slide.generate("## 見出し\n\n平文です\n\n* 項目\n", "/dev/null")
        gen = gen_with_slide

        marks = {text: char for _level, text, char in self._bullets(gen)}
        assert marks["平文です"] is None
        assert marks["項目"] == self._master_char(gen, 0)

    def test_marked_character_matches_the_template(self, gen_with_slide):
        """書き込む記号は継承値と同じにする（見た目を変えないため）"""
        gen_with_slide.generate("## 見出し\n\n* 親\n    * 子\n", "/dev/null")
        gen = gen_with_slide

        chars = {level: char for level, _text, char in self._bullets(gen)}
        assert chars[0] == self._master_char(gen, 0)
        assert chars[1] == self._master_char(gen, 1)
        assert chars[0] != chars[1]  # 階層ごとに記号が違うことの確認

    def test_bullet_font_is_carried_over(self, gen_with_slide):
        """記号の書体も継承値をそのまま使う"""
        gen_with_slide.generate("## 見出し\n\n* 項目\n", "/dev/null")

        p = [p for p in gen_with_slide.current_body.paragraphs if p.text == "項目"][0]
        font = p._element.find(qn("a:pPr")).find(qn("a:buFont"))
        assert font is not None and font.get("typeface") == "Arial"

    def test_ordered_items_keep_auto_numbering(self, gen_with_slide):
        """番号付きリストは従来どおり自動採番のまま"""
        gen_with_slide.generate("## 見出し\n\n1. 手順\n", "/dev/null")

        p = [p for p in gen_with_slide.current_body.paragraphs if p.text == "手順"][0]
        assert p._element.find(qn("a:pPr")).find(qn("a:buAutoNum")) is not None
        assert p._element.find(qn("a:pPr")).find(qn("a:buChar")) is None

    def test_nothing_is_written_when_the_template_hides_bullets(self, gen_with_slide, mocker):
        """行頭記号を出さないテンプレートでは、記号を足さない"""
        mocker.patch("utils.inherited_bullet", return_value=None)

        gen_with_slide.generate("## 見出し\n\n* 項目\n", "/dev/null")

        assert all(char is None for _l, _t, char in self._bullets(gen_with_slide))


class TestInheritedBullet:
    """継承される行頭記号の解決（utils.inherited_bullet）"""

    def _body_style(self, gen):
        return gen.prs.slide_master._element.find(qn("p:txStyles")).find(qn("p:bodyStyle"))

    def _level(self, gen, level=0):
        return self._body_style(gen).find(qn(f"a:lvl{level + 1}pPr"))

    def test_finds_the_character_for_the_level(self, gen_with_slide):
        char, font = inherited_bullet(gen_with_slide.current_slide, 0)

        assert char.get("char") == "•"
        assert font.get("typeface") == "Arial"

    def test_returns_none_when_bullets_are_disabled(self, gen_with_slide):
        """buNone のテンプレートでは記号が無いので None"""
        level = self._level(gen_with_slide)
        level.remove(level.find(qn("a:buChar")))
        level.append(OxmlElement("a:buNone"))

        assert inherited_bullet(gen_with_slide.current_slide, 0) is None

    def test_returns_none_for_picture_bullets(self, gen_with_slide):
        """画像の行頭記号は文字で表せないので None"""
        level = self._level(gen_with_slide)
        level.remove(level.find(qn("a:buChar")))
        level.append(OxmlElement("a:buBlip"))

        assert inherited_bullet(gen_with_slide.current_slide, 0) is None

    def test_returns_none_when_nothing_is_defined(self, gen_with_slide):
        """どこにも定義が無ければ None（勝手に記号を作らない）"""
        body_style = self._body_style(gen_with_slide)
        body_style.getparent().remove(body_style)

        assert inherited_bullet(gen_with_slide.current_slide, 0) is None

    def test_layout_takes_precedence_over_the_master(self, gen_with_slide):
        """レイアウト側の指定があれば、そちらを優先する"""
        layout = gen_with_slide.current_slide.slide_layout
        list_style = find_body_placeholder(layout).text_frame._txBody.find(qn("a:lstStyle"))
        level = OxmlElement("a:lvl1pPr")
        char = OxmlElement("a:buChar")
        char.set("char", "▶")
        level.append(char)
        list_style.append(level)

        found, _font = inherited_bullet(gen_with_slide.current_slide, 0)

        assert found.get("char") == "▶"


class TestBulletsFromOtherTools:
    """md2pptx 以外で作られた資料の扱い（後方互換）"""

    def _deck(self, text_frame_filler):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "見出し"
        text_frame_filler(find_body_placeholder(slide).text_frame)
        return prs

    def test_unmarked_paragraphs_stay_bullets(self):
        """行頭記号の指定が無い資料は、従来どおり箇条書きとして書き出す

        本文枠の既定でどの段落にも記号が付いて見えるため、
        平文として書き出すと見た目と食い違う。
        """
        def fill(tf):
            tf.text = "一つ目"
            tf.add_paragraph().text = "二つ目"

        markdown = extract(self._deck(fill)).markdown

        assert "* 一つ目" in markdown
        assert "* 二つ目" in markdown

    def test_explicit_marking_switches_the_interpretation(self):
        """1つでも明示があれば、指定の無い段落は平文として扱う"""
        def fill(tf):
            tf.text = "平文"
            marked = tf.add_paragraph()
            marked.text = "箇条書き"
            char = OxmlElement("a:buChar")
            char.set("char", "•")
            utils._set_bullet(marked, char)

        markdown = extract(self._deck(fill)).markdown

        assert "* 箇条書き" in markdown
        assert "* 平文" not in markdown
        assert "\n平文\n" in markdown


# =============================================================================
# font_metrics.py（フォントの実測値）
# =============================================================================


def build_test_font(directory, family="Metrics Test"):
    """既知のメトリクスを持つフォントを作る

    実在するフォントは環境によって有無が変わるため、テストでは合成する。
    行送り比 = (800 + 200 + 200) / 1000 = 1.2、'A' の送り幅 = 0.6em、
    'あ' = 1.0em、それ以外の文字は持たない。
    """
    from fontTools.fontBuilder import FontBuilder
    from fontTools.pens.ttGlyphPen import TTGlyphPen

    glyphs = [".notdef", "A", "uni3042"]
    builder = FontBuilder(1000, isTTF=True)
    builder.setupGlyphOrder(glyphs)
    builder.setupCharacterMap({0x41: "A", 0x3042: "uni3042"})
    builder.setupGlyf({name: TTGlyphPen(None).glyph() for name in glyphs})
    builder.setupHorizontalMetrics(
        {".notdef": (500, 0), "A": (600, 0), "uni3042": (1000, 0)}
    )
    builder.setupHorizontalHeader(ascent=800, descent=-200, lineGap=200)
    builder.setupNameTable({
        "familyName": family, "styleName": "Regular",
        "fullName": f"{family} Regular", "psName": f"{family}-Regular",
        "version": "1.0", "uniqueFontIdentifier": family,
    })
    builder.setupOS2()
    builder.setupPost()
    path = os.path.join(str(directory), "TestFont.ttf")
    builder.save(path)
    return path


@pytest.fixture
def installed_font(tmp_path, monkeypatch):
    """合成したフォントだけが入っている環境を用意する"""
    build_test_font(tmp_path)
    monkeypatch.setattr(font_metrics, "font_directories", lambda: [str(tmp_path)])
    font_metrics._font_index.cache_clear()
    font_metrics.metrics_for.cache_clear()
    yield "Metrics Test"
    font_metrics._font_index.cache_clear()
    font_metrics.metrics_for.cache_clear()


class TestFontLookup:
    """フォントファイルの探索"""

    def test_name_is_matched_loosely(self):
        """空白や大小の違いは無視して照合する"""
        assert font_metrics.normalize("Yu Gothic") == font_metrics.normalize("yugothic")

    def test_finds_a_font_by_its_family_name(self, installed_font):
        assert font_metrics.metrics_for(installed_font) is not None

    def test_unknown_font_is_not_found(self, installed_font):
        """手元に無いフォントは None（既定値に任せる）"""
        assert font_metrics.metrics_for("存在しないフォント") is None

    def test_no_name_is_not_looked_up(self):
        assert font_metrics.metrics_for(None) is None

    def test_missing_directories_are_skipped(self, monkeypatch, tmp_path):
        monkeypatch.setattr(
            font_metrics, "font_directories", lambda: [str(tmp_path / "無い")]
        )
        assert font_metrics.font_files() == []

    def test_broken_font_files_are_skipped(self, tmp_path, monkeypatch):
        """壊れたファイルがあっても探索は続く"""
        (tmp_path / "broken.ttf").write_text("フォントではありません", encoding="utf-8")
        build_test_font(tmp_path)
        monkeypatch.setattr(font_metrics, "font_directories", lambda: [str(tmp_path)])
        font_metrics._font_index.cache_clear()
        font_metrics.metrics_for.cache_clear()

        assert font_metrics.metrics_for("Metrics Test") is not None

    def test_unreadable_font_falls_back(self, installed_font, mocker):
        """索引には載ったが開けない場合も既定値に落とす"""
        assert font_metrics.metrics_for(installed_font) is not None  # 先に索引を作る

        mocker.patch("font_metrics._open_fonts", side_effect=OSError())
        font_metrics.metrics_for.cache_clear()  # 索引は残したまま読み込みだけやり直す

        assert font_metrics.metrics_for(installed_font) is None

    def test_unreadable_name_table_is_skipped(self):
        """名前のテーブルを読めないフォントは、名前なしとして扱う"""
        class Broken:
            def __getitem__(self, key):
                raise KeyError(key)

        assert font_metrics._names_of(Broken()) == []

    def test_works_without_fonttools(self, mocker):
        """fontTools が無い環境でも変換自体は動く"""
        mocker.patch("font_metrics.TTFont", None)
        font_metrics.metrics_for.cache_clear()
        font_metrics._font_index.cache_clear()
        try:
            assert font_metrics.metrics_for("Metrics Test") is None
            assert font_metrics._font_index() == {}
        finally:
            font_metrics.metrics_for.cache_clear()
            font_metrics._font_index.cache_clear()

    def test_font_directories_differ_by_platform(self, monkeypatch):
        monkeypatch.setattr(font_metrics.sys, "platform", "darwin")
        assert any("Library/Fonts" in d for d in font_metrics.font_directories())

        monkeypatch.setattr(font_metrics.sys, "platform", "win32")
        monkeypatch.setattr(font_metrics.os, "name", "nt")
        monkeypatch.setenv("LOCALAPPDATA", r"C:\Users\me\AppData\Local")
        directories = font_metrics.font_directories()
        assert any(d.endswith("Fonts") for d in directories)
        assert any("AppData" in d for d in directories)  # 利用者ごとのフォントも見る

        monkeypatch.setattr(font_metrics.sys, "platform", "linux")
        monkeypatch.setattr(font_metrics.os, "name", "posix")
        assert "/usr/share/fonts" in font_metrics.font_directories()


class TestMeasuredMetrics:
    """実測値を使った概算"""

    def test_line_height_comes_from_the_font(self, installed_font):
        """行の高さはフォントの ascent + descent + lineGap から決まる"""
        assert line_height_ratio(installed_font) == pytest.approx(1.2)

    def test_default_is_used_without_a_font(self):
        assert line_height_ratio("存在しないフォント") == LINE_HEIGHT_RATIO

    def test_character_width_comes_from_the_font(self, installed_font):
        """'A' は 0.5em ではなく実測の 0.6em として扱う"""
        width = estimate_text_width_pt("A", 100, installed_font)

        assert width == pytest.approx(60.0)

    def test_characters_absent_from_the_font_fall_back(self, installed_font):
        """フォントに無い文字は全角/半角から概算する"""
        assert estimate_text_width_pt("Z", 100, installed_font) == pytest.approx(50.0)
        assert estimate_text_width_pt("あ", 100, installed_font) == pytest.approx(100.0)

    def test_line_count_uses_measured_widths(self, installed_font):
        """折り返し行数の判定にも実測値が効く"""
        # 'A' 10文字 = 600pt。550pt の幅なら2行になる（既定値の 0.5em なら1行）
        assert estimate_line_count("A" * 10, 100, 550, installed_font) == 2
        assert estimate_line_count("A" * 10, 100, 550) == 1

    def test_paragraph_height_uses_the_font(self, installed_font):
        """段落の高さも実測の行送りで計算される"""
        paragraph = ParagraphMetrics(text="A", font_size_pt=100, font_name=installed_font)

        assert estimate_height_pt([paragraph], 1000) == pytest.approx(120.0)

    def test_table_row_height_uses_the_font(self, installed_font):
        """表の行の高さも同様"""
        row = TableRowMetrics(texts=["A"], font_size_pt=100, font_name=installed_font)

        (height,) = estimate_row_heights_pt([row], 1000, 0.0)
        assert height == pytest.approx(120.0)


class TestFontNameIsCarriedOver:
    """書体の指定が概算まで届いているか"""

    def test_paragraph_font_is_picked_up(self, gen_with_slide):
        """段落から書体名を取り出す"""
        gen_with_slide.generate("## 見出し\n\n本文です\n", "/dev/null")
        paragraph = [
            p for p in gen_with_slide.current_body.paragraphs if p.text == "本文です"
        ][0]

        assert utils._paragraph_metrics(paragraph).font_name == paragraph.runs[0].font.name

    def test_table_font_is_picked_up(self):
        """表の設定から書体名を取り出す"""
        rows = parse_md("| A |\n|---|\n| 1 |").find_all("tr")
        metrics = table_row_metrics(
            rows, {"name": "見出し書体", "size_pt": 14}, {"name": "本文書体", "size_pt": 12}
        )

        assert [row.font_name for row in metrics] == ["見出し書体", "本文書体"]
