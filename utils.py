"""描画のヘルパー関数群

このモジュールは PPTXGenerator に依存しない純粋なヘルパーのみを置く。
スライドや設定を引数で受け取り、ジェネレーターの状態は参照しない。
ジェネレーターの状態に依存する処理は processors.py 側に置くこと。
"""

from __future__ import annotations

import collections
import collections.abc
import os
from io import BytesIO
from typing import IO, TYPE_CHECKING, Any, Union

from PIL import Image
from pptx.util import Emu, Inches, Length, Pt

from text_metrics import (
    DEFAULT_FONT_SIZE_PT,
    DEFAULT_LINE_SPACING,
    ParagraphMetrics,
    fit_scale,
)
from pptx.dml.color import RGBColor
from bs4 import NavigableString, Tag
from pygments import lex
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.styles import get_style_by_name

if TYPE_CHECKING:
    from pptx.shapes.picture import Picture
    from pptx.slide import Slide
    from pptx.text.text import TextFrame, _Paragraph, _Run

#: config.yaml のフォント設定（name / size_pt / bold / color_rgb）
FontConfig = dict[str, Any]

#: 画像データとして受け付ける型（ローカルパス、またはメモリ上のバイト列）
ImageSource = Union[str, IO[bytes]]

#: 画像を再サンプリングする際の既定解像度（PowerPointの「図の圧縮」の印刷品質相当）
DEFAULT_IMAGE_DPI = 150

#: JPEGを再エンコードする際の品質
JPEG_QUALITY = 90

#: インラインコードの既定フォント・色（config.yaml に inline_code が無い場合）
DEFAULT_INLINE_CODE_FONT = 'Consolas'
DEFAULT_INLINE_CODE_COLOR = [220, 20, 60]

#: コードブロック枠の内側の余白と行送り
CODE_BOX_MARGIN_INCHES = 0.2
CODE_BOX_LINE_SPACING = 1.1


def hex_to_rgb(hex_str: str | None) -> RGBColor | None:
    if not hex_str: return None
    hex_str = hex_str.lstrip('#')
    return RGBColor(*tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4)))

def apply_syntax_highlight(
    p: _Paragraph, text: str, language: str | None, font_conf: FontConfig | None
) -> None:
    try:
        if language:
            lexer = get_lexer_by_name(language, stripall=False)
        else:
            lexer = guess_lexer(text)
    except Exception:
        # 未知の言語や判定不能なコードはハイライトなしのプレーンテキストとして扱う
        lexer = get_lexer_by_name('text')

    style = get_style_by_name('monokai') # 濃い背景に合うmonokaiを使用

    for token, content in lex(text, lexer):
        if not content: continue
        run = p.add_run()
        run.text = content
        apply_font_style(run, font_conf)

        token_style = style.style_for_token(token)
        if token_style['color']:
            run.font.color.rgb = hex_to_rgb(token_style['color'])
        else:
            run.font.color.rgb = RGBColor(248, 248, 242) # 背景に合う白をデフォルトに

        if token_style['bold']: run.font.bold = True
        if token_style['italic']: run.font.italic = True

def apply_font_style(run: _Run, font_config: FontConfig | None) -> None:
    """フォントスタイルの適用"""
    if not font_config: return
    font = run.font
    if 'name' in font_config: font.name = font_config['name']
    if 'size_pt' in font_config: font.size = Pt(font_config['size_pt'])
    if 'bold' in font_config: font.bold = font_config['bold']
    if 'color_rgb' in font_config:
        rgb = font_config['color_rgb']
        font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

def _rewind(img_data: ImageSource) -> ImageSource:
    """ファイルライクな画像データを先頭に巻き戻して返す"""
    if hasattr(img_data, 'seek'):
        img_data.seek(0)
    return img_data

def _source_size(img_data: ImageSource) -> int | None:
    """元画像のバイト数を返す（取得できない場合は None）"""
    try:
        if isinstance(img_data, str):
            return os.path.getsize(img_data)
        position = img_data.tell()
        img_data.seek(0, os.SEEK_END)
        size = img_data.tell()
        img_data.seek(position)
        return size
    except Exception:
        return None

def downscale_image(
    img_data: ImageSource,
    max_width: Length,
    max_height: Length,
    dpi: int = DEFAULT_IMAGE_DPI,
) -> ImageSource:
    """表示サイズに対して過大な画像を再サンプリングし、埋め込みサイズを削減する

    スライド上の表示サイズと指定DPIから必要な画素数を求め、それを超える分だけ縮小する。
    拡大は行わない（元が小さい画像はそのまま返す）。

    再エンコードでかえってサイズが増える画像（色数の少ない図版など、元のPNGが
    よく圧縮されているケース）では元データをそのまま使う。
    画像として読めない・保存できない場合も同様に元データを返し、変換自体は止めない。
    """
    target_w = max(1, int(Emu(int(max_width)).inches * dpi))
    target_h = max(1, int(Emu(int(max_height)).inches * dpi))

    try:
        _rewind(img_data)

        with Image.open(img_data) as img:
            # 既に十分小さい画像は再エンコードしない（無駄な劣化とCPUを避ける）
            if img.width <= target_w and img.height <= target_h:
                return _rewind(img_data)

            image_format = img.format or 'PNG'
            resized = img.copy()
            resized.thumbnail((target_w, target_h), Image.Resampling.LANCZOS)

            buffer = BytesIO()
            save_options: dict[str, Any] = {'dpi': (dpi, dpi)}
            if image_format == 'JPEG':
                save_options['quality'] = JPEG_QUALITY
            elif image_format == 'PNG':
                save_options['optimize'] = True
            resized.save(buffer, format=image_format, **save_options)

        original_size = _source_size(img_data)
        if original_size is not None and buffer.getbuffer().nbytes >= original_size:
            return _rewind(img_data)

        buffer.seek(0)
        return buffer
    except Exception as e:
        print(f"Warning: 画像の縮小に失敗したため元の画像を使用します: {e}")
        return _rewind(img_data)

def insert_image_fit(
    slide: Slide,
    img_data: ImageSource,
    left: Length,
    top: Length,
    max_width: Length,
    max_height: Length,
) -> Picture:
    """画像を最大枠に収まるようにアスペクト比を保って自動縮小・中央配置する"""
    pic = slide.shapes.add_picture(img_data, left, top)
    ratio_w = max_width / pic.width
    ratio_h = max_height / pic.height
    ratio = min(ratio_w, ratio_h)
    ratio = min(ratio, 1.5) # 極端な拡大を防止

    pic.width = Emu(int(pic.width * ratio))
    pic.height = Emu(int(pic.height * ratio))
    pic.left = Emu(int(left + (max_width - pic.width) / 2))
    pic.top = Emu(int(top + (max_height - pic.height) / 2))
    return pic

def add_runs_from_tag(
    element: Tag,
    paragraph: _Paragraph,
    default_font_conf: FontConfig | None,
    inline_code_conf: FontConfig | None = None,
) -> None:
    """インライン装飾を解釈しながらテキストを追加（再帰処理）"""
    code_conf = inline_code_conf or {}

    for child in element:
        if isinstance(child, NavigableString):
            text = str(child).replace('\n', ' ')
            if text.strip() or text == ' ':
                run = paragraph.add_run()
                run.text = text
                apply_font_style(run, default_font_conf)
        elif isinstance(child, Tag):
            if child.name in ['ul', 'ol', 'pre', 'img', 'table', 'blockquote']: continue
            if child.name in ['p', 'div', 'span', 'li', 'th', 'td']:
                add_runs_from_tag(child, paragraph, default_font_conf, inline_code_conf)
            else:
                run = paragraph.add_run()
                run.text = child.get_text().replace('\n', ' ')
                apply_font_style(run, default_font_conf)

                if child.name in ['strong', 'b']: run.font.bold = True
                elif child.name in ['em', 'i']: run.font.italic = True
                elif child.name == 'code':
                    run.font.name = code_conf.get('name', DEFAULT_INLINE_CODE_FONT)
                    rgb = code_conf.get('color_rgb', DEFAULT_INLINE_CODE_COLOR)
                    run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

def shrink_body_shape(
    slide: Slide,
    width: Length,
    max_height: Length | None = None,
) -> None:
    """テキスト枠を指定サイズ（EMU）に縮める（レイアウト調整用ヘルパー）

    注意: 下の left/top の自己代入は削除しないこと。プレースホルダーの位置・サイズは
    スライドレイアウトからの継承値であり、一部だけを書き換えると継承が切れて
    残りの値が 0 になる。継承値を明示的に書き戻してから変更する必要がある。
    """
    try:
        body_shape = slide.placeholders[1]
        body_shape.left, body_shape.top = body_shape.left, body_shape.top
        body_shape.width = width
        if max_height:
            body_shape.height = max_height
    except Exception:
        pass

def append_text_block(
    text_frame: TextFrame,
    content: Tag,
    reuse_first_paragraph: bool = False,
    level: int = 0,
    font_conf: FontConfig | None = None,
    inline_code_conf: FontConfig | None = None,
) -> None:
    """段落オブジェクトを追加し、テキストまたはタグ構造を書き込むヘルパー

    reuse_first_paragraph が真で、かつ枠が空のときは、既存の空段落を再利用する
    （先頭に空行が入るのを避けるため）。
    """
    is_empty = len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text
    if reuse_first_paragraph and is_empty:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    p.level = level
    p.space_after = Pt(12)  # 段落後の余白を追加してレイアウトを美しく
    p.line_spacing = 1.2    # 行間を1.2倍に設定

    add_runs_from_tag(content, p, font_conf, inline_code_conf)

def create_code_textbox(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    content: str,
    language: str | None,
    font_conf: FontConfig | None,
    background_rgb: list[int],
) -> None:
    """背景色付きのテキストボックスを作り、コードをハイライトして書き込む"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(*background_rgb)

    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(CODE_BOX_MARGIN_INCHES)
    tf.margin_top = Inches(CODE_BOX_MARGIN_INCHES)
    tf.margin_right = Inches(CODE_BOX_MARGIN_INCHES)
    tf.margin_bottom = Inches(CODE_BOX_MARGIN_INCHES)

    p = tf.paragraphs[0]
    p.line_spacing = CODE_BOX_LINE_SPACING

    apply_syntax_highlight(p, content, language, font_conf)

def _paragraph_metrics(paragraph: _Paragraph) -> ParagraphMetrics:
    """段落から高さ概算に必要な情報を取り出す"""
    sizes = [run.font.size.pt for run in paragraph.runs if run.font.size]
    if paragraph.font.size:
        sizes.append(paragraph.font.size.pt)
    font_size = max(sizes) if sizes else DEFAULT_FONT_SIZE_PT

    # line_spacing は倍率（float）か絶対値（Length）のどちらもあり得る
    spacing = paragraph.line_spacing
    line_spacing = float(spacing) if isinstance(spacing, (int, float)) else DEFAULT_LINE_SPACING

    return ParagraphMetrics(
        text=paragraph.text,
        font_size_pt=font_size,
        level=paragraph.level,
        line_spacing=line_spacing,
        space_after_pt=paragraph.space_after.pt if paragraph.space_after else 0.0,
    )


def auto_shrink_text(slide: Slide | None) -> None:
    """本文が枠に収まらない場合、フォントサイズと余白を自動で縮小する

    段落数ではなく、文字幅（全角/半角）から**折り返し後の行数**を概算して判定する。
    長い1段落が何行にも折り返してはみ出すケースを拾うため。
    """
    if not slide: return
    try:
        if len(slide.placeholders) <= 1: return
        body = slide.placeholders[1]
        if not body.has_text_frame: return

        tf = body.text_frame
        paragraphs = list(tf.paragraphs)
        if not any(p.text.strip() for p in paragraphs): return

        available_width = Emu(int(body.width - tf.margin_left - tf.margin_right)).pt
        available_height = Emu(int(body.height - tf.margin_top - tf.margin_bottom)).pt

        metrics = [_paragraph_metrics(p) for p in paragraphs]
        scale = fit_scale(metrics, available_width, available_height)
        if scale >= 1.0: return

        for p in paragraphs:
            if p.space_after:
                p.space_after = Pt(p.space_after.pt * scale)
            for run in p.runs:
                if run.font.size:
                    run.font.size = Pt(int(run.font.size.pt * scale))
    except Exception:
        pass
