"""Markdown（HTML）の各タグをスライド上の要素へ変換する処理"""

from __future__ import annotations

import math
import re
import requests
from io import BytesIO
from typing import TYPE_CHECKING, cast

from bs4 import BeautifulSoup, Tag
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Length, Pt
from pptx.dml.color import RGBColor

from mermaid_renderer import mermaid_conf, render_mermaid
from text_metrics import (
    TableRowMetrics,
    estimate_row_heights_pt,
    fit_table_scale,
    paginate_row_heights,
)
from utils import (
    DEFAULT_IMAGE_DPI,
    FontConfig,
    ImageSource,
    apply_auto_numbering,
    apply_font_style,
    disable_bullet,
    downscale_image,
    fit_shape_into,
    insert_image_fit,
    set_alt_text,
    shrink_body_shape,
    add_runs_from_tag,
    append_text_block,
    create_code_textbox,
    auto_shrink_text
)

if TYPE_CHECKING:
    from pptx.shapes.picture import Picture

    from generator import PPTXGenerator

# 画像取得（HTTP）のタイムアウト秒数
HTTP_TIMEOUT_SEC = 15

# スピーカーノートで1つのまとまりとして扱うブロック要素
NOTE_BLOCK_TAGS = ['p', 'li']

# 箇条書きのインデントレベルの上限（PowerPointの仕様）
MAX_BULLET_LEVEL = 8

# 画像を横一列に並べる上限の枚数（これを超えたらグリッド配置にする）
MAX_IMAGES_IN_A_ROW = 3

# 並べた画像どうしの間隔
IMAGE_GUTTER_INCHES = 0.2

# slides.h3_as に指定できる値。h3 をスライド内の小見出しにするか、新規スライドにするか
H3_AS_SUBHEADING = 'subheading'
H3_AS_SLIDE = 'slide'

# config.yaml に該当設定が無い場合のコードブロックの既定値
DEFAULT_CODE_BLOCK_FONT = {'name': 'Consolas', 'size_pt': 12}
DEFAULT_CODE_BG_COLOR = [40, 44, 52]

# 同、表の既定値
DEFAULT_TABLE_HEADER_FONT: FontConfig = {
    'name': 'Meiryo', 'size_pt': 14, 'bold': True, 'color_rgb': [255, 255, 255],
}
DEFAULT_TABLE_BODY_FONT: FontConfig = {'name': 'Meiryo', 'size_pt': 12}

# 分割された表の続きスライドのタイトルに付ける接尾辞
TABLE_CONTINUATION_SUFFIX = '（続き）'

# Markdownの列揃え指定と PowerPoint の配置の対応
CELL_ALIGNMENTS = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
}

# python-pptx が設定するセルの既定マージン（ポイント）
CELL_MARGIN_H_PT = 0.1 * 72
CELL_MARGIN_V_PT = 0.05 * 72

def process_heading(generator: PPTXGenerator, tag: Tag) -> None:
    """見出しタグの処理とスライド作成"""
    if generator.current_slide:
        auto_shrink_text(generator.current_slide)
        
    layout_idx = 0 if tag.name == 'h1' else 1
    generator.current_slide = generator.prs.slides.add_slide(generator.prs.slide_layouts[layout_idx])
    
    from pptx.enum.text import MSO_AUTO_SIZE
    title_shape = generator.current_slide.shapes.title
    if title_shape:
        title_shape.text_frame.word_wrap = True
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_shape.text = tag.get_text()
        
        style_key = 'title_h1' if tag.name == 'h1' else 'title_h2'
        for run in title_shape.text_frame.paragraphs[0].runs:
            apply_font_style(run, generator.fonts_conf.get(style_key, generator.fonts_conf.get('title')))

    generator.current_body = generator.current_slide.placeholders[1].text_frame
    generator.current_body.word_wrap = True
    generator.current_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    generator.current_body.text = "" 
    generator.slide_has_text = False
    generator.current_images = []

    # デフォルト枠のはみ出し補正
    if not generator.slides_conf.get('template_path'):
        try:
            layout = generator.layout
            body_shape = generator.current_slide.placeholders[1]
            # 継承値は書き込む前にすべて読み出す。一部だけ書き換えると継承が切れ、
            # 残りの値が0になってしまうため（shrink_body_shape の注意書きも参照）
            original_top = body_shape.top
            # 本文枠をコンテンツ領域に合わせ、画像・表と左右の端を揃える
            # （プレースホルダーは画角を変えても追従しないので、ここで合わせる）
            body_shape.left = layout.content_left
            body_shape.top = original_top
            body_shape.width = layout.content_width
            body_shape.height = layout.body_height_for(original_top)
        except Exception:
            pass

def h3_creates_slide(generator: PPTXGenerator) -> bool:
    """h3 を新規スライドとして扱うか（config.yaml の slides.h3_as）"""
    setting = generator.slides_conf.get('h3_as', H3_AS_SUBHEADING)
    return str(setting).lower() == H3_AS_SLIDE

def process_h3(generator: PPTXGenerator, tag: Tag) -> None:
    """H3見出しの処理

    既定ではスライド内の小見出しとして扱う。
    slides.h3_as: "slide" を指定した場合は h2 と同じく新規スライドを作る。
    """
    if not tag.get_text(strip=True): return

    if h3_creates_slide(generator):
        # 新しいスライドに移るので、直前のレイアウト指定は引き継がない
        generator.forced_layout = None
        process_heading(generator, tag)
        return

    if generator.current_body is None:
        return  # 見出しより前に現れた h3 は配置先が無い

    from pptx.util import Pt
    
    p = generator.current_body.add_paragraph()
    p.level = 0
    p.space_before = Pt(10)
    p.space_after = Pt(2)
    disable_bullet(p)  # 見出しなので行頭記号は出さない
    
    run = p.add_run()
    run.text = tag.get_text()
    
    font_conf = generator.fonts_conf.get('title_h3', {'name': 'Meiryo', 'size_pt': 20, 'bold': True})
    apply_font_style(run, font_conf)
    
    # 段落レベルでフォントサイズを固定（はみ出し防止のベース）
    if 'size_pt' in font_conf:
        p.font.size = Pt(font_conf['size_pt'])
    generator.slide_has_text = True

def process_hr(generator: PPTXGenerator, tag: Tag) -> None:
    """水平線（---）による新しいスライド（タイトルなし）の生成"""
    if generator.current_slide:
        from utils import auto_shrink_text
        auto_shrink_text(generator.current_slide)
        
    generator.current_slide = generator.prs.slides.add_slide(generator.prs.slide_layouts[1])
    
    # タイトルシェイプを削除して上部から広く使えるようにする
    if generator.current_slide.shapes.title:
        sp = generator.current_slide.shapes.title._element
        sp.getparent().remove(sp)
        
    generator.current_body = generator.current_slide.placeholders[1].text_frame
    generator.current_body.word_wrap = True
    
    from pptx.enum.text import MSO_AUTO_SIZE
    generator.current_body.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    generator.current_body.text = "" 
    generator.slide_has_text = False
    generator.current_images = []

    if not generator.slides_conf.get('template_path'):
        try:
            layout = generator.layout
            body_shape = generator.current_slide.placeholders[1]
            new_top = Inches(0.5)
            body_shape.left = layout.content_left
            body_shape.top = new_top
            body_shape.width = layout.content_width
            body_shape.height = layout.body_height_for(new_top)
        except Exception:
            pass

def extract_note_text(tag: Tag) -> str:
    """引用ブロックから、段落構造を保ったままノート用のテキストを取り出す

    tag.get_text() は段落や箇条書きの区切りを落として全体を連結してしまうため、
    ブロック単位で取り出して改行でつなぐ。
    """
    # リスト項目の中の p は、項目側でまとめて扱うため除外する
    blocks = [
        block for block in tag.find_all(NOTE_BLOCK_TAGS)
        if not (block.name == 'p' and block.find_parent('li'))
    ]
    if not blocks:
        return tag.get_text(strip=True)

    parts: list[str] = []
    for block in blocks:
        text = block.get_text().strip()
        if not text:
            continue
        if parts:
            # 箇条書きは詰めて、段落どうしは空行を挟む
            parts.append('\n' if block.name == 'li' else '\n\n')
        parts.append(text)
    return ''.join(parts)

def process_blockquote(generator: PPTXGenerator, tag: Tag) -> None:
    """スピーカーノートの処理"""
    text_frame = generator.current_slide.notes_slide.notes_text_frame
    note_text = extract_note_text(tag)
    if not note_text:
        return
    text_frame.text = text_frame.text + "\n\n" + note_text if text_frame.text else note_text

def inline_code_conf(generator: PPTXGenerator) -> FontConfig | None:
    """インラインコード（`code`）のフォント設定を取り出す"""
    return generator.fonts_conf.get('inline_code')

def append_code_textbox(
    generator: PPTXGenerator, content: str, language: str | None = None
) -> None:
    """レイアウトを決めて、背景色付きのコード枠をスライドに追加する"""
    layout = generator.layout

    if generator.slide_has_text or generator.forced_layout == '2-column':
        shrink_body_shape(generator.current_slide, layout.split_body_width)
        box = (layout.split_right_left, layout.content_top,
               layout.split_right_width, layout.content_height)
    elif generator.forced_layout == 'center':
        box = (layout.center_left, layout.content_top,
               layout.center_width, layout.content_height)
    else:
        box = (layout.content_left, layout.content_top,
               layout.content_width, layout.content_height)

    theme = generator.config.get('theme') or {}
    create_code_textbox(
        generator.current_slide, *box,
        content=content,
        language=language,
        font_conf=generator.fonts_conf.get('code_block', DEFAULT_CODE_BLOCK_FONT),
        background_rgb=theme.get('code_bg_color', DEFAULT_CODE_BG_COLOR),
    )
    generator.slide_has_text = True

def image_dpi(generator: PPTXGenerator) -> int | None:
    """埋め込み画像の解像度を返す（images.downscale が false の場合は None＝縮小しない）"""
    if not generator.images_conf.get('downscale', True):
        return None
    return generator.images_conf.get('dpi', DEFAULT_IMAGE_DPI)

def place_image(
    generator: PPTXGenerator,
    img_data: ImageSource,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
) -> Picture:
    """設定に応じて画像を縮小したうえで、指定した枠に収めて配置する"""
    dpi = image_dpi(generator)
    if dpi:
        img_data = downscale_image(img_data, width, height, dpi)
    return insert_image_fit(generator.current_slide, img_data, left, top, width, height)

def place_image_full(generator: PPTXGenerator, img_data: ImageSource) -> Picture:
    """コンテンツ領域いっぱいに図を配置する"""
    layout = generator.layout
    return place_image(
        generator, img_data,
        layout.content_left, layout.content_top, layout.content_width, layout.content_height,
    )

def place_image_split(generator: PPTXGenerator, img_data: ImageSource) -> Picture:
    """本文枠を左に縮め、図を右半分に配置する（2カラム）"""
    layout = generator.layout
    shrink_body_shape(generator.current_slide, layout.split_body_width)
    return place_image(
        generator, img_data,
        layout.split_right_left, layout.content_top,
        layout.split_right_width, layout.content_height,
    )

def image_grid(count: int) -> tuple[int, int]:
    """画像の枚数から (列数, 行数) を決める

    3枚までは横一列に並べる。スライドは横長なので、その方が1枚あたりが大きくなる。
    4枚以上は正方形に近いグリッドに収める。
    """
    if count <= MAX_IMAGES_IN_A_ROW:
        return count, 1
    columns = math.ceil(math.sqrt(count))
    return columns, math.ceil(count / columns)

def arrange_images(generator: PPTXGenerator) -> None:
    """同じスライド上の画像が重ならないように並べ直す

    画像は現れた順に1枚ずつ配置するため、何枚になるかは事前に分からない。
    そこで追加のたびに、その時点の枚数で全体を並べ直す。
    """
    images = generator.current_images
    if not images:
        return

    layout = generator.layout
    if generator.slide_has_text or generator.forced_layout == '2-column':
        area_left, area_width = layout.split_right_left, layout.split_right_width
    else:
        area_left, area_width = layout.content_left, layout.content_width
    area_top, area_height = layout.content_top, layout.content_height

    columns, rows = image_grid(len(images))
    gutter = Inches(IMAGE_GUTTER_INCHES)
    cell_width = (area_width - gutter * (columns - 1)) / columns
    cell_height = (area_height - gutter * (rows - 1)) / rows

    for index, picture in enumerate(images):
        row, column = divmod(index, columns)
        fit_shape_into(
            picture,
            Emu(int(area_left + (cell_width + gutter) * column)),
            Emu(int(area_top + (cell_height + gutter) * row)),
            Emu(int(cell_width)),
            Emu(int(cell_height)),
        )

def load_image(src: str) -> ImageSource:
    """画像URLならダウンロードし、ローカルパスならそのまま返す"""
    if src.startswith(('http://', 'https://')):
        response = requests.get(src, timeout=HTTP_TIMEOUT_SEC)
        response.raise_for_status()
        return BytesIO(response.content)
    return src

def process_image(generator: PPTXGenerator, tag: Tag) -> None:
    """画像の挿入処理"""
    img_url = tag.get('src')
    if not img_url:
        print("Warning: src属性が無い画像をスキップしました。")
        return

    try:
        img_data = load_image(cast(str, img_url))
        pos = generator.images_conf.get('position_inches')
        
        if pos and len(pos) >= 2:
            # YAMLの固定位置（幅は縦横比に従うため、上限はスライド幅とする）
            fixed_height = Inches(generator.images_conf.get('default_height_inches', 3.5))
            dpi = image_dpi(generator)
            if dpi:
                img_data = downscale_image(img_data, generator.prs.slide_width, fixed_height, dpi)
            picture = generator.current_slide.shapes.add_picture(
                img_data, Inches(pos[0]), Inches(pos[1]), height=fixed_height
            )
        elif generator.forced_layout == 'center':
            picture = place_image_full(generator, img_data)
        else:
            # オートレイアウト
            if generator.slide_has_text or generator.forced_layout == '2-column':
                picture = place_image_split(generator, img_data)
            else:
                picture = place_image_full(generator, img_data)

        # Markdownの代替テキスト（![説明](...) の「説明」）を引き継ぐ
        set_alt_text(picture, str(tag.get('alt') or ''))

        # 固定位置の指定がある場合は利用者の指定を尊重し、自動配置の対象にしない
        if not (pos and len(pos) >= 2):
            generator.current_images.append(picture)
            arrange_images(generator)
    except Exception as e:
        print(f"Warning: 画像の挿入に失敗しました: {e}")

def scaled_font(font_conf: FontConfig, scale: float) -> FontConfig:
    """フォント設定のサイズだけを縮小した写しを返す"""
    if scale >= 1.0 or 'size_pt' not in font_conf:
        return font_conf
    return {**font_conf, 'size_pt': max(1, int(font_conf['size_pt'] * scale))}

def cell_alignment(cell: Tag) -> PP_ALIGN | None:
    """Markdownの列揃え指定（`:---:` など）をPowerPointの配置に変換する

    markdown拡張は `style="text-align: center;"` を、古い版は `align="center"` を出力する。
    指定が無い場合は None を返し、PowerPoint側の既定に任せる。
    """
    value = cell.get('align')
    if not value:
        match = re.search(r'text-align:\s*(left|center|right)', str(cell.get('style') or ''))
        value = match.group(1) if match else None
    return CELL_ALIGNMENTS.get(str(value)) if value else None

def table_row_metrics(
    rows: list[Tag], header_conf: FontConfig, body_conf: FontConfig
) -> list[TableRowMetrics]:
    """表の各行から、高さの概算に必要な情報を取り出す"""
    metrics = []
    for row in rows:
        cells = row.find_all(['th', 'td'])
        is_header = any(cell.name == 'th' for cell in cells)
        conf = header_conf if is_header else body_conf
        metrics.append(TableRowMetrics(
            texts=[cell.get_text() for cell in cells],
            font_size_pt=float(conf.get('size_pt') or DEFAULT_TABLE_BODY_FONT['size_pt']),
        ))
    return metrics

def start_table_continuation_slide(generator: PPTXGenerator) -> None:
    """表の続きを載せる新しいスライドを作る（見出しは元のタイトル＋「続き」）"""
    title_shape = generator.current_slide.shapes.title if generator.current_slide else None
    base_title = title_shape.text if title_shape and title_shape.text else '表'
    base_title = base_title.removesuffix(TABLE_CONTINUATION_SUFFIX)

    heading = BeautifulSoup('<h2></h2>', 'html.parser').h2
    assert heading is not None  # 直前に組み立てたHTMLなので必ず存在する
    heading.string = f"{base_title}{TABLE_CONTINUATION_SUFFIX}"
    process_heading(generator, heading)

def render_table_page(
    generator: PPTXGenerator,
    page_rows: list[Tag],
    row_heights: list[float],
    num_cols: int,
    table_top: Length,
    scale: float,
    header_conf: FontConfig,
    body_conf: FontConfig,
) -> None:
    """1スライドぶんの表を描画する"""
    layout = generator.layout
    table_shape = generator.current_slide.shapes.add_table(
        len(page_rows), num_cols,
        layout.content_left, table_top, layout.content_width, Pt(sum(row_heights)),
    )
    table = table_shape.table

    # 行の高さを明示する。既定では総高さを行数で均等割りするため、行数が多いと
    # 1行あたりが極端に小さくなり、描画時にPowerPointが押し広げてはみ出す。
    for row_idx, height in enumerate(row_heights):
        table.rows[row_idx].height = Pt(height)

    for row_idx, row in enumerate(page_rows):
        for col_idx, col in enumerate(row.find_all(['th', 'td'])):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = cell_alignment(col)

            if col.name == 'th':
                font_conf = scaled_font(header_conf, scale)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 73, 125) # 濃い青をデフォルトに設定
            else:
                font_conf = scaled_font(body_conf, scale)

            add_runs_from_tag(col, p, font_conf, inline_code_conf(generator))

def process_table(generator: PPTXGenerator, tag: Tag) -> None:
    """表の挿入処理

    行数が多い場合は、まず文字を縮めて1枚に収めることを試み、
    それでも収まらなければ見出し行を繰り返しつつ複数スライドに分割する。
    """
    layout = generator.layout
    rows = tag.find_all('tr')
    if not rows: return

    num_cols = max(len(row.find_all(['th', 'td'])) for row in rows)
    if num_cols == 0: return

    if generator.slide_has_text:
        shrink_body_shape(
            generator.current_slide,
            layout.content_width,
            max_height=layout.table_split_body_height,
        )
        table_top = layout.table_split_top
    else:
        table_top = layout.content_top

    header_conf = generator.fonts_conf.get('table_header', DEFAULT_TABLE_HEADER_FONT)
    body_conf = generator.fonts_conf.get('table_body', DEFAULT_TABLE_BODY_FONT)

    metrics = table_row_metrics(rows, header_conf, body_conf)
    column_width = Emu(int(layout.content_width / num_cols)).pt - CELL_MARGIN_H_PT * 2
    cell_margin = CELL_MARGIN_V_PT * 2
    available_height = Emu(int(layout.content_top + layout.content_height - table_top)).pt

    # まず縮小で1枚に収めることを試みる
    scale = fit_table_scale(metrics, column_width, cell_margin, available_height)
    row_heights = estimate_row_heights_pt(metrics, column_width, cell_margin, scale)
    if sum(row_heights) > available_height:
        # 縮小しても収まらない場合は、読みづらくするより分割を優先する
        scale = 1.0
        row_heights = estimate_row_heights_pt(metrics, column_width, cell_margin, scale)

    has_header = any(cell.name == 'th' for cell in rows[0].find_all(['th', 'td']))
    pages = paginate_row_heights(row_heights, available_height, repeat_first_row=has_header)
    if len(pages) > 1:
        print(f"INFO: 表が1枚に収まらないため、{len(pages)}枚のスライドに分割しました。")

    for page_no, indices in enumerate(pages):
        if page_no > 0:
            start_table_continuation_slide(generator)
            table_top = layout.content_top

        render_table_page(
            generator,
            [rows[i] for i in indices],
            [row_heights[i] for i in indices],
            num_cols, table_top, scale, header_conf, body_conf,
        )
    generator.slide_has_text = True

def process_code_or_mermaid(generator: PPTXGenerator, tag: Tag) -> None:
    """コードブロックまたはMermaid図形の処理"""
    code_tag = tag.find('code')
    # class属性が無いコードブロックでは get('class') が None を返すため、必ず空リストへ倒す
    classes: list[str] = cast(
        "list[str]", (code_tag.get('class') or []) if isinstance(code_tag, Tag) else []
    )
    is_mermaid = 'language-mermaid' in classes or 'mermaid' in classes

    language = None
    for cls in classes:
        if cls.startswith('language-'):
            language = cls.replace('language-', '')
            break

    if is_mermaid and isinstance(code_tag, Tag):
        try:
            print("INFO: Mermaid図形を生成中...")
            image = render_mermaid(mermaid_conf(generator), code_tag.get_text())
            if image is None:  # renderer: off
                return

            if generator.slide_has_text:
                picture = place_image_split(generator, BytesIO(image))
            else:
                picture = place_image_full(generator, BytesIO(image))
            generator.current_images.append(picture)
            arrange_images(generator)
        except Exception as e:
            print(f"Warning: Mermaid図形の生成に失敗しました: {e}")
    else:
        append_code_textbox(generator, tag.get_text(), language=language)

def bullet_font_conf(
    fonts_conf: dict[str, FontConfig], level: int, ordered: bool = False
) -> FontConfig | None:
    """指定レベルのリストのフォント設定を返す

    番号付きリスト（ol）は ordered_level_N を優先し、無ければ
    箇条書き（ul）の bullet_level_N を使う。手順書だけ書式を変えたい場合に
    ordered_level_N を足すだけで済むようにするため。

    設定が無いレベルは、より浅いレベルの設定を引き継ぐ。
    いきなり body へフォールバックすると、ネストした途端に書体やサイズが
    変わってしまうため（body は本文用で、リストとは役割が異なる）。
    """
    prefixes = ('ordered_level_', 'bullet_level_') if ordered else ('bullet_level_',)
    for prefix in prefixes:
        for candidate in range(level + 1, 0, -1):
            conf = fonts_conf.get(f'{prefix}{candidate}')
            if conf:
                return conf
    return fonts_conf.get('body')

def is_ordered_item(tag: Tag) -> bool:
    """リスト項目が番号付きリスト（ol）に属するかを判定する"""
    parent = tag.find_parent(['ul', 'ol'])
    return parent is not None and parent.name == 'ol'

def process_text(generator: PPTXGenerator, tag: Tag) -> None:
    """段落・リストの処理"""
    if not tag.get_text(strip=True): return

    ordered = False
    if tag.name == 'li':
        level = min(len(tag.find_parents(['ul', 'ol'])) - 1, MAX_BULLET_LEVEL)
        ordered = is_ordered_item(tag)
        font_conf = bullet_font_conf(generator.fonts_conf, level, ordered)
    else:
        level = 0
        font_conf = generator.fonts_conf.get('body')


    paragraph = append_text_block(
        generator.current_body, tag,
        reuse_first_paragraph=not generator.slide_has_text,
        level=level,
        font_conf=font_conf,
        inline_code_conf=inline_code_conf(generator),
    )
    if ordered:
        # PowerPoint側で採番させる（項目を入れ替えても番号が崩れない）
        apply_auto_numbering(paragraph)
    generator.slide_has_text = True
