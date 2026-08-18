"""表をスライドに描画する

行数が多い表は、まず文字を縮めて1枚に収めることを試み、それでも収まらなければ
見出し行を繰り返しながら複数のスライドに分割する。この判断と描画をまとめている。
"""

from __future__ import annotations

import re
from typing import TYPE_CHECKING

from bs4 import BeautifulSoup, Tag
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Length, Pt

from processors import inline_code_conf, process_heading
from text_metrics import (
    TableRowMetrics,
    estimate_row_heights_pt,
    fit_table_scale,
    paginate_row_heights,
)
from utils import FontConfig, add_runs_from_tag, shrink_body_shape

if TYPE_CHECKING:
    from generator import PPTXGenerator

# 表の既定の書式（config.yaml に該当設定が無い場合）
DEFAULT_TABLE_HEADER_FONT: FontConfig = {
    'name': 'Meiryo', 'size_pt': 14, 'bold': True, 'color_rgb': [255, 255, 255],
}
DEFAULT_TABLE_BODY_FONT: FontConfig = {'name': 'Meiryo', 'size_pt': 12}

# 分割された表の続きスライドのタイトルに付ける接尾辞
TABLE_CONTINUATION_SUFFIX = '（続き）'

# Markdownの列揃え指定と PowerPoint の配置の対応
CELL_ALIGNMENTS = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
}

# python-pptx が設定するセルの既定マージン（ポイント）
CELL_MARGIN_H_PT = 0.1 * 72
CELL_MARGIN_V_PT = 0.05 * 72

# 見出し行の背景色
HEADER_FILL_RGB = (31, 73, 125)


def scaled_font(font_conf: FontConfig, scale: float) -> FontConfig:
    """フォント設定のサイズだけを縮小した写しを返す"""
    if scale >= 1.0 or 'size_pt' not in font_conf:
        return font_conf
    return {**font_conf, 'size_pt': max(1, int(font_conf['size_pt'] * scale))}

def cell_alignment(cell: Tag) -> PP_ALIGN | None:
    """Markdownの列揃え指定（`:---:` など）をPowerPointの配置に変換する

    markdown拡張は `style="text-align: center;"` を、古い版は `align="center"` を出力する。
    指定が無い場合は None を返し、PowerPoint側の既定に任せる。
    """
    value = cell.get('align')
    if not value:
        match = re.search(r'text-align:\s*(left|center|right)', str(cell.get('style') or ''))
        value = match.group(1) if match else None
    return CELL_ALIGNMENTS.get(str(value)) if value else None

def table_row_metrics(
    rows: list[Tag], header_conf: FontConfig, body_conf: FontConfig
) -> list[TableRowMetrics]:
    """表の各行から、高さの概算に必要な情報を取り出す"""
    metrics = []
    for row in rows:
        cells = row.find_all(['th', 'td'])
        is_header = any(cell.name == 'th' for cell in cells)
        conf = header_conf if is_header else body_conf
        metrics.append(TableRowMetrics(
            texts=[cell.get_text() for cell in cells],
            font_size_pt=float(conf.get('size_pt') or DEFAULT_TABLE_BODY_FONT['size_pt']),
            font_name=conf.get('name') or DEFAULT_TABLE_BODY_FONT['name'],
        ))
    return metrics

def start_table_continuation_slide(generator: PPTXGenerator) -> None:
    """表の続きを載せる新しいスライドを作る（見出しは元のタイトル＋「続き」）"""
    title_shape = generator.current_slide.shapes.title if generator.current_slide else None
    base_title = title_shape.text if title_shape and title_shape.text else '表'
    base_title = base_title.removesuffix(TABLE_CONTINUATION_SUFFIX)

    heading = BeautifulSoup('<h2></h2>', 'html.parser').h2
    assert heading is not None  # 直前に組み立てたHTMLなので必ず存在する
    heading.string = f"{base_title}{TABLE_CONTINUATION_SUFFIX}"
    process_heading(generator, heading)

def render_table_page(
    generator: PPTXGenerator,
    page_rows: list[Tag],
    row_heights: list[float],
    num_cols: int,
    table_top: Length,
    scale: float,
    header_conf: FontConfig,
    body_conf: FontConfig,
) -> None:
    """1スライドぶんの表を描画する"""
    layout = generator.layout
    table_shape = generator.current_slide.shapes.add_table(
        len(page_rows), num_cols,
        layout.content_left, table_top, layout.content_width, Pt(sum(row_heights)),
    )
    table = table_shape.table

    # 行の高さを明示する。既定では総高さを行数で均等割りするため、行数が多いと
    # 1行あたりが極端に小さくなり、描画時にPowerPointが押し広げてはみ出す。
    for row_idx, height in enumerate(row_heights):
        table.rows[row_idx].height = Pt(height)

    for row_idx, row in enumerate(page_rows):
        for col_idx, col in enumerate(row.find_all(['th', 'td'])):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = cell_alignment(col)

            if col.name == 'th':
                font_conf = scaled_font(header_conf, scale)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*HEADER_FILL_RGB)
            else:
                font_conf = scaled_font(body_conf, scale)

            add_runs_from_tag(col, p, font_conf, inline_code_conf(generator))

def process_table(generator: PPTXGenerator, tag: Tag) -> None:
    """表の挿入処理

    行数が多い場合は、まず文字を縮めて1枚に収めることを試み、
    それでも収まらなければ見出し行を繰り返しつつ複数スライドに分割する。
    """
    layout = generator.layout
    rows = tag.find_all('tr')
    if not rows: return

    num_cols = max(len(row.find_all(['th', 'td'])) for row in rows)
    if num_cols == 0: return

    if generator.slide_has_text:
        shrink_body_shape(
            generator.current_slide,
            layout.content_width,
            max_height=layout.table_split_body_height,
        )
        table_top = layout.table_split_top
    else:
        table_top = layout.content_top

    header_conf = generator.fonts_conf.get('table_header', DEFAULT_TABLE_HEADER_FONT)
    body_conf = generator.fonts_conf.get('table_body', DEFAULT_TABLE_BODY_FONT)

    metrics = table_row_metrics(rows, header_conf, body_conf)
    column_width = Emu(int(layout.content_width / num_cols)).pt - CELL_MARGIN_H_PT * 2
    cell_margin = CELL_MARGIN_V_PT * 2
    available_height = Emu(int(layout.content_top + layout.content_height - table_top)).pt

    # まず縮小で1枚に収めることを試みる
    scale = fit_table_scale(metrics, column_width, cell_margin, available_height)
    row_heights = estimate_row_heights_pt(metrics, column_width, cell_margin, scale)
    if sum(row_heights) > available_height:
        # 縮小しても収まらない場合は、読みづらくするより分割を優先する
        scale = 1.0
        row_heights = estimate_row_heights_pt(metrics, column_width, cell_margin, scale)

    has_header = any(cell.name == 'th' for cell in rows[0].find_all(['th', 'td']))
    pages = paginate_row_heights(row_heights, available_height, repeat_first_row=has_header)
    if len(pages) > 1:
        print(f"INFO: 表が1枚に収まらないため、{len(pages)}枚のスライドに分割しました。")

    for page_no, indices in enumerate(pages):
        if page_no > 0:
            start_table_continuation_slide(generator)
            table_top = layout.content_top

        render_table_page(
            generator,
            [rows[i] for i in indices],
            [row_heights[i] for i in indices],
            num_cols, table_top, scale, header_conf, body_conf,
        )
    generator.slide_has_text = True
