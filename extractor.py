"""PowerPoint（.pptx）から Markdown を復元する

md2pptx が生成した構造を読み取り、元のMarkdownに近い形へ戻す。
スライドの配置や配色は再現できないため、復元するのは内容（見出し・本文・
リスト・表・画像・コード・ノート）に限る。

他のツールで作られたPPTXも読み取れるが、その場合は本ツールの慣習
（レイアウト0＝タイトルスライド、下端のテキストはフッター等）に沿って
解釈するため、意図どおりにならないことがある。
"""

from __future__ import annotations

import hashlib
import os
from dataclasses import dataclass, field

from pptx.enum.dml import MSO_FILL
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

#: スライド下端のこの範囲にある独立テキストは、フッターやページ番号とみなす
FOOTER_BAND_INCHES = 0.8

#: インラインコードとみなす等幅フォント
MONOSPACE_FONTS = (
    'consolas', 'courier new', 'courier', 'menlo', 'monaco', 'sf mono',
    'dejavu sans mono', 'ms gothic', 'osaka-mono',
)

#: 表の列揃えと Markdown の記法の対応
ALIGNMENT_MARKS = {
    PP_ALIGN.LEFT: ':---',
    PP_ALIGN.CENTER: ':---:',
    PP_ALIGN.RIGHT: '---:',
}
DEFAULT_ALIGNMENT_MARK = '---'

#: 箇条書き1段あたりのインデント
INDENT = '    '

#: 代替テキストとみなさない拡張子（python-pptx は元のファイル名を既定で入れるため）
IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')


@dataclass
class ExtractedImage:
    """PPTXから取り出した画像"""

    filename: str
    blob: bytes


@dataclass
class ExtractionResult:
    """変換結果"""

    markdown: str = ''
    images: list[ExtractedImage] = field(default_factory=list)


def is_footer_shape(shape, slide_height) -> bool:
    """スライド下端にある装飾（日付・フッター文言・ページ番号）かどうか"""
    if shape.is_placeholder or not shape.has_text_frame or shape.top is None:
        return False
    band_top = slide_height - Emu(int(FOOTER_BAND_INCHES * 914400))
    return shape.top >= band_top


def is_code_box(shape) -> bool:
    """コードブロック用のテキストボックスかどうか（背景色で塗られている）"""
    if shape.is_placeholder or not shape.has_text_frame:
        return False
    try:
        return shape.fill.type == MSO_FILL.SOLID
    except Exception:
        return False


def is_monospace(run) -> bool:
    """等幅フォント（インラインコード）かどうか"""
    name = (run.font.name or '').lower()
    return name in MONOSPACE_FONTS


def has_bullet_tag(paragraph, tag: str) -> bool:
    """段落の行頭記号の指定を調べる"""
    return f'<a:{tag}' in paragraph._element.xml


def run_to_markdown(run, plain: bool = False) -> str:
    """1つの run を装飾付きのMarkdownに戻す

    plain が真の場合は装飾を付けない（表の見出し行など、Markdown側で
    既に強調されている箇所で二重にならないようにするため）。
    """
    text = run.text
    if not text.strip():
        return text
    if plain:
        return text

    if is_monospace(run):
        return f'`{text}`'
    if run.font.bold:
        return f'**{text}**'
    if run.font.italic:
        return f'*{text}*'
    return text


def paragraph_to_markdown(paragraph, plain: bool = False) -> str:
    """段落のテキストを装飾付きで組み立てる"""
    return ''.join(run_to_markdown(run, plain) for run in paragraph.runs).strip()


def body_to_markdown(text_frame, as_list: bool = True) -> list[str]:
    """本文プレースホルダーを Markdown の行に戻す

    本文プレースホルダーは既定ですべての段落に行頭記号が付くため、
    平文の段落と第1階層の箇条書きはXML上で区別できない。
    スライド上は箇条書きとして見えているので、リスト項目として書き出す。
    （タイトルスライドのサブタイトルなど、リストでない枠は as_list=False）
    """
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = paragraph_to_markdown(paragraph)
        if not text:
            continue

        if has_bullet_tag(paragraph, 'buNone'):
            # 行頭記号を消した段落は小見出し（h3）として書き出していたもの。
            # 見出しは Markdown 側で強調されるため、装飾は付け直さない
            lines.append(f'### {paragraph_to_markdown(paragraph, plain=True)}')
        elif has_bullet_tag(paragraph, 'buAutoNum'):
            lines.append(f'{INDENT * paragraph.level}1. {text}')
        elif as_list:
            lines.append(f'{INDENT * paragraph.level}* {text}')
        else:
            lines.append(text)
        lines.append('')
    return lines


def table_to_markdown(table) -> list[str]:
    """表を Markdown の表に戻す"""
    if not table.rows:
        return []

    lines = []
    for row_index, row in enumerate(table.rows):
        is_header = row_index == 0
        cells = [
            paragraph_to_markdown(cell.text_frame.paragraphs[0], plain=is_header).replace('|', '\\|')
            for cell in row.cells
        ]
        lines.append('| ' + ' | '.join(cells) + ' |')

        if is_header:
            marks = [
                ALIGNMENT_MARKS.get(cell.text_frame.paragraphs[0].alignment, DEFAULT_ALIGNMENT_MARK)
                for cell in row.cells
            ]
            lines.append('| ' + ' | '.join(marks) + ' |')

    lines.append('')
    return lines


def code_box_to_markdown(shape) -> list[str]:
    """コード枠を フェンス付きコードブロックに戻す"""
    text = '\n'.join(
        ''.join(run.text for run in paragraph.runs)
        for paragraph in shape.text_frame.paragraphs
    )
    return ['```', text.rstrip(), '```', '']


class ImageCollector:
    """取り出した画像を集める

    同じ内容の画像は1つのファイルにまとめ、同じファイル名を使い回す。
    連番を振ってしまうと、まとめた結果として存在しないファイルを
    参照することになるため。
    """

    def __init__(self) -> None:
        self._filenames: dict[str, str] = {}
        self.images: list[ExtractedImage] = []

    def add(self, blob: bytes, ext: str) -> str:
        """画像を登録し、参照するファイル名を返す"""
        digest = hashlib.sha1(blob).hexdigest()
        if digest not in self._filenames:
            filename = f'image{len(self.images) + 1}.{ext}'
            self._filenames[digest] = filename
            self.images.append(ExtractedImage(filename, blob))
        return self._filenames[digest]


def picture_to_markdown(shape, collector: ImageCollector) -> str:
    """画像を Markdown の記法に戻す（画像自体は collector に集める）"""
    image = shape.image
    filename = collector.add(image.blob, image.ext)
    alt = shape._element._nvXxPr.cNvPr.get('descr') or ''
    # python-pptx は代替テキストが無い場合に元のファイル名を入れる。
    # 説明として意味がないため、ファイル名に見える値は落とす
    if alt.lower().endswith(IMAGE_EXTENSIONS):
        alt = ''
    return f'![{alt}]({{image_dir}}/{filename})' 


def notes_to_markdown(slide) -> list[str]:
    """スピーカーノートを引用記法に戻す"""
    if not slide.has_notes_slide:
        return []
    text = slide.notes_slide.notes_text_frame.text.strip()
    if not text:
        return []

    lines = [f'> {line}' if line else '>' for line in text.split('\n')]
    lines.append('')
    return lines


def slide_to_markdown(
    slide, is_title_layout: bool, slide_height, collector: ImageCollector
) -> list[str]:
    """1枚のスライドを Markdown の行に戻す"""
    lines: list[str] = []

    title_shape = slide.shapes.title
    if title_shape and title_shape.text_frame.text.strip():
        prefix = '#' if is_title_layout else '##'
        lines.append(f'{prefix} {title_shape.text_frame.text.strip()}')
        lines.append('')
    else:
        # タイトルの無いスライドは水平線で区切っていたもの
        lines.append('---')
        lines.append('')

    for shape in slide.shapes:
        if shape == title_shape or is_footer_shape(shape, slide_height):
            continue

        if shape.is_placeholder and shape.has_text_frame:
            # タイトルスライドの2つ目の枠はサブタイトルなので、箇条書きにしない
            lines.extend(body_to_markdown(shape.text_frame, as_list=not is_title_layout))
        elif shape.has_table:
            lines.extend(table_to_markdown(shape.table))
        elif is_code_box(shape):
            lines.extend(code_box_to_markdown(shape))
        elif shape.shape_type is not None and 'PICTURE' in str(shape.shape_type):
            lines.append(picture_to_markdown(shape, collector))
            lines.append('')

    lines.extend(notes_to_markdown(slide))
    return lines


def extract(prs, image_dir: str = 'images') -> ExtractionResult:
    """プレゼンテーション全体を Markdown に戻す"""
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) else None

    lines: list[str] = []
    collector = ImageCollector()

    for slide in prs.slides:
        lines.extend(slide_to_markdown(
            slide, slide.slide_layout == title_layout, prs.slide_height, collector
        ))

    markdown = '\n'.join(lines).replace('{image_dir}', image_dir)
    # 空行が続く箇所を1つにまとめる
    while '\n\n\n' in markdown:
        markdown = markdown.replace('\n\n\n', '\n\n')

    return ExtractionResult(markdown.strip() + '\n', collector.images)


def write_images(images: list[ExtractedImage], directory: str) -> None:
    """取り出した画像をディレクトリへ書き出す"""
    if not images:
        return
    os.makedirs(directory, exist_ok=True)
    for image in images:
        with open(os.path.join(directory, image.filename), 'wb') as f:
            f.write(image.blob)
