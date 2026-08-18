"""Markdown全体を読み取り、プレゼンテーションを組み立てるジェネレーター"""

from __future__ import annotations

import os
from typing import Any

import yaml
from pptx import Presentation
from pptx.util import Inches
import markdown
from bs4 import BeautifulSoup, Comment, Tag
from layout import SlideLayout
from utils import FontConfig, apply_font_style, find_body_placeholder

from processors import (
    add_slide_footers,
    apply_layout_directives,
    finalize_slide,
    process_heading,
    process_h3,
    process_hr,
    process_blockquote,
    process_image,
    process_code_or_mermaid,
    process_text
)
from table_render import process_table

# スライドへ変換する対象のHTMLタグ
TARGET_TAGS = ['h1', 'h2', 'h3', 'hr', 'p', 'li', 'img', 'pre', 'table', 'blockquote']

# スライドレイアウトの用途と、名前が指定されない場合に使う既定の位置
DEFAULT_LAYOUT_INDEXES = {'title': 0, 'content': 1}

# use_template_fonts の指定時も残すフォント設定
# 等幅フォントや表見出しの文字色は、体裁ではなく可読性のために必要なもの
FUNCTIONAL_FONT_KEYS = ('inline_code', 'code_block')


class TemplateError(Exception):
    """テンプレートが本ツールの想定を満たしていないことを表す"""


def _has_body_placeholder(layout) -> bool:
    """本文を書き込めるプレースホルダーを持つレイアウトかどうか

    判定は find_body_placeholder に委ねる。ここだけ idx=1 で判定すると、
    実際には書き込めるテンプレートを門前払いしてしまう。
    """
    return find_body_placeholder(layout) is not None


class PPTXGenerator:
    """MarkdownをPowerPointに変換するジェネレータークラス"""

    def __init__(self, config: dict[str, Any] | None) -> None:
        self.config: dict[str, Any] = config or {}
        self.slides_conf: dict[str, Any] = self.config.get('slides') or {}
        self.fonts_conf: dict[str, FontConfig] = self.config.get('fonts') or {}
        self.images_conf: dict[str, Any] = self.config.get('images') or {}

        # スライドサイズから導出した配置寸法（_init_presentation で設定）
        self.layout: SlideLayout = None  # type: ignore[assignment]
        # 見出しの種類ごとに使うスライドレイアウト（_init_presentation で設定）
        self.slide_layouts: dict[str, Any] = {}

        # python-pptx のオブジェクト群（型が動的なため Any で保持する）
        self.prs: Any = None
        self.current_slide: Any = None
        self.current_body: Any = None
        # 現在のスライドに置いた画像（重ならないよう並べ直すために保持する）
        self.current_images: list[Any] = []
        self.slide_has_text: bool = False
        self.forced_layout: str | None = None
        # このスライドをダークテーマにするか（<!-- layout: dark-theme -->）
        self.dark_slide: bool = False

        self._init_presentation()

    def _init_presentation(self) -> None:
        """プレゼンテーションの初期化（テンプレート読み込み・サイズ設定）"""
        template_path = self.slides_conf.get('template_path')
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            # 完成済みの資料をテンプレートに指定すると、その中身がそのまま残る。
            # 気付かないまま「枚数が合わない」となりやすいので知らせる。
            if len(self.prs.slides) > 0:
                print(
                    f"Warning: テンプレートに既存のスライドが{len(self.prs.slides)}枚あります。"
                    "生成したスライドはその後ろに追加されます。\n"
                    "         書式だけを使いたい場合は、中身を消したファイルか .potx を指定してください。"
                )
        else:
            self.prs = Presentation()
            width, height = self._get_slide_size(self.slides_conf.get('layout', '16:9'))
            self.prs.slide_width = width
            self.prs.slide_height = height

        if self.slides_conf.get('use_template_fonts'):
            # テンプレート（テーマ）のフォントに任せる。
            # ただし等幅フォントなど、可読性のために必要な指定は残す
            self.fonts_conf = {
                key: value for key, value in self.fonts_conf.items()
                if key in FUNCTIONAL_FONT_KEYS
            }

        self._resolve_layouts()

        # 画像・表・コード枠の配置は、この寸法を基準に決める
        # （基準は、実際に本文で使うレイアウトの本文枠）
        self.layout = SlideLayout.from_presentation(self.prs, self.slide_layouts.get('content'))

    def _all_layouts(self) -> list[Any]:
        """全てのスライドマスターのレイアウト

        prs.slide_layouts は1つ目のマスターのものしか返さない。Googleスライドから
        書き出したファイルなどはマスターが複数あり、本文用のレイアウトが
        2つ目以降のマスターに入っていることがある。
        """
        return [layout for master in self.prs.slide_masters for layout in master.slide_layouts]

    def _resolve_layouts(self) -> None:
        """見出しの種類ごとに使うスライドレイアウトを決める

        config.yaml で名前が指定されていればそれを使い、無ければ位置で選ぶ。
        レイアウトの並び順はテンプレートによって異なるため、社内テンプレートでは
        名前で指定できる方が確実である。
        """
        names = self.slides_conf.get('layouts') or {}

        for kind, index in DEFAULT_LAYOUT_INDEXES.items():
            name = names.get(kind)
            self.slide_layouts[kind] = (
                self._layout_by_name(kind, name) if name
                else self._layout_by_index(index)
            )

        self._ensure_body_placeholder()

    def _layout_by_name(self, kind: str, name: str):
        for layout in self._all_layouts():
            if layout.name == name:
                return layout
        raise TemplateError(
            f"slides.layouts.{kind}: テンプレートに '{name}' というレイアウトがありません。\n"
            f"       使用できるレイアウト: {self._layout_names()}"
        )

    def _layout_by_index(self, index: int):
        layouts = self._all_layouts()
        if index < len(layouts):
            return layouts[index]
        raise TemplateError(
            f"テンプレートのレイアウトが足りません（{len(layouts)}個）。\n"
            f"       slides.layouts で使用するレイアウト名を指定してください。\n"
            f"       使用できるレイアウト: {self._layout_names()}"
        )

    def _ensure_body_placeholder(self) -> None:
        """本文用レイアウトに本文プレースホルダーがあることを確かめる

        無いレイアウト（「タイトルのみ」など）が選ばれていると、本文を書き込む
        段階で原因の分かりにくいエラーになるため、ここで気付けるようにする。
        """
        if _has_body_placeholder(self.slide_layouts['content']):
            return

        fallback = next(
            (l for l in self._all_layouts() if _has_body_placeholder(l)), None
        )
        if fallback is None:
            raise TemplateError(
                "テンプレートに本文を書き込めるレイアウトがありません"
                "（本文プレースホルダーを持つレイアウトが必要です）。\n"
                f"       使用できるレイアウト: {self._layout_names()}"
            )

        print(
            f"Warning: レイアウト '{self.slide_layouts['content'].name}' には本文枠が無いため、"
            f"'{fallback.name}' を使用します。\n"
            f"         slides.layouts.content で明示的に指定できます。"
        )
        self.slide_layouts['content'] = fallback

    def _layout_names(self) -> str:
        return ' / '.join(layout.name for layout in self._all_layouts())

    def _get_slide_size(self, layout_str: str) -> tuple[int, int]:
        sizes = {
            "16:9": (Inches(10), Inches(5.625)),
            "4:3":  (Inches(10), Inches(7.5)),
            "16:10": (Inches(10), Inches(6.25)),
            "A4":   (Inches(11.69), Inches(8.27))
        }
        return sizes.get(layout_str, sizes["16:9"])

    def generate(self, md_text: str, output_file: str) -> None:
        """MarkdownをパースしてPPTXを生成するメイン処理"""

        # 1. フロントマターの解析
        front_matter: dict[str, Any] = {}
        if md_text.startswith('---'):
            parts = md_text.split('---', 2)
            if len(parts) >= 3:
                try:
                    front_matter = yaml.safe_load(parts[1])
                    md_text = parts[2]
                except Exception as e:
                    print(f"Warning: フロントマターの解析に失敗しました: {e}")
        
        # タイトルスライドの自動生成（フロントマターがある場合）
        if front_matter and front_matter.get('title'):
            from pptx.enum.text import MSO_AUTO_SIZE
            
            self.current_slide = self.prs.slides.add_slide(self.slide_layouts['title'])
            
            title_shape = self.current_slide.shapes.title
            if title_shape:
                title_shape.text_frame.word_wrap = True
                title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                title_shape.text = str(front_matter.get('title'))
                
                for run in title_shape.text_frame.paragraphs[0].runs:
                    apply_font_style(run, self.fonts_conf.get('title_h1', self.fonts_conf.get('title')))
            
            subtitle_text = []
            if 'subtitle' in front_matter: subtitle_text.append(str(front_matter['subtitle']))
            if 'author' in front_matter: subtitle_text.append(str(front_matter['author']))
            if 'date' in front_matter: subtitle_text.append(str(front_matter['date']))
            
            sub_shape = find_body_placeholder(self.current_slide) if subtitle_text else None
            if subtitle_text and sub_shape is None:
                print(
                    "Warning: 表紙のレイアウトに副題を書き込める枠が無いため、"
                    "subtitle / author / date を省略しました。\n"
                    "         slides.layouts.title で別のレイアウトを指定できます。"
                )
            if sub_shape is not None:
                sub_shape.text_frame.word_wrap = True
                sub_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                sub_shape.text = "\n".join(subtitle_text)
                
                # サブタイトルには少し小さめのフォントサイズを適用
                sub_font_conf = self.fonts_conf.get('body', {'name': 'Meiryo', 'size_pt': 20})
                for p in sub_shape.text_frame.paragraphs:
                    for run in p.runs:
                        apply_font_style(run, sub_font_conf)

        html = markdown.markdown(md_text, extensions=['extra', 'fenced_code', 'sane_lists'])
        soup = BeautifulSoup(html, 'html.parser')

        # 2. タグとコメントの処理
        #
        # コメント（<!-- layout: ... -->）はTagではなくNavigableStringのため、
        # find_all(名前リスト) では拾えない。文書順を保ったまま両方を扱うために
        # descendants を走査する。
        for element in list(soup.descendants):
            if isinstance(element, Comment):
                text = element.strip()
                if text.startswith('layout:'):
                    apply_layout_directives(self, text.split(':', 1)[1])
                continue

            if not isinstance(element, Tag) or element.name not in TARGET_TAGS:
                continue

            tag = element
            # 引用の中身はノートとしてまとめて扱うため、本文側では処理しない
            # （p だけを除いていた頃は、引用内の箇条書きが本文にも出力されていた）
            if tag.name != 'blockquote' and tag.find_parent('blockquote'):
                continue
            # リスト項目の中の p は、項目側でまとめて処理する
            if tag.name == 'p' and tag.find_parent('li'):
                continue

            if tag.name in ['h1', 'h2']:
                self.forced_layout = None # 新しいスライドでリセット
                process_heading(self, tag)
            elif tag.name == 'h3':
                # 新規スライドにするか小見出しにするかは process_h3 が設定を見て判断する
                process_h3(self, tag)
            elif tag.name == 'hr':
                self.forced_layout = None
                process_hr(self, tag)
            elif tag.name == 'blockquote' and self.current_slide:
                process_blockquote(self, tag)
            elif tag.name == 'img' and self.current_slide:
                process_image(self, tag)
            elif tag.name == 'table' and self.current_slide:
                process_table(self, tag)
            elif tag.name == 'pre' and self.current_body:
                process_code_or_mermaid(self, tag)
            elif tag.name in ['li', 'p'] and self.current_body:
                process_text(self, tag)

        finalize_slide(self)
            
        # 日付・フッター文言・ページ番号をまとめて挿入する
        add_slide_footers(self)

        self.prs.save(output_file)
