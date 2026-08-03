"""MarkdownをPowerPoint（.pptx）に変換するコンバーター。"""

from __future__ import annotations

import argparse
import base64
import os
import sys
import traceback
import zlib
from dataclasses import dataclass
from io import BytesIO
from typing import Any, Union

import markdown
import requests
import yaml
from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# --- 定数 ---

#: スライドの画角名とサイズ（幅, 高さ）の対応表
SLIDE_SIZES: dict[str, tuple[int, int]] = {
    "16:9": (Inches(10), Inches(5.625)),
    "4:3": (Inches(10), Inches(7.5)),
    "16:10": (Inches(10), Inches(6.25)),
    "A4": (Inches(11.69), Inches(8.27)),
}
DEFAULT_LAYOUT_NAME = "16:9"

#: 画像の極端な拡大を防ぐための倍率上限
MAX_IMAGE_SCALE = 1.5

#: 外部API（画像取得・Kroki）のタイムアウト秒数
HTTP_TIMEOUT_SEC = 15

#: Mermaid図形をPNG化するKrokiのエンドポイント
KROKI_MERMAID_PNG_URL = "https://kroki.io/mermaid/png/"

#: スライドへ変換する対象のHTMLタグ
TARGET_TAGS = ["h1", "h2", "p", "li", "img", "pre", "table", "blockquote"]

#: インライン装飾として解釈しないブロック要素（別処理で扱うためスキップする）
BLOCK_TAGS = {"ul", "ol", "pre", "img", "table", "blockquote"}

#: 中身を再帰的にたどるコンテナ要素
CONTAINER_TAGS = {"p", "div", "span", "li", "th", "td"}

#: 箇条書きのインデントレベルの上限（PowerPointの仕様）
MAX_BULLET_LEVEL = 8

#: config.yaml に該当キーが無い場合に補完されるフォント設定
DEFAULT_FONTS: dict[str, dict[str, Any]] = {
    "inline_code": {"name": "Consolas", "color_rgb": [220, 20, 60]},
    "code_block": {"name": "Consolas", "size_pt": 12, "color_rgb": [0, 80, 160]},
    "table_header": {
        "name": "Meiryo",
        "size_pt": 14,
        "bold": True,
        "color_rgb": [255, 255, 255],
    },
    "table_body": {"name": "Meiryo", "size_pt": 12},
}

#: 画像データとして受け付ける型（ローカルパス、またはメモリ上のバイト列）
ImageSource = Union[str, BytesIO]


@dataclass(frozen=True)
class LayoutMetrics:
    """オートレイアウトで使用する配置寸法（インチ単位）。"""

    content_top: float = 1.5
    content_height: float = 3.8
    full_left: float = 1.0
    full_width: float = 8.0
    split_body_width: float = 4.8
    split_image_left: float = 5.2
    split_image_width: float = 4.5
    table_left: float = 1.0
    table_width: float = 8.0
    table_height: float = 0.8
    table_split_top: float = 2.8
    table_split_body_height: float = 2.0
    body_bottom_margin: float = 0.5
    default_image_height: float = 3.5


LAYOUT = LayoutMetrics()


class PPTXGenerator:
    """MarkdownをPowerPointに変換するジェネレータークラス"""

    def __init__(self, config: dict[str, Any] | None) -> None:
        self.config: dict[str, Any] = config or {}
        self.slides_conf: dict[str, Any] = self.config.get("slides") or {}
        self.fonts_conf: dict[str, Any] = self.config.get("fonts") or {}
        self.images_conf: dict[str, Any] = self.config.get("images") or {}

        # python-pptx のオブジェクト群（型が動的なため Any で保持する）
        self.prs: Any = None
        self.current_slide: Any = None
        self.current_body: Any = None
        self.current_body_shape: Any = None
        self.slide_has_text: bool = False

        self._init_presentation()

    # --- 初期化 ---

    def _init_presentation(self) -> None:
        """プレゼンテーションの初期化（テンプレート読み込み・サイズ設定）"""
        template_path = self.slides_conf.get("template_path")
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            width, height = self._get_slide_size(
                self.slides_conf.get("layout", DEFAULT_LAYOUT_NAME)
            )
            self.prs.slide_width = width
            self.prs.slide_height = height

    @staticmethod
    def _get_slide_size(layout_str: str) -> tuple[int, int]:
        """画角名からスライドサイズを取得する（未知の値はデフォルトにフォールバック）"""
        return SLIDE_SIZES.get(layout_str, SLIDE_SIZES[DEFAULT_LAYOUT_NAME])

    def _font(self, key: str) -> dict[str, Any]:
        """組み込みのデフォルト値にconfig.yamlの設定を重ねたフォント設定を返す"""
        user_conf = self.fonts_conf.get(key) or {}
        return {**DEFAULT_FONTS.get(key, {}), **user_conf}

    # --- 描画ヘルパー ---

    @staticmethod
    def apply_font_style(run, font_config: dict[str, Any] | None) -> None:
        """フォントスタイルの適用"""
        if not font_config:
            return
        font = run.font
        if "name" in font_config:
            font.name = font_config["name"]
        if "size_pt" in font_config:
            font.size = Pt(font_config["size_pt"])
        if "bold" in font_config:
            font.bold = font_config["bold"]
        if "color_rgb" in font_config:
            rgb = font_config["color_rgb"]
            font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])

    @staticmethod
    def insert_image_fit(
        slide,
        img_data: ImageSource,
        left: int,
        top: int,
        max_width: int,
        max_height: int,
    ):
        """画像を最大枠に収まるようにアスペクト比を保って自動縮小・中央配置する"""
        pic = slide.shapes.add_picture(img_data, left, top)
        ratio = min(max_width / pic.width, max_height / pic.height, MAX_IMAGE_SCALE)

        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(left + (max_width - pic.width) / 2)
        pic.top = int(top + (max_height - pic.height) / 2)
        return pic

    def _add_runs_from_tag(self, element, paragraph, default_font_conf) -> None:
        """インライン装飾を解釈しながらテキストを追加（再帰処理）"""
        for child in element:
            if isinstance(child, NavigableString):
                text = str(child).replace("\n", " ")
                if text.strip() or text == " ":
                    run = paragraph.add_run()
                    run.text = text
                    self.apply_font_style(run, default_font_conf)
            elif isinstance(child, Tag):
                if child.name in BLOCK_TAGS:
                    continue
                if child.name in CONTAINER_TAGS:
                    self._add_runs_from_tag(child, paragraph, default_font_conf)
                else:
                    run = paragraph.add_run()
                    run.text = child.get_text().replace("\n", " ")
                    self.apply_font_style(run, default_font_conf)

                    if child.name in ("strong", "b"):
                        run.font.bold = True
                    elif child.name in ("em", "i"):
                        run.font.italic = True
                    elif child.name == "code":
                        self.apply_font_style(run, self._font("inline_code"))

    @staticmethod
    def _resize_shape(shape, width: int | None = None, height: int | None = None) -> None:
        """プレースホルダーのサイズを変更する。

        プレースホルダーの位置・サイズはスライドレイアウトから継承されるため、
        一部だけを書き換えると残りの値が 0 になってしまう。
        継承値をいったん明示的に書き戻してから変更する。
        """
        left, top = shape.left, shape.top
        current_width, current_height = shape.width, shape.height
        shape.left, shape.top = left, top
        shape.width = current_width if width is None else width
        shape.height = current_height if height is None else height

    def _shrink_body_shape(
        self, width_inches: float, max_height_inches: float | None = None
    ) -> None:
        """テキスト枠を指定サイズに縮める（レイアウト調整用ヘルパー）"""
        if self.current_body_shape is None:
            return
        self._resize_shape(
            self.current_body_shape,
            width=Inches(width_inches),
            height=Inches(max_height_inches) if max_height_inches else None,
        )

    def _place_image(self, img_data: ImageSource) -> None:
        """テキストの有無に応じて画像を右半分（2カラム）または中央に配置する"""
        if self.slide_has_text:
            self._shrink_body_shape(LAYOUT.split_body_width)
            left, width = LAYOUT.split_image_left, LAYOUT.split_image_width
        else:
            left, width = LAYOUT.full_left, LAYOUT.full_width

        self.insert_image_fit(
            self.current_slide,
            img_data,
            Inches(left),
            Inches(LAYOUT.content_top),
            Inches(width),
            Inches(LAYOUT.content_height),
        )

    # --- 外部リソースの取得 ---

    @staticmethod
    def _load_image(src: str) -> ImageSource:
        """画像URLならダウンロードし、ローカルパスならそのまま返す"""
        if src.startswith(("http://", "https://")):
            response = requests.get(src, timeout=HTTP_TIMEOUT_SEC)
            response.raise_for_status()
            return BytesIO(response.content)
        return src

    @staticmethod
    def _render_mermaid(code: str) -> BytesIO:
        """Mermaid記法をKroki APIでPNG画像に変換する"""
        compressed = zlib.compress(code.encode("utf-8"), 9)
        encoded = base64.urlsafe_b64encode(compressed).decode("ascii")
        response = requests.get(
            f"{KROKI_MERMAID_PNG_URL}{encoded}", timeout=HTTP_TIMEOUT_SEC
        )
        response.raise_for_status()
        return BytesIO(response.content)

    # --- タグごとの処理 ---

    def _process_heading(self, tag) -> None:
        """見出しタグの処理とスライド作成"""
        is_title_slide = tag.name == "h1"
        layout_idx = 0 if is_title_slide else 1
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        self.current_slide.shapes.title.text = tag.get_text()

        style_key = "title_h1" if is_title_slide else "title_h2"
        font_conf = self._font(style_key) or self._font("title")
        for run in self.current_slide.shapes.title.text_frame.paragraphs[0].runs:
            self.apply_font_style(run, font_conf)

        body_shape = self.current_slide.placeholders[1]
        self.current_body_shape = body_shape
        self.current_body = body_shape.text_frame
        self.current_body.text = ""
        self.slide_has_text = False

        # デフォルト枠のはみ出し補正（テンプレート利用時は元のレイアウトを尊重する）
        if not self.slides_conf.get("template_path") and body_shape.top is not None:
            new_height = (
                self.prs.slide_height - body_shape.top - Inches(LAYOUT.body_bottom_margin)
            )
            self._resize_shape(body_shape, height=new_height)

    def _process_blockquote(self, tag) -> None:
        """スピーカーノートの処理"""
        text_frame = self.current_slide.notes_slide.notes_text_frame
        note_text = tag.get_text(strip=True)
        text_frame.text = (
            f"{text_frame.text}\n\n{note_text}" if text_frame.text else note_text
        )

    def _process_image(self, tag) -> None:
        """画像の挿入処理"""
        img_url = tag.get("src")
        if not img_url:
            print("Warning: src属性が無い画像をスキップしました。")
            return

        try:
            img_data = self._load_image(img_url)
            pos = self.images_conf.get("position_inches")

            if pos and len(pos) >= 2:
                # YAMLの固定位置
                height = self.images_conf.get(
                    "default_height_inches", LAYOUT.default_image_height
                )
                self.current_slide.shapes.add_picture(
                    img_data, Inches(pos[0]), Inches(pos[1]), height=Inches(height)
                )
            else:
                # オートレイアウト
                self._place_image(img_data)
        except Exception as e:
            print(f"Warning: 画像の挿入に失敗しました: {e}")

    def _process_table(self, tag) -> None:
        """表の挿入処理"""
        rows = tag.find_all("tr")
        if not rows:
            return

        num_rows = len(rows)
        num_cols = max(len(row.find_all(["th", "td"])) for row in rows)
        if num_cols == 0:
            return

        if self.slide_has_text:
            self._shrink_body_shape(
                LAYOUT.full_width, max_height_inches=LAYOUT.table_split_body_height
            )
            table_top = Inches(LAYOUT.table_split_top)
        else:
            table_top = Inches(LAYOUT.content_top)

        table_shape = self.current_slide.shapes.add_table(
            num_rows,
            num_cols,
            Inches(LAYOUT.table_left),
            table_top,
            Inches(LAYOUT.table_width),
            Inches(LAYOUT.table_height),
        )
        table = table_shape.table

        for row_idx, row in enumerate(rows):
            for col_idx, col in enumerate(row.find_all(["th", "td"])):
                cell = table.cell(row_idx, col_idx)
                cell.text = ""
                font_conf = self._font(
                    "table_header" if col.name == "th" else "table_body"
                )
                self._add_runs_from_tag(col, cell.text_frame.paragraphs[0], font_conf)

        self.slide_has_text = True

    def _process_code_or_mermaid(self, tag) -> None:
        """コードブロックまたはMermaid図形の処理"""
        code_tag = tag.find("code")

        if self._is_mermaid(code_tag):
            try:
                print("INFO: Mermaid図形をAPIで生成中...")
                self._place_image(self._render_mermaid(code_tag.get_text()))
            except Exception as e:
                print(f"Warning: Mermaid図形の生成に失敗しました: {e}")
        else:
            self._append_text_block(tag.get_text(), is_code=True)
            self.slide_has_text = True

    @staticmethod
    def _is_mermaid(code_tag) -> bool:
        """コードブロックがMermaid記法かどうかを判定する"""
        if code_tag is None:
            return False
        classes = code_tag.get("class") or []
        return "language-mermaid" in classes or "mermaid" in classes

    def _process_text(self, tag) -> None:
        """段落・リストの処理"""
        if not tag.get_text(strip=True):
            return

        if tag.name == "li":
            level = min(len(tag.find_parents(["ul", "ol"])) - 1, MAX_BULLET_LEVEL)
            font_conf = self._font(f"bullet_level_{level + 1}") or self._font("body")
        else:
            level = 0
            font_conf = self._font("body")

        self._append_text_block(tag, is_code=False, level=level, font_conf=font_conf)
        self.slide_has_text = True

    def _append_text_block(
        self,
        content,
        is_code: bool = False,
        level: int = 0,
        font_conf: dict[str, Any] | None = None,
    ) -> None:
        """段落オブジェクトを追加し、テキストまたはタグ構造を書き込むヘルパー"""
        first_paragraph = self.current_body.paragraphs[0]
        is_empty_body = (
            not self.slide_has_text
            and len(self.current_body.paragraphs) == 1
            and not first_paragraph.text
        )
        p = first_paragraph if is_empty_body else self.current_body.add_paragraph()
        p.level = level

        if is_code:
            run = p.add_run()
            run.text = content
            self.apply_font_style(run, self._font("code_block"))
        else:
            self._add_runs_from_tag(content, p, font_conf)

    # --- エントリポイント ---

    def generate(self, md_text: str, output_file: str) -> None:
        """MarkdownをパースしてPPTXを生成するメイン処理"""
        html = markdown.markdown(
            md_text, extensions=["extra", "fenced_code", "sane_lists"]
        )
        soup = BeautifulSoup(html, "html.parser")

        for tag in soup.find_all(TARGET_TAGS):
            # リスト・引用の内側のpタグは親要素側でまとめて処理する
            if tag.name == "p" and (
                tag.find_parent("li") or tag.find_parent("blockquote")
            ):
                continue
            self._dispatch(tag)

        self.prs.save(output_file)

    def _dispatch(self, tag) -> None:
        """タグの種類に応じた処理へ振り分ける"""
        if tag.name in ("h1", "h2"):
            self._process_heading(tag)
            return

        # 見出しより前に現れた要素は配置先のスライドが無いため無視する
        if self.current_slide is None:
            return

        if tag.name == "blockquote":
            self._process_blockquote(tag)
        elif tag.name == "img":
            self._process_image(tag)
        elif tag.name == "table":
            self._process_table(tag)
        elif tag.name == "pre":
            self._process_code_or_mermaid(tag)
        else:
            self._process_text(tag)


def load_config(path: str) -> dict[str, Any]:
    """YAML設定ファイルを読み込む（空ファイルの場合は空の設定を返す）"""
    with open(path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}


def read_text_file(path: str) -> str:
    """UTF-8のテキストファイルを読み込む"""
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    """コマンドライン引数を解析する"""
    parser = argparse.ArgumentParser(
        description="MarkdownファイルをPowerPointに変換します。"
    )
    parser.add_argument("input", help="変換するMarkdownファイルのパス")
    parser.add_argument("-o", "--output", help="出力ファイル名", default="output.pptx")
    parser.add_argument(
        "-c", "--config", help="YAML設定ファイルのパス", default="config.yaml"
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    """CLIのエントリポイント（0=成功, 1=失敗 の終了コードを返す）"""
    print("INFO: 変換処理を開始します...")
    args = parse_args(argv)

    if not os.path.exists(args.input):
        print(f"Error: 入力ファイル '{args.input}' が見つかりません。")
        return 1

    if not os.path.exists(args.config):
        print(f"Error: 設定ファイル '{args.config}' が見つかりません。")
        return 1

    try:
        config = load_config(args.config)
        content = read_text_file(args.input)

        generator = PPTXGenerator(config)
        generator.generate(content, args.output)

        print(f"Success: '{args.output}' の生成が完了しました！")
        return 0

    except PermissionError:
        print(
            f"Error: '{args.output}' に書き込めません！"
            "PowerPointでファイルを開いたままにしていませんか？閉じてから再実行してください。"
        )
        return 1
    except Exception as e:
        print(f"Error: 予期せぬエラーが発生しました: {e}")
        traceback.print_exc()
        return 1


if __name__ == "__main__":  # pragma: no cover
    sys.exit(main())
